import streamlit as st
import pandas as pd
import plotly.express as px
import plotly.graph_objects as go
from io import StringIO, BytesIO
import PyPDF2
from docx import Document
from docx.shared import Inches, Pt, RGBColor
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.enum.table import WD_TABLE_ALIGNMENT
from wordcloud import WordCloud
import matplotlib.pyplot as plt
import json
import requests
import base64
from datetime import datetime
import os
import pickle
from pathlib import Path

# 数据保存目录 - 使用脚本所在目录
SCRIPT_DIR = Path(__file__).parent if '__file__' in dir() else Path('.')
SAVE_DIR = SCRIPT_DIR / "saved_surveys"
SAVE_DIR.mkdir(exist_ok=True)

# 导出相关库
try:
    from pptx import Presentation
    from pptx.util import Inches as PptxInches, Pt as PptxPt
    from pptx.dml.color import RGBColor as PptxRGBColor
    from pptx.enum.text import PP_ALIGN
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False

try:
    from reportlab.lib import colors as rl_colors
    from reportlab.lib.colors import Color, HexColor
    from reportlab.lib.pagesizes import A4, landscape
    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
    from reportlab.lib.units import inch
    from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Image, Table, TableStyle, PageBreak
    from reportlab.pdfbase import pdfmetrics
    from reportlab.pdfbase.ttfonts import TTFont
    PDF_AVAILABLE = True
except ImportError:
    PDF_AVAILABLE = False

try:
    import kaleido
    KALEIDO_AVAILABLE = True
except ImportError:
    KALEIDO_AVAILABLE = False

# --- 页面配置 ---
st.set_page_config(
    page_title="用户调研智能看板 Pro", 
    layout="wide",
    page_icon="📊",
    initial_sidebar_state="expanded"
)

# === 关键修复：禁用文件监视和自动刷新 ===
import streamlit.config as _config
try:
    _config.set_option('server.fileWatcherType', 'none')
    _config.set_option('server.runOnSave', False)
except:
    pass

# --- shadcn UI 风格设计系统 - Ptengine 调研分析工具 ---
st.markdown("""
<style>
    /* ========== shadcn UI 设计系统 ========== */
    
    /* 导入 Inter 字体 - shadcn 默认字体 */
    @import url('https://fonts.googleapis.com/css2?family=Inter:wght@300;400;500;600;700&family=Noto+Sans+SC:wght@300;400;500;700&display=swap');
    
    /* CSS 变量 - shadcn 风格 */
    :root {
        /* 背景色 */
        --background: #ffffff;
        --foreground: #0a0a0a;
        --card: #ffffff;
        --card-foreground: #0a0a0a;
        --popover: #ffffff;
        --popover-foreground: #0a0a0a;
        
        /* 主色调 - Ptengine 品牌色 */
        --primary: #18181b;
        --primary-foreground: #fafafa;
        
        /* 辅助色 */
        --secondary: #f4f4f5;
        --secondary-foreground: #18181b;
        --muted: #f4f4f5;
        --muted-foreground: #71717a;
        --accent: #f4f4f5;
        --accent-foreground: #18181b;
        
        /* 边框和输入框 */
        --border: #e4e4e7;
        --input: #e4e4e7;
        --ring: #18181b;
        
        /* 语义色 */
        --destructive: #ef4444;
        --destructive-foreground: #fafafa;
        --success: #22c55e;
        --success-foreground: #fafafa;
        --warning: #f59e0b;
        --warning-foreground: #fafafa;
        
        /* 品牌色 - Ptengine 蓝 */
        --brand: #2563eb;
        --brand-foreground: #ffffff;
        --brand-muted: #dbeafe;
        
        /* 圆角 - shadcn 风格 */
        --radius: 0.5rem;
        --radius-sm: 0.375rem;
        --radius-lg: 0.75rem;
        --radius-xl: 1rem;
    }
    
    /* 隐藏 Streamlit 默认元素 */
    #MainMenu, footer, header {visibility: hidden;}
    
    /* 全局字体 - Inter */
    html, body, [class*="css"] {
        font-family: 'Inter', 'Noto Sans SC', -apple-system, BlinkMacSystemFont, sans-serif;
        -webkit-font-smoothing: antialiased;
        font-feature-settings: "cv02", "cv03", "cv04", "cv11";
        color: var(--foreground);
    }
    
    /* 主容器 */
    .main {
        background: #fafafa;
    }
    
    .block-container {
        padding: 1.5rem 2rem 2rem;
        max-width: 100%;
    }
    
    /* ========== 卡片组件 - shadcn Card ========== */
    .kpi-card {
        background: var(--card);
        border-radius: var(--radius-lg);
        padding: 1.5rem;
        border: 1px solid var(--border);
        box-shadow: 0 1px 2px 0 rgb(0 0 0 / 0.05);
    }
    
    .kpi-card:hover {
        box-shadow: 0 4px 6px -1px rgb(0 0 0 / 0.1);
    }
    
    .kpi-label {
        font-size: 0.875rem;
        font-weight: 500;
        color: var(--muted-foreground);
        margin-bottom: 0.5rem;
    }
    
    .kpi-value {
        font-size: 2rem;
        font-weight: 700;
        color: var(--foreground);
        line-height: 1;
        letter-spacing: -0.025em;
    }
    
    .kpi-delta {
        font-size: 0.75rem;
        font-weight: 500;
        display: inline-flex;
        align-items: center;
        gap: 0.25rem;
        padding: 0.25rem 0.625rem;
        border-radius: 9999px;
        margin-top: 0.5rem;
    }
    
    .kpi-delta.positive {
        background: #dcfce7;
        color: #166534;
    }
    
    .kpi-delta.negative {
        background: #fee2e2;
        color: #991b1b;
    }
    
    .kpi-delta.neutral {
        background: var(--secondary);
        color: var(--muted-foreground);
    }
    
    /* ========== 侧边栏 - shadcn Sidebar ========== */
    [data-testid="stSidebar"] {
        background: var(--card);
        border-right: 1px solid var(--border);
    }
    
    [data-testid="stSidebar"] > div:first-child {
        padding: 1.25rem 1rem !important;
    }
    
    [data-testid="stSidebar"] [data-testid="stVerticalBlock"] {
        gap: 0.75rem !important;
    }
    
    /* 侧边栏内组件间距 */
    [data-testid="stSidebar"] [data-testid="stExpander"] {
        margin-bottom: 0.75rem;
    }
    
    [data-testid="stSidebar"] .stSelectbox,
    [data-testid="stSidebar"] .stTextInput {
        margin-bottom: 0.5rem;
    }
    
    /* 侧边栏标签文字 */
    [data-testid="stSidebar"] label {
        font-size: 0.8rem !important;
        font-weight: 500 !important;
        color: #71717a !important;
        margin-bottom: 0.375rem !important;
    }
    
    /* 侧边栏分组标题 - 旧版（保留兼容） */
    .sidebar-section-title-old {
        font-size: 0.75rem;
        font-weight: 500;
        color: var(--muted-foreground);
        text-transform: uppercase;
        letter-spacing: 0.05em;
        margin: 1rem 0 0.5rem;
        padding-left: 0.5rem;
    }
    
    /* 文件上传成功 */
    .upload-success {
        background: #f0fdf4;
        border: 1px solid #bbf7d0;
        border-radius: var(--radius);
        padding: 0.75rem 1rem;
        margin: 0.5rem 0;
    }
    
    .upload-success-icon {
        display: inline-flex;
        width: 1.25rem;
        height: 1.25rem;
        background: #22c55e;
        border-radius: 50%;
        align-items: center;
        justify-content: center;
        color: white;
        font-size: 0.625rem;
        margin-right: 0.5rem;
    }
    
    .upload-success-text {
        font-weight: 600;
        color: #166534;
        font-size: 0.875rem;
    }
    
    .upload-file-info {
        color: #15803d;
        font-size: 0.75rem;
        margin-top: 0.25rem;
        padding-left: 1.75rem;
    }
    
    /* ========== 标签页 - shadcn Tabs ========== */
    .stTabs [data-baseweb="tab-list"] {
        gap: 0;
        background: var(--muted);
        padding: 0.25rem;
        border-radius: var(--radius-lg);
        border: none;
    }
    
    .stTabs [data-baseweb="tab"] {
        height: 2.25rem;
        background: transparent;
        border-radius: var(--radius);
        padding: 0 1rem;
        font-weight: 500;
        font-size: 0.875rem;
        color: var(--muted-foreground);
        border: none;
    }
    
    .stTabs [data-baseweb="tab"]:hover {
        color: var(--foreground);
    }
    
    .stTabs [aria-selected="true"] {
        background: var(--background) !important;
        color: var(--foreground) !important;
        box-shadow: 0 1px 2px 0 rgb(0 0 0 / 0.05);
    }
    
    /* ========== 按钮 - shadcn Button ========== */
    .stButton>button {
        border-radius: var(--radius);
        font-weight: 500;
        font-family: 'Inter', 'Noto Sans SC', sans-serif;
        border: 1px solid #e4e4e7;
        background: #ffffff;
        color: #18181b;
        padding: 0.5rem 1.25rem;
        font-size: 0.875rem;
        height: auto;
        box-shadow: 0 1px 2px rgba(0,0,0,0.05);
    }
    
    .stButton>button:hover {
        background: #f4f4f5 !important;
        border-color: #18181b !important;
        color: #18181b !important;
    }
    
    .stDownloadButton>button {
        background: #2563eb !important;
        color: #ffffff !important;
        border-color: #2563eb !important;
    }
    
    .stDownloadButton>button:hover {
        background: #1d4ed8 !important;
        border-color: #1d4ed8 !important;
        color: #ffffff !important;
    }
    
    /* Action 按钮 */
    .action-btn {
        display: inline-flex;
        align-items: center;
        gap: 0.5rem;
        padding: 0.5rem 1rem;
        height: 2.25rem;
        background: var(--primary);
        color: var(--primary-foreground);
        border-radius: var(--radius);
        font-size: 0.875rem;
        font-weight: 500;
        text-decoration: none;
        border: none;
        cursor: pointer;
        display: inline-flex;
        align-items: center;
        justify-content: center;
    }
    
    /* 所有按钮 hover 统一样式 */
    .stButton>button:hover,
    .stButton>button:focus,
    .stButton>button:active {
        background: #f4f4f5 !important;
        border-color: #18181b !important;
        color: #18181b !important;
    }
    
    /* 主要按钮（蓝色） */
    .stButton>button[kind="primary"] {
        background: #2563eb !important;
        color: #ffffff !important;
        border: 1px solid #2563eb !important;
    }
    
    .stButton>button[kind="primary"]:hover {
        background: #1d4ed8 !important;
        border-color: #1d4ed8 !important;
        color: #ffffff !important;
    }
    
    /* ========== 数据表格 - shadcn Table ========== */
    [data-testid="stDataFrame"] {
        border-radius: var(--radius);
        overflow: hidden;
        border: 1px solid var(--border);
    }
    
    [data-testid="stDataFrame"] table {
        font-family: 'Inter', 'Noto Sans SC', sans-serif;
        font-size: 0.875rem;
    }
    
    [data-testid="stDataFrame"] thead tr th {
        background: var(--muted);
        color: var(--muted-foreground);
        font-weight: 500;
        padding: 0.75rem 1rem;
        border-bottom: 1px solid var(--border);
        text-align: left;
        font-size: 0.75rem;
    }
    
    [data-testid="stDataFrame"] tbody tr td {
        padding: 0.75rem 1rem;
        border-bottom: 1px solid var(--border);
        color: var(--foreground);
    }
    
    [data-testid="stDataFrame"] tbody tr:hover td {
        background: var(--muted);
    }
    
    /* ========== 输入组件 - shadcn Input ========== */
    .stTextInput>div>div>input,
    .stTextArea>div>div>textarea,
    .stSelectbox>div>div,
    .stMultiSelect>div>div {
        border-radius: var(--radius);
        border: 1px solid var(--input);
        font-family: 'Inter', 'Noto Sans SC', sans-serif;
        font-size: 0.875rem;
        background: var(--background);
    }
    
    .stTextInput>div>div>input:focus,
    .stTextArea>div>div>textarea:focus {
        border-color: var(--ring);
        box-shadow: 0 0 0 2px var(--ring);
        outline: none;
    }
    
    /* ========== MultiSelect - shadcn Badge ========== */
    .stMultiSelect [data-baseweb="tag"] {
        background-color: var(--secondary) !important;
        border: 1px solid var(--border) !important;
        color: var(--secondary-foreground) !important;
        border-radius: var(--radius-sm) !important;
        font-size: 0.75rem !important;
        font-weight: 500 !important;
        padding: 0.125rem 0.5rem !important;
        margin: 2px !important;
    }
    
    .stMultiSelect [data-baseweb="tag"]:hover {
        background-color: var(--accent) !important;
    }
    
    .stMultiSelect [data-baseweb="tag"] span {
        color: var(--secondary-foreground) !important;
    }
    
    .stMultiSelect [data-baseweb="tag"] [data-baseweb="icon"] {
        color: var(--muted-foreground) !important;
    }
    
    .stMultiSelect>div>div:focus-within {
        border-color: var(--ring) !important;
        box-shadow: 0 0 0 2px var(--ring) !important;
    }
    
    /* ========== 展开器 - shadcn Accordion ========== */
    [data-testid="stExpander"] {
        border: 1px solid var(--border);
        border-radius: var(--radius);
        overflow: hidden;
        background: var(--card);
    }
    
    .streamlit-expanderHeader {
        font-weight: 500;
        font-size: 0.875rem;
        color: var(--foreground);
        padding: 1rem;
        background: transparent;
    }
    
    .streamlit-expanderHeader:hover {
        background: var(--muted);
    }
    
    /* ========== 提示消息 - shadcn Alert ========== */
    .stAlert {
        border-radius: var(--radius);
        border: 1px solid var(--border);
        padding: 1rem;
        font-size: 0.875rem;
    }
    
    /* ========== 图表容器 ========== */
    .chart-container {
        background: var(--card);
        border-radius: var(--radius);
        padding: 1.5rem;
        border: 1px solid var(--border);
        margin-bottom: 1rem;
    }
    
    .chart-title {
        font-size: 0.875rem;
        font-weight: 600;
        color: var(--foreground);
        margin-bottom: 1rem;
        display: flex;
        align-items: center;
        gap: 0.5rem;
    }
    
    /* ========== 问题卡片 ========== */
    .question-card {
        background: var(--card);
        border: 1px solid var(--border);
        border-radius: var(--radius);
        padding: 1.5rem;
        margin-bottom: 1rem;
    }
    
    .question-card:hover {
        border-color: var(--ring);
    }
    
    .question-header {
        display: flex;
        align-items: flex-start;
        gap: 0.75rem;
        margin-bottom: 1rem;
    }
    
    .question-number {
        background: var(--primary);
        color: white;
        min-width: 1.75rem;
        height: 1.75rem;
        border-radius: var(--radius-sm);
        display: inline-flex;
        align-items: center;
        justify-content: center;
        font-size: 0.8rem;
        font-weight: 700;
        flex-shrink: 0;
    }
    
    .question-title {
        font-size: 0.95rem;
        font-weight: 600;
        color: var(--gray-900);
        line-height: 1.4;
    }
    
    .question-type-badge {
        display: inline-flex;
        align-items: center;
        padding: 0.2rem 0.6rem;
        border-radius: 9999px;
        font-size: 0.65rem;
        font-weight: 600;
        text-transform: uppercase;
        letter-spacing: 0.03em;
        margin-left: auto;
    }
    
    .question-type-badge.single {
        background: var(--info-light);
        color: var(--info);
    }
    
    .question-type-badge.multi {
        background: var(--warning-light);
        color: var(--warning);
    }
    
    .question-type-badge.nps {
        background: var(--success-light);
        color: var(--success);
    }
    
    /* ========== NPS 仪表盘 ========== */
    .nps-gauge-container {
        text-align: center;
        padding: 1.5rem;
    }
    
    .nps-score {
        font-size: 3.5rem;
        font-weight: 800;
        color: var(--primary);
        line-height: 1;
    }
    
    .nps-label {
        font-size: 0.875rem;
        color: var(--gray-600);
        margin-top: 0.5rem;
    }
    
    .nps-breakdown {
        display: flex;
        justify-content: center;
        gap: 2rem;
        margin-top: 1.5rem;
        padding-top: 1.5rem;
        border-top: 1px solid var(--gray-200);
    }
    
    .nps-segment {
        text-align: center;
    }
    
    .nps-segment-value {
        font-size: 1.5rem;
        font-weight: 700;
    }
    
    .nps-segment-label {
        font-size: 0.75rem;
        color: var(--gray-600);
        margin-top: 0.25rem;
    }
    
    .nps-promoters .nps-segment-value { color: var(--success); }
    .nps-passives .nps-segment-value { color: var(--warning); }
    .nps-detractors .nps-segment-value { color: var(--danger); }
    
    /* ========== 交叉分析热力图 ========== */
    .cross-analysis-header {
        background: linear-gradient(135deg, var(--primary-light), var(--white));
        border-radius: var(--radius-lg);
        padding: 1.5rem;
        margin-bottom: 1.5rem;
        border: 1px solid var(--gray-200);
    }
    
    .cross-analysis-title {
        font-size: 1.25rem;
        font-weight: 700;
        color: var(--gray-900);
        margin-bottom: 0.5rem;
    }
    
    .cross-analysis-desc {
        font-size: 0.875rem;
        color: var(--gray-600);
    }
    
    /* ========== 用户原声卡片 ========== */
    .feedback-card {
        background: var(--white);
        border: 1px solid var(--gray-200);
        border-radius: var(--radius-lg);
        padding: 1.25rem;
        margin-bottom: 1rem;
        transition: all 0.2s ease;
    }
    
    .feedback-card:hover {
        border-color: var(--primary);
        box-shadow: var(--shadow-md);
    }
    
    .feedback-card.negative {
        border-left: 4px solid var(--danger);
    }
    
    .feedback-card.positive {
        border-left: 4px solid var(--success);
    }
    
    .feedback-header {
        display: flex;
        align-items: center;
        justify-content: space-between;
        margin-bottom: 0.75rem;
    }
    
    .feedback-sentiment {
        font-size: 1.25rem;
    }
    
    .feedback-user {
        font-size: 0.8rem;
        color: var(--gray-500);
    }
    
    .feedback-text {
        font-size: 0.9rem;
        color: var(--gray-800);
        line-height: 1.6;
        margin-bottom: 1rem;
    }
    
    .feedback-actions {
        display: flex;
        gap: 0.75rem;
    }
    
    /* ========== 洞察提示 ========== */
    .insight-box {
        background: linear-gradient(135deg, var(--info-light), var(--white));
        border: 1px solid var(--info);
        border-radius: var(--radius);
        padding: 0.875rem 1rem;
        margin-top: 0.75rem;
        font-size: 0.8rem;
        color: var(--gray-800);
    }
    
    .insight-box-icon {
        color: var(--info);
        margin-right: 0.5rem;
    }
    
    /* ========== 徽章 ========== */
    .badge {
        display: inline-flex;
        align-items: center;
        padding: 0.25rem 0.75rem;
        border-radius: 9999px;
        font-size: 0.7rem;
        font-weight: 600;
        letter-spacing: 0.02em;
    }
    
    .badge-primary {
        background: var(--primary);
        color: white;
    }
    
    .badge-success {
        background: var(--success-light);
        color: var(--success);
    }
    
    .badge-warning {
        background: var(--warning-light);
        color: var(--warning);
    }
    
    .badge-danger {
        background: var(--danger-light);
        color: var(--danger);
    }
    
    .badge-gray {
        background: var(--gray-100);
        color: var(--gray-600);
    }
    
    /* ========== 滚动条 ========== */
    ::-webkit-scrollbar {
        width: 8px;
        height: 8px;
    }
    
    ::-webkit-scrollbar-track {
        background: var(--gray-100);
        border-radius: 4px;
    }
    
    ::-webkit-scrollbar-thumb {
        background: var(--gray-300);
        border-radius: 4px;
    }
    
    ::-webkit-scrollbar-thumb:hover {
        background: var(--primary);
    }
    
    /* ========== Metric 组件样式覆盖 ========== */
    [data-testid="stMetricValue"] {
        font-size: 1.75rem;
        font-weight: 700;
        color: var(--primary);
    }
    
    [data-testid="stMetricDelta"] {
        font-size: 0.8rem;
    }
    
    /* ========== 分隔线 ========== */
    hr {
        margin: 1.5rem 0;
        border: none;
        height: 1px;
        background: var(--gray-200);
    }
    
    /* ========== 页面标题区 - shadcn 风格 ========== */
    .page-header {
        background: transparent;
        padding: 0 0 1.5rem;
        margin-bottom: 1.5rem;
        border-bottom: 1px solid var(--border);
    }
    
    .page-title {
        font-size: 1.875rem;
        font-weight: 700;
        color: var(--foreground);
        letter-spacing: -0.025em;
        line-height: 1.2;
    }
    
    .page-subtitle {
        font-size: 0.875rem;
        color: var(--muted-foreground);
        margin-top: 0.25rem;
    }
    
    /* ========== 统计网格 ========== */
    .stats-grid {
        display: grid;
        grid-template-columns: repeat(auto-fit, minmax(200px, 1fr));
        gap: 1rem;
        margin-bottom: 1.5rem;
    }
    
    .stat-item {
        background: var(--card);
        border: 1px solid var(--border);
        border-radius: var(--radius-lg);
        padding: 1.25rem;
        transition: all 0.2s ease;
    }
    
    .stat-item:hover {
        border-color: var(--green-200);
    }
    
    .stat-label {
        font-size: 0.7rem;
        color: var(--gray-500);
        text-transform: uppercase;
        letter-spacing: 0.08em;
        margin-bottom: 0.375rem;
        font-weight: 600;
    }
    
    .stat-value {
        font-size: 1.625rem;
        font-weight: 700;
        color: var(--black);
        font-family: 'JetBrains Mono', monospace;
    }
    
    /* ========== 文件上传区域 - 精致样式 ========== */
    [data-testid="stFileUploader"] {
        background: var(--white);
        border: 1.5px dashed var(--gray-300);
        border-radius: var(--radius);
        transition: all 0.2s ease;
    }
    
    [data-testid="stFileUploader"]:hover {
        border-color: var(--primary);
        background: var(--primary-light);
        border-style: dashed;
    }
    
    /* 上传区域内部样式 - 紧凑布局 */
    [data-testid="stFileUploader"] > div {
        padding: 0.75rem 0.625rem !important;
    }
    
    [data-testid="stFileUploader"] section {
        padding: 0 !important;
        gap: 0.25rem !important;
    }
    
    [data-testid="stFileUploader"] section > div {
        gap: 0.125rem !important;
    }
    
    /* 拖放区域文字 - 更小更精致 */
    [data-testid="stFileUploader"] [data-testid="stMarkdownContainer"] {
        margin: 0 !important;
    }
    
    [data-testid="stFileUploader"] [data-testid="stMarkdownContainer"] p {
        color: var(--gray-400) !important;
        font-size: 0.7rem !important;
        margin: 0.125rem 0 !important;
        line-height: 1.3 !important;
    }
    
    /* "Drag and drop file here" 主文字 - 缩小 */
    [data-testid="stFileUploader"] section > div:first-child,
    [data-testid="stFileUploader"] [data-testid="stMarkdownContainer"]:first-of-type p {
        font-size: 0.75rem !important;
        color: var(--gray-500) !important;
        font-weight: 400 !important;
    }
    
    /* 所有上传区域内的文字统一缩小 */
    [data-testid="stFileUploader"] span,
    [data-testid="stFileUploader"] p,
    [data-testid="stFileUploader"] div {
        font-size: 0.7rem !important;
    }
    
    /* 限制说明文字 - 更小 */
    [data-testid="stFileUploader"] section small,
    [data-testid="stFileUploader"] small {
        color: var(--gray-400) !important;
        font-size: 0.6rem !important;
        line-height: 1.2 !important;
    }
    
    /* Browse files 按钮 - 深色样式 */
    [data-testid="stFileUploader"] button,
    [data-testid="stFileUploader"] button[kind="secondary"] {
        background: #18181b !important;
        border: 1px solid #18181b !important;
        border-radius: 0.5rem !important;
        color: #ffffff !important;
        font-size: 0.75rem !important;
        font-weight: 500 !important;
        padding: 0.5rem 1rem !important;
        margin-top: 0.5rem !important;
        transition: all 0.15s ease !important;
        box-shadow: 0 1px 2px rgba(0,0,0,0.1) !important;
    }
    
    [data-testid="stFileUploader"] button:hover,
    [data-testid="stFileUploader"] button[kind="secondary"]:hover {
        background: #3f3f46 !important;
        border-color: #3f3f46 !important;
        color: #ffffff !important;
        box-shadow: 0 2px 4px rgba(0,0,0,0.15) !important;
    }
    
    /* 已上传文件列表 - 紧凑 */
    [data-testid="stFileUploader"] [data-testid="stFileUploaderFile"] {
        background: var(--white) !important;
        border: 1px solid var(--gray-200) !important;
        border-radius: var(--radius-sm) !important;
        padding: 0.375rem 0.625rem !important;
        margin-top: 0.5rem !important;
    }
    
    /* 文件图标和名称 */
    [data-testid="stFileUploader"] [data-testid="stFileUploaderFile"] span {
        color: var(--gray-700) !important;
        font-size: 0.7rem !important;
    }
    
    /* 文件大小 */
    [data-testid="stFileUploader"] [data-testid="stFileUploaderFile"] small {
        color: var(--gray-400) !important;
        font-size: 0.6rem !important;
    }
    
    /* 删除按钮 */
    [data-testid="stFileUploader"] [data-testid="stFileUploaderFile"] button {
        background: transparent !important;
        border: none !important;
        color: var(--gray-400) !important;
        padding: 0.125rem !important;
        min-width: auto !important;
        width: 1.25rem !important;
        height: 1.25rem !important;
    }
    
    [data-testid="stFileUploader"] [data-testid="stFileUploaderFile"] button:hover {
        color: var(--danger) !important;
        background: var(--danger-light) !important;
        border-radius: 50% !important;
    }
    
    /* 隐藏上传区域的标签文字 */
    [data-testid="stFileUploader"] label {
        font-size: 0.75rem !important;
        color: var(--gray-600) !important;
        font-weight: 500 !important;
        margin-bottom: 0.375rem !important;
    }
    
    /* 侧边栏内的上传区域特殊处理 */
    [data-testid="stSidebar"] [data-testid="stFileUploader"] {
        background: var(--white);
        border-color: var(--gray-200);
    }
    
    [data-testid="stSidebar"] [data-testid="stFileUploader"]:hover {
        border-color: var(--primary);
        background: var(--primary-light);
    }
    
    /* ========== 进度条 ========== */
    .stProgress > div > div > div > div {
        background: linear-gradient(90deg, var(--green-500), var(--green-400));
    }
    
    /* ========== 成功消息框 ========== */
    .success-box {
        background: var(--green-50);
        border: 1px solid var(--green-200);
        border-left: 4px solid var(--green-500);
        border-radius: var(--radius);
        padding: 1rem 1.25rem;
        margin: 1rem 0;
    }
    
    .success-box-title {
        color: var(--green-700);
        font-weight: 600;
        font-size: 0.9rem;
        margin-bottom: 0.25rem;
    }
    
    .success-box-text {
        color: var(--green-600);
        font-size: 0.85rem;
    }
    
    /* ========== 响应式 ========== */
    @media (max-width: 768px) {
        .main-title {
            font-size: 1.75rem;
        }
        
        .metric-card {
            padding: 1rem;
        }
        
        .metric-value {
            font-size: 1.5rem;
        }
    }
    
    /* ========== Plotly 图表优化 ========== */
    .plotly .modebar {
        top: 0.5rem !important;
        right: 0.5rem !important;
    }
    
    .plotly .modebar-btn {
        font-size: 14px !important;
    }
</style>
""", unsafe_allow_html=True)

# 解决 Matplotlib 中文乱码
plt.rcParams['font.sans-serif'] = ['SimHei', 'Arial Unicode MS', 'sans-serif'] 
plt.rcParams['axes.unicode_minus'] = False

# ========== 导出功能函数 ==========

def fig_to_image_bytes(fig, format='png', width=800, height=500):
    """将 Plotly 图表转换为图片字节"""
    if not KALEIDO_AVAILABLE:
        st.warning("⚠️ kaleido 库未安装，无法导出图表图片")
        return None
    
    try:
        # 确保图表有白色背景
        fig.update_layout(
            plot_bgcolor='white',
            paper_bgcolor='white'
        )
        img_bytes = fig.to_image(format=format, width=width, height=height, scale=2, engine="kaleido")
        if img_bytes and len(img_bytes) > 0:
            return img_bytes
        else:
            return None
    except Exception as e:
        st.warning(f"⚠️ 图表导出失败: {str(e)}")
        return None

def create_word_report(title, sections, include_charts=True):
    """创建 Word 报告"""
    doc = Document()
    
    # 设置标题
    title_para = doc.add_heading(title, 0)
    title_para.alignment = WD_ALIGN_PARAGRAPH.CENTER
    
    # 添加生成时间
    doc.add_paragraph(f"生成时间: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}")
    doc.add_paragraph("")
    
    for section in sections:
        # 添加问题标题
        doc.add_heading(section.get('question', ''), level=1)
        
        # 添加统计摘要
        if 'summary' in section:
            doc.add_paragraph(section['summary'])
        
        # 添加数据表格
        if 'data' in section and section['data'] is not None:
            df = section['data']
            if len(df) > 0:
                table = doc.add_table(rows=len(df) + 1, cols=len(df.columns))
                table.style = 'Table Grid'
                
                # 表头
                for j, col in enumerate(df.columns):
                    table.rows[0].cells[j].text = str(col)
                    table.rows[0].cells[j].paragraphs[0].runs[0].bold = True
                
                # 数据行
                for i, row in df.iterrows():
                    for j, val in enumerate(row):
                        table.rows[i + 1].cells[j].text = str(val)
        
        # 添加图表图片
        if include_charts and 'chart_image' in section and section['chart_image'] is not None:
            doc.add_paragraph("")
            img_stream = BytesIO(section['chart_image'])
            doc.add_picture(img_stream, width=Inches(6))
        
        doc.add_paragraph("")
    
    # 保存到字节流
    doc_bytes = BytesIO()
    doc.save(doc_bytes)
    doc_bytes.seek(0)
    return doc_bytes

def create_ppt_report(title, sections, include_charts=True):
    """创建 PPT 报告"""
    if not PPTX_AVAILABLE:
        return None
    
    prs = Presentation()
    prs.slide_width = PptxInches(13.333)
    prs.slide_height = PptxInches(7.5)
    
    # 标题页
    title_slide_layout = prs.slide_layouts[6]  # 空白布局
    slide = prs.slides.add_slide(title_slide_layout)
    
    # 添加标题文本框
    left = PptxInches(0.5)
    top = PptxInches(2.5)
    width = PptxInches(12.333)
    height = PptxInches(1.5)
    
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = PptxPt(44)
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    
    # 添加副标题
    top2 = PptxInches(4)
    txBox2 = slide.shapes.add_textbox(left, top2, width, PptxInches(0.5))
    tf2 = txBox2.text_frame
    p2 = tf2.paragraphs[0]
    p2.text = f"生成时间: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}"
    p2.font.size = PptxPt(18)
    p2.font.color.rgb = PptxRGBColor(128, 128, 128)
    p2.alignment = PP_ALIGN.CENTER
    
    # 为每个问题创建幻灯片
    for section in sections:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        
        # 问题标题
        txBox = slide.shapes.add_textbox(PptxInches(0.5), PptxInches(0.3), PptxInches(12.333), PptxInches(0.8))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = section.get('question', '')[:80]  # 限制长度
        p.font.size = PptxPt(24)
        p.font.bold = True
        
        # 添加图表 - 占据更大空间
        if include_charts and 'chart_image' in section and section['chart_image'] is not None:
            img_stream = BytesIO(section['chart_image'])
            # 图表占据整个幻灯片宽度，更美观
            slide.shapes.add_picture(img_stream, PptxInches(0.5), PptxInches(1.2), width=PptxInches(12.333))
        
        # 数据摘要改为右上角小标签形式（英文避免乱码）
        if 'summary' in section:
            # 将中文摘要转换为英文格式
            summary_text = section['summary']
            # 替换中文为英文
            summary_text = summary_text.replace('总样本数:', 'Total:')
            summary_text = summary_text.replace('总样本:', 'Total:')
            summary_text = summary_text.replace('最常见选项:', 'Top:')
            summary_text = summary_text.replace('最常见:', 'Top:')
            summary_text = summary_text.replace('数量:', 'Count:')
            summary_text = summary_text.replace('唯一值数量:', 'Unique:')
            summary_text = summary_text.replace('唯一值:', 'Unique:')
            summary_text = summary_text.replace('种', '')
            
            # 右上角添加简洁的数据标签
            txBox = slide.shapes.add_textbox(PptxInches(9.5), PptxInches(0.3), PptxInches(3.5), PptxInches(0.6))
            tf = txBox.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            # 只显示核心数据，一行展示
            summary_parts = summary_text.split(',')
            p.text = ', '.join(summary_parts[:2]) if len(summary_parts) > 1 else summary_text[:60]
            p.font.size = PptxPt(10)
            p.font.color.rgb = PptxRGBColor(128, 128, 128)
            p.alignment = PP_ALIGN.RIGHT
    
    # 保存到字节流
    ppt_bytes = BytesIO()
    prs.save(ppt_bytes)
    ppt_bytes.seek(0)
    return ppt_bytes

def create_pdf_report(title, sections, include_charts=True):
    """创建 PDF 报告"""
    if not PDF_AVAILABLE:
        return None
    
    pdf_bytes = BytesIO()
    doc = SimpleDocTemplate(pdf_bytes, pagesize=A4, topMargin=0.5*inch, bottomMargin=0.5*inch)
    
    # 获取样式
    styles = getSampleStyleSheet()
    
    # 自定义样式
    title_style = ParagraphStyle(
        'CustomTitle',
        parent=styles['Heading1'],
        fontSize=24,
        spaceAfter=20,
        alignment=1  # 居中
    )
    
    heading_style = ParagraphStyle(
        'CustomHeading',
        parent=styles['Heading2'],
        fontSize=14,
        spaceAfter=10,
        spaceBefore=15
    )
    
    normal_style = ParagraphStyle(
        'CustomNormal',
        parent=styles['Normal'],
        fontSize=10,
        spaceAfter=8
    )
    
    story = []
    
    # 标题
    story.append(Paragraph(title, title_style))
    story.append(Paragraph(f"生成时间: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}", normal_style))
    story.append(Spacer(1, 20))
    
    for section in sections:
        # 问题标题
        question_text = section.get('question', '')[:100]
        story.append(Paragraph(question_text, heading_style))
        
        # 摘要
        if 'summary' in section:
            story.append(Paragraph(section['summary'], normal_style))
        
        # 数据表格
        if 'data' in section and section['data'] is not None:
            df = section['data']
            if len(df) > 0 and len(df) <= 20:  # 限制行数
                table_data = [df.columns.tolist()] + df.values.tolist()
                t = Table(table_data)
                t.setStyle(TableStyle([
                    ('BACKGROUND', (0, 0), (-1, 0), HexColor('#374151')),
                    ('TEXTCOLOR', (0, 0), (-1, 0), HexColor('#ffffff')),
                    ('ALIGN', (0, 0), (-1, -1), 'CENTER'),
                    ('FONTSIZE', (0, 0), (-1, -1), 8),
                    ('BOTTOMPADDING', (0, 0), (-1, 0), 8),
                    ('BACKGROUND', (0, 1), (-1, -1), HexColor('#f9fafb')),
                    ('GRID', (0, 0), (-1, -1), 1, HexColor('#e5e7eb'))
                ]))
                story.append(t)
        
        # 图表图片
        if include_charts and 'chart_image' in section and section['chart_image'] is not None:
            img_stream = BytesIO(section['chart_image'])
            img = Image(img_stream, width=5*inch, height=3*inch)
            story.append(Spacer(1, 10))
            story.append(img)
        
        story.append(Spacer(1, 20))
    
    doc.build(story)
    pdf_bytes.seek(0)
    return pdf_bytes

def get_export_data_for_question(question_name, df, value_counts_df, fig=None):
    """为单个问题准备导出数据"""
    section = {
        'question': question_name,
        'data': value_counts_df.copy() if value_counts_df is not None else None,
        'summary': '',
        'chart_image': None
    }
    
    # 生成摘要
    if value_counts_df is not None and len(value_counts_df) > 0:
        total = value_counts_df['count'].sum()
        top_item = value_counts_df.iloc[0]
        top_name = str(top_item.iloc[0])[:30]
        top_count = top_item['count']
        top_pct = (top_count / total * 100)
        
        section['summary'] = f"Total: {total}, Top: {top_name}, Count: {top_count} ({top_pct:.1f}%), Unique: {len(value_counts_df)}"
    
    # 生成图表图片
    if fig is not None and KALEIDO_AVAILABLE:
        try:
            section['chart_image'] = fig_to_image_bytes(fig)
        except:
            pass
    
    return section

# 初始化导出数据存储
if 'export_sections' not in st.session_state:
    st.session_state.export_sections = []

if 'analyzed_data' not in st.session_state:
    st.session_state.analyzed_data = {}

# ========== 长标签换行函数 ==========
def wrap_label(text, max_len=20):
    """将长标签换行显示"""
    text = str(text)
    if len(text) <= max_len:
        return text
    words = text.split(' ')
    lines = []
    current = ""
    for word in words:
        if len(current) + len(word) + 1 > max_len:
            if current:
                lines.append(current)
            current = word
        else:
            current = current + " " + word if current else word
    if current:
        lines.append(current)
    return "<br>".join(lines[:2]) + ("..." if len(lines) > 2 else "")

# ========== 数据缓存函数 - 防止闪烁 ==========
@st.cache_data(show_spinner=False)
def load_csv_data(file_content, file_name):
    """缓存 CSV 数据加载，避免重复读取"""
    import io
    lines = file_content.split('\n')
    
    # 检测真正的表头行
    header_row = 0
    for i, line in enumerate(lines[:10]):
        comma_count = line.count(',')
        if comma_count >= 5 and not line.strip().endswith(':'):
            header_row = i
            break
    
    # 读取数据
    df = pd.read_csv(io.StringIO(file_content), skiprows=header_row, on_bad_lines='skip')
    
    # 清理空列和空行
    df = df.dropna(how='all', axis=1)
    df = df.dropna(how='all', axis=0)
    
    return df

@st.cache_data(show_spinner=False)
def load_excel_data(file_content, file_name):
    """缓存 Excel 数据加载"""
    import io
    return pd.read_excel(io.BytesIO(file_content))

@st.cache_data(show_spinner=False)
def compute_basic_stats(df_hash, total_rows, total_cols, null_count):
    """缓存基础统计计算"""
    completeness = (1 - null_count / (total_rows * total_cols)) * 100 if total_rows * total_cols > 0 else 0
    return completeness

def add_export_section(section):
    """添加导出章节"""
    existing_questions = [s['question'] for s in st.session_state.export_sections]
    if section['question'] not in existing_questions:
        st.session_state.export_sections.append(section)

def clear_export_sections():
    """清空导出章节"""
    st.session_state.export_sections = []

def generate_all_export_sections(df, selected_columns):
    """一键生成所有选中变量的导出数据"""
    sections = []
    
    for col in selected_columns:
        if col not in df.columns:
            continue
            
        section = {
            'question': col,
            'data': None,
            'summary': '',
            'chart_image': None
        }
        
        is_numeric = pd.api.types.is_numeric_dtype(df[col])
        unique_values = df[col].nunique()
        
        # 判断是否为评分型数据
        is_rating_data = False
        if is_numeric and unique_values <= 10:
            non_null_values = df[col].dropna()
            if len(non_null_values) > 0:
                is_rating_data = (non_null_values == non_null_values.astype(int)).all()
        
        if is_rating_data:
            # 评分型数据
            mean_val = df[col].mean()
            median_val = df[col].median()
            min_val = int(df[col].min())
            max_val = int(df[col].max())
            total_responses = df[col].count()
            
            # 统计各评分的频次
            rating_counts = df[col].value_counts().sort_index()
            all_ratings = list(range(min_val, max_val + 1))
            rating_df = pd.DataFrame({
                '评分': all_ratings,
                '人数': [rating_counts.get(r, 0) for r in all_ratings],
                '占比': [(rating_counts.get(r, 0) / total_responses * 100) for r in all_ratings]
            })
            rating_df['占比'] = rating_df['占比'].round(1).astype(str) + '%'
            
            section['data'] = rating_df
            
            # 计算满意度指标
            high_score_count = df[col][df[col] >= (max_val - 1)].count()
            high_score_pct = (high_score_count / total_responses * 100) if total_responses > 0 else 0
            
            section['summary'] = f"平均分: {mean_val:.2f}/{max_val}, 中位数: {median_val:.1f}, 高分率(≥{max_val-1}分): {high_score_pct:.1f}%, 样本数: {total_responses}"
            
            # 生成评分分布图
            try:
                colors = ['#ef4444', '#f97316', '#eab308', '#84cc16', '#22c55e']
                if max_val - min_val + 1 <= len(colors):
                    bar_colors = colors[-(max_val - min_val + 1):]
                else:
                    bar_colors = px.colors.qualitative.Set2[:len(all_ratings)]
                
                export_fig = go.Figure()
                rating_counts_list = [rating_counts.get(r, 0) for r in all_ratings]
                rating_pcts = [(c / total_responses * 100) for c in rating_counts_list]
                
                for i, rating in enumerate(all_ratings):
                    export_fig.add_trace(go.Bar(
                        x=[str(rating)],
                        y=[rating_counts_list[i]],
                        marker_color=bar_colors[i % len(bar_colors)],
                        text=f"{rating_counts_list[i]}<br>({rating_pcts[i]:.1f}%)",
                        textposition='outside'
                    ))
                
                export_fig.update_layout(
                    title=f"⭐ {col} 评分分布 (1-{max_val}分) - 平均分: {mean_val:.2f}",
                    xaxis_title="评分",
                    yaxis_title="人数",
                    showlegend=False,
                    height=400, width=700,
                    plot_bgcolor='white', paper_bgcolor='white',
                    bargap=0.3
                )
                img_bytes = fig_to_image_bytes(export_fig, width=700, height=400)
                if img_bytes:
                    section['chart_image'] = img_bytes
            except Exception as e:
                pass
                
        elif is_numeric:
            # 纯数值型变量（连续变量）
            mean_val = df[col].mean()
            median_val = df[col].median()
            std_val = df[col].std()
            min_val = df[col].min()
            max_val = df[col].max()
            
            stats_df = pd.DataFrame({
                '统计量': ['均值', '中位数', '标准差', '最小值', '最大值', '样本数'],
                '数值': [f"{mean_val:.2f}", f"{median_val:.2f}", f"{std_val:.2f}", 
                        f"{min_val:.2f}", f"{max_val:.2f}", str(df[col].count())]
            })
            section['data'] = stats_df
            section['summary'] = f"均值: {mean_val:.2f}, 中位数: {median_val:.2f}, 标准差: {std_val:.2f}, 范围: {min_val:.2f} ~ {max_val:.2f}"
            
            # 生成图表
            try:
                from plotly.subplots import make_subplots
                export_fig = make_subplots(rows=1, cols=2, subplot_titles=('分布直方图', '箱线图'))
                export_fig.add_trace(go.Histogram(x=df[col].dropna(), marker_color='#667eea', nbinsx=20), row=1, col=1)
                export_fig.add_trace(go.Box(y=df[col].dropna(), marker_color='#764ba2'), row=1, col=2)
                export_fig.update_layout(title=f"{col} 数值分布", showlegend=False, height=400, width=800,
                                        plot_bgcolor='white', paper_bgcolor='white')
                img_bytes = fig_to_image_bytes(export_fig, width=800, height=400)
                if img_bytes:
                    section['chart_image'] = img_bytes
            except Exception as e:
                pass
        else:
            # 类别型变量
            value_counts = df[col].value_counts()
            total = len(df[col].dropna())
            
            vc_df = value_counts.reset_index()
            vc_df.columns = [col, 'count']
            vc_df['percentage'] = (vc_df['count'] / total * 100).round(1)
            
            section['data'] = vc_df.head(15)
            
            top_val = value_counts.index[0] if len(value_counts) > 0 else "N/A"
            top_count = value_counts.values[0] if len(value_counts) > 0 else 0
            top_pct = (top_count / total * 100) if total > 0 else 0
            
            section['summary'] = f"Total: {total}, Top: {str(top_val)[:30]}, Count: {top_count} ({top_pct:.1f}%), Unique: {len(value_counts)}"
            
            # 生成饼图
            try:
                pie_data = vc_df.head(8).copy()
                if len(vc_df) > 8:
                    other_count = vc_df.iloc[8:]['count'].sum()
                    other_row = pd.DataFrame({col: ['其他'], 'count': [other_count], 'percentage': [(other_count/total*100)]})
                    pie_data = pd.concat([pie_data, other_row], ignore_index=True)
                
                pie_data['label'] = pie_data[col].apply(lambda x: str(x)[:20] + "..." if len(str(x)) > 20 else str(x))
                pie_data['text'] = pie_data['percentage'].apply(lambda x: f"{x:.1f}%")
                
                colors = px.colors.qualitative.Set2[:len(pie_data)]
                export_fig = go.Figure(data=[go.Pie(
                    labels=pie_data['label'],
                    values=pie_data['count'],
                    hole=0.4,
                    text=pie_data['text'],
                    textposition='inside',
                    textinfo='text',
                    textfont=dict(size=11, color='white'),
                    marker=dict(colors=colors, line=dict(color='white', width=2))
                )])
                export_fig.update_layout(
                    title=f"{col} 分布",
                    showlegend=True,
                    legend=dict(orientation="h", yanchor="top", y=-0.15, xanchor="center", x=0.5, font=dict(size=9)),
                    height=450, width=650,
                    plot_bgcolor='white', paper_bgcolor='white',
                    margin=dict(t=60, b=120, l=20, r=20)
                )
                img_bytes = fig_to_image_bytes(export_fig, width=650, height=450)
                if img_bytes:
                    section['chart_image'] = img_bytes
            except Exception as e:
                pass
        
        sections.append(section)
    
    return sections

# --- 侧边栏 - 专业 SaaS 风格 ---
with st.sidebar:
    # 数据上传区
    st.caption("📁 数据上传")
    uploaded_file = st.file_uploader(
        "上传调研数据 (Excel/CSV)", 
        type=["csv", "xlsx", "pdf", "docx"],
        help="支持格式: CSV, Excel, PDF, Word | 最大 200MB"
    )
    
    if uploaded_file:
        st.markdown(f"""
        <div class="upload-success">
            <span class="upload-success-icon">✓</span>
            <span class="upload-success-text">文件已加载</span>
            <div class="upload-file-info">
                📄 {uploaded_file.name}<br>
                💾 {uploaded_file.size / 1024:.1f} KB
            </div>
        </div>
        """, unsafe_allow_html=True)
    
    # ========== 数据保存/加载功能 ==========
    with st.expander("💾 保存/加载数据", expanded=False):
        st.caption("保存当前数据，下次无需重新上传")
        
        # 获取已保存的文件列表
        saved_files = list(SAVE_DIR.glob("*.pkl"))
        saved_names = [f.stem for f in saved_files]
        
        # 保存当前数据
        st.markdown("**保存当前数据**")
        save_name = st.text_input(
            "保存名称",
            placeholder="例如：用户满意度调研_2024Q1",
            key="save_name_input",
            label_visibility="collapsed"
        )
        
        if st.button("💾 保存数据", key="save_data_btn", use_container_width=True):
            # 检查是否有数据可保存
            has_data = False
            save_data = {}
            
            # 检查上传的文件数据
            if uploaded_file:
                file_type = uploaded_file.name.split('.')[-1].lower()
                if file_type in ['csv', 'xlsx']:
                    uploaded_file.seek(0)
                    if file_type == 'csv':
                        df_to_save = pd.read_csv(uploaded_file)
                    else:
                        df_to_save = pd.read_excel(uploaded_file)
                    save_data['df'] = df_to_save
                    save_data['source'] = 'file'
                    save_data['filename'] = uploaded_file.name
                    has_data = True
            
            # 检查 URL 获取的数据
            if 'url_df' in st.session_state and st.session_state['url_df'] is not None:
                save_data['df'] = st.session_state['url_df']
                save_data['source'] = 'url'
                has_data = True
            
            # 保存问题映射
            if st.session_state.get('question_map'):
                save_data['question_map'] = st.session_state['question_map']
            
            if has_data and save_name:
                save_data['saved_at'] = datetime.now().strftime('%Y-%m-%d %H:%M:%S')
                save_path = SAVE_DIR / f"{save_name}.pkl"
                
                try:
                    with open(save_path, 'wb') as f:
                        pickle.dump(save_data, f)
                    st.success(f"✅ 数据已保存为: {save_name}")
                    st.rerun()
                except Exception as e:
                    st.error(f"❌ 保存失败: {e}")
            elif not save_name:
                st.warning("⚠️ 请输入保存名称")
            else:
                st.warning("⚠️ 没有可保存的数据，请先上传文件或获取数据")
        
        st.markdown("---")
        
        # 加载已保存的数据
        st.markdown("**加载已保存数据**")
        
        if saved_names:
            # 显示已保存的数据列表
            selected_save = st.selectbox(
                "选择数据",
                options=saved_names,
                key="load_data_select",
                label_visibility="collapsed"
            )
            
            # 显示保存信息
            if selected_save:
                save_path = SAVE_DIR / f"{selected_save}.pkl"
                try:
                    with open(save_path, 'rb') as f:
                        preview_data = pickle.load(f)
                    
                    saved_at = preview_data.get('saved_at', '未知')
                    source = preview_data.get('source', '未知')
                    row_count = len(preview_data.get('df', pd.DataFrame()))
                    
                    st.caption(f"📅 保存时间: {saved_at}")
                    st.caption(f"📊 数据量: {row_count} 条")
                except:
                    pass
            
            col_load, col_delete = st.columns([2, 1])
            
            with col_load:
                if st.button("📂 加载数据", key="load_data_btn", use_container_width=True):
                    save_path = SAVE_DIR / f"{selected_save}.pkl"
                    try:
                        with open(save_path, 'rb') as f:
                            loaded_data = pickle.load(f)
                        
                        # 加载数据到 session_state
                        if 'df' in loaded_data:
                            st.session_state['saved_df'] = loaded_data['df']
                            st.session_state['saved_df_name'] = selected_save
                        
                        if 'question_map' in loaded_data:
                            st.session_state['question_map'] = loaded_data['question_map']
                        
                        st.success(f"✅ 已加载: {selected_save}")
                        st.rerun()
                    except Exception as e:
                        st.error(f"❌ 加载失败: {e}")
            
            with col_delete:
                if st.button("🗑️", key="delete_data_btn", help="删除此保存"):
                    save_path = SAVE_DIR / f"{selected_save}.pkl"
                    try:
                        os.remove(save_path)
                        st.success("已删除")
                        st.rerun()
                    except Exception as e:
                        st.error(f"删除失败: {e}")
        else:
            st.info("💡 暂无已保存的数据")
        
        # 显示当前加载状态
        if 'saved_df' in st.session_state and st.session_state.get('saved_df') is not None:
            st.markdown("---")
            st.markdown(f"""
            <div style="background: rgba(255,255,255,0.1); padding: 0.75rem; border-radius: 8px; border: 1px solid rgba(255,255,255,0.2);">
                <div style="font-size: 0.75rem; color: rgba(255,255,255,0.7);">当前已加载</div>
                <div style="font-size: 0.9rem; font-weight: 600; color: white;">📊 {st.session_state.get('saved_df_name', '未命名')}</div>
                <div style="font-size: 0.75rem; color: rgba(255,255,255,0.6);">{len(st.session_state['saved_df'])} 条数据</div>
            </div>
            """, unsafe_allow_html=True)
            
            if st.button("❌ 清除已加载数据", key="clear_saved_df", use_container_width=True):
                del st.session_state['saved_df']
                if 'saved_df_name' in st.session_state:
                    del st.session_state['saved_df_name']
                st.rerun()
    
    # CSV 链接获取数据
    with st.expander("🔗 或：通过链接获取", expanded=False):
        st.caption("粘贴 Ptengine CSV 报告链接")
        csv_url = st.text_input("CSV链接", placeholder="https://ecbi.ptengine.com/public/question/xxx.csv", key="csv_url_input", label_visibility="collapsed")
        if csv_url and ".csv" in csv_url:
            if st.button("🔄 获取数据", key="fetch_csv_btn", use_container_width=True):
                try:
                    resp = requests.get(csv_url, timeout=30)
                    if resp.status_code == 200:
                        df_from_url = pd.read_csv(StringIO(resp.text))
                        if len(df_from_url) > 0:
                            st.session_state['url_df'] = df_from_url
                            st.success(f"✅ 成功获取 {len(df_from_url)} 条数据！")
                except Exception as e:
                    st.error(f"❌ {str(e)}")
        if 'url_df' in st.session_state:
            st.success(f"✅ 已加载 {len(st.session_state['url_df'])} 条数据")
    
    # 问题映射配置
    with st.expander("🔗 问题映射", expanded=False):
        import re as re_module
        
        st.markdown("**① 粘贴表单链接**")
        form_url = st.text_input(
            "链接", placeholder="https://comp.ptengine.com/assets/xxx/latest/index.html",
            key="form_url_input", label_visibility="collapsed"
        )
        
        if form_url:
            match = re_module.search(r'/assets/([^/]+)/', form_url)
            if match:
                config_url = f"https://comp.ptengine.com/assets/{match.group(1)}/latest/config.json"
                st.markdown("**② 打开链接，Ctrl+A 全选，Ctrl+C 复制**")
                st.markdown(f"[🔗 点击打开 config.json]({config_url})")
        
        st.markdown("**③ 粘贴完整内容**")
        st.caption("⚠️ 必须 Ctrl+A 全选后复制，不要只复制部分！")
        config_text = st.text_area(
            "内容", height=80, key="config_text_input", 
            label_visibility="collapsed", placeholder="粘贴完整的 JSON 内容..."
        )
        
        # 添加解析按钮
        if st.button("🔍 解析配置", key="parse_config_btn", use_container_width=True):
            if config_text and len(config_text) > 100:
                question_map = {}
                
                try:
                    # 方法1：尝试解析 JSON
                    import json
                    data = json.loads(config_text)
                    
                    # 递归提取
                    def extract(obj):
                        if isinstance(obj, dict):
                            name = obj.get('name', '')
                            question = obj.get('question', '')
                            if name and question and not name.startswith('$') and not name.startswith('表单页'):
                                question_map[name] = question
                            for v in obj.values():
                                extract(v)
                        elif isinstance(obj, list):
                            for item in obj:
                                extract(item)
                    
                    extract(data)
                except:
                    # 方法2：正则提取（支持格式化JSON）
                    names = re_module.findall(r'"name"\s*:\s*"([^"]+)"', config_text)
                    questions = re_module.findall(r'"question"\s*:\s*"([^"]+)"', config_text)
                    
                    # 按顺序配对（假设 name 和 question 是成对出现的）
                    valid_names = [n for n in names if not n.startswith('$') and not n.startswith('表单页') and not n.startswith('结束页')]
                    
                    for i, name in enumerate(valid_names):
                        if i < len(questions):
                            question_map[name] = questions[i]
                
                if question_map:
                    st.session_state['question_map'] = question_map
                    st.success(f"✅ 成功加载 {len(question_map)} 个问题映射！")
                else:
                    st.error("❌ 未找到问题。请检查内容是否完整。")
            else:
                st.warning("⚠️ 内容太短，请确保完整复制")
        
        # 显示当前映射
        if st.session_state.get('question_map'):
            st.success(f"📋 当前已加载 {len(st.session_state['question_map'])} 个映射")
            with st.expander("查看映射详情"):
                for k, v in st.session_state['question_map'].items():
                    st.caption(f"**{k}**: {v[:60]}...")
            if st.button("🗑️ 清除映射", key="clear_map"):
                st.session_state['question_map'] = {}
    
    # ========== 简洁浅色导航栏 - 参考 GOODFOOD 风格 ==========
    st.markdown("""
    <style>
    /* 浅色主题侧边栏 - 2种主色：浅灰背景 + 紫色强调 */
    [data-testid="stSidebar"] {
        background: #F8F9FC !important;
    }
    
    [data-testid="stSidebar"] > div:first-child {
        background: transparent !important;
    }
    
    /* 侧边栏内所有文字默认深灰色 */
    [data-testid="stSidebar"] * {
        color: #64748b !important;
    }
    
    [data-testid="stSidebar"] label {
        color: #94a3b8 !important;
    }
    
    /* Logo 区域 */
    .sidebar-logo {
        display: flex;
        align-items: center;
        gap: 0.75rem;
        padding: 1rem 0.5rem;
        margin-bottom: 1.5rem;
    }
    .sidebar-logo-icon {
        width: 36px;
        height: 36px;
        background: linear-gradient(135deg, #818cf8 0%, #6366f1 100%);
        border-radius: 10px;
        display: flex;
        align-items: center;
        justify-content: center;
        font-size: 1.25rem;
    }
    .sidebar-logo-text {
        font-size: 1.1rem;
        font-weight: 700;
        color: #6366f1 !important;
        letter-spacing: -0.5px;
    }
    
    /* 分组标题 */
    .nav-section-title {
        font-size: 0.7rem;
        font-weight: 600;
        color: #64748b !important;
        text-transform: uppercase;
        letter-spacing: 0.5px;
        padding: 0.75rem 0.75rem 0.5rem;
        margin-top: 0.5rem;
        margin-bottom: 0.25rem;
    }
    
    /* 导航按钮默认样式 */
    [data-testid="stSidebar"] .stButton > button {
        background: transparent !important;
        border: none !important;
        color: #64748b !important;
        padding: 0.6rem 0.75rem !important;
        font-size: 0.875rem !important;
        font-weight: 500 !important;
        border-radius: 8px !important;
        transition: all 0.15s ease !important;
        text-align: left !important;
        justify-content: flex-start !important;
    }
    
    [data-testid="stSidebar"] .stButton > button:hover {
        background: #f1f5f9 !important;
        color: #475569 !important;
    }
    
    /* 选中状态的导航按钮 - 深蓝色背景，高对比度 */
    [data-testid="stSidebar"] .stButton > button[kind="primary"] {
        background: #4f46e5 !important;
        color: #ffffff !important;
        font-weight: 600 !important;
        box-shadow: 0 2px 8px rgba(79, 70, 229, 0.25) !important;
        border: none !important;
    }
    
    [data-testid="stSidebar"] .stButton > button[kind="primary"]:hover {
        background: #4338ca !important;
        color: #ffffff !important;
    }
    
    /* 侧边栏 selectbox 样式 */
    [data-testid="stSidebar"] .stSelectbox > div > div {
        background: white !important;
        border: 1px solid #e2e8f0 !important;
        border-radius: 8px !important;
        color: #475569 !important;
    }
    
    [data-testid="stSidebar"] .stSelectbox > div > div > div {
        color: #475569 !important;
    }
    
    [data-testid="stSidebar"] .stSelectbox svg {
        fill: #94a3b8 !important;
    }
    
    [data-testid="stSidebar"] .stSelectbox > div > div:hover {
        border-color: #6366f1 !important;
    }
    
    /* 侧边栏 expander 样式 */
    [data-testid="stSidebar"] [data-testid="stExpander"] {
        background: white !important;
        border: 1px solid #e2e8f0 !important;
        border-radius: 8px !important;
    }
    
    [data-testid="stSidebar"] [data-testid="stExpander"] > details {
        background: transparent !important;
    }
    
    [data-testid="stSidebar"] [data-testid="stExpander"] > details > summary {
        background: transparent !important;
        color: #475569 !important;
        font-weight: 500 !important;
    }
    
    [data-testid="stSidebar"] [data-testid="stExpander"] > details > summary > span,
    [data-testid="stSidebar"] [data-testid="stExpander"] > details > summary > span > span,
    [data-testid="stSidebar"] [data-testid="stExpander"] > details > summary p {
        color: #475569 !important;
    }
    
    /* expander 箭头图标 */
    [data-testid="stSidebar"] [data-testid="stExpander"] svg {
        fill: #94a3b8 !important;
        stroke: #94a3b8 !important;
    }
    
    [data-testid="stSidebar"] .streamlit-expanderHeader {
        background: transparent !important;
        color: #475569 !important;
    }
    
    [data-testid="stSidebar"] .streamlit-expanderContent {
        background: transparent !important;
        padding: 0.75rem !important;
    }
    
    /* 侧边栏 expander 内的所有文字 */
    [data-testid="stSidebar"] .streamlit-expanderContent p,
    [data-testid="stSidebar"] .streamlit-expanderContent span,
    [data-testid="stSidebar"] .streamlit-expanderContent label,
    [data-testid="stSidebar"] .streamlit-expanderContent div,
    [data-testid="stSidebar"] [data-testid="stExpander"] p,
    [data-testid="stSidebar"] [data-testid="stExpander"] span,
    [data-testid="stSidebar"] [data-testid="stExpander"] label,
    [data-testid="stSidebar"] [data-testid="stExpander"] div {
        color: #475569 !important;
    }
    
    [data-testid="stSidebar"] .streamlit-expanderContent strong,
    [data-testid="stSidebar"] [data-testid="stExpander"] strong {
        color: #6366f1 !important;
    }
    
    /* 侧边栏输入框样式 */
    [data-testid="stSidebar"] .stTextInput > div > div > input,
    [data-testid="stSidebar"] .stTextInput input,
    [data-testid="stSidebar"] input[type="text"],
    [data-testid="stSidebar"] [data-baseweb="input"] input {
        background: white !important;
        background-color: white !important;
        border: 1px solid #e2e8f0 !important;
        color: #1e293b !important;
        border-radius: 8px !important;
        -webkit-text-fill-color: #1e293b !important;
        caret-color: #6366f1 !important;
    }
    
    [data-testid="stSidebar"] .stTextInput > div > div > input::placeholder,
    [data-testid="stSidebar"] input::placeholder {
        color: #94a3b8 !important;
        -webkit-text-fill-color: #94a3b8 !important;
    }
    
    [data-testid="stSidebar"] .stTextInput > div > div > input:focus,
    [data-testid="stSidebar"] input:focus {
        border-color: #6366f1 !important;
        box-shadow: 0 0 0 2px rgba(99, 102, 241, 0.1) !important;
    }
    
    /* 侧边栏 text_area 样式 */
    [data-testid="stSidebar"] .stTextArea textarea,
    [data-testid="stSidebar"] textarea {
        background: white !important;
        background-color: white !important;
        color: #1e293b !important;
        border: 1px solid #e2e8f0 !important;
        -webkit-text-fill-color: #1e293b !important;
        caret-color: #6366f1 !important;
        border-radius: 8px !important;
    }
    
    [data-testid="stSidebar"] .stTextArea textarea::placeholder,
    [data-testid="stSidebar"] textarea::placeholder {
        color: #94a3b8 !important;
        -webkit-text-fill-color: #94a3b8 !important;
    }
    
    /* 自动补全下拉框样式 */
    [data-testid="stSidebar"] [data-baseweb="popover"] {
        background: white !important;
        border: 1px solid #e2e8f0 !important;
        border-radius: 8px !important;
    }
    
    [data-testid="stSidebar"] [data-baseweb="popover"] li,
    [data-testid="stSidebar"] [data-baseweb="menu"] li {
        color: #475569 !important;
    }
    
    [data-testid="stSidebar"] [data-baseweb="popover"] li:hover,
    [data-testid="stSidebar"] [data-baseweb="menu"] li:hover {
        background: #f1f5f9 !important;
    }
    
    /* 侧边栏 caption 样式 */
    [data-testid="stSidebar"] .stCaption {
        color: #94a3b8 !important;
    }
    
    /* 侧边栏 markdown 文字 */
    [data-testid="stSidebar"] .stMarkdown p {
        color: #64748b !important;
    }
    
    /* 上传成功提示 */
    [data-testid="stSidebar"] .stSuccess {
        background: #f0fdf4 !important;
        border: 1px solid #bbf7d0 !important;
        color: #166534 !important;
    }
    
    /* 文件上传组件样式 */
    [data-testid="stSidebar"] [data-testid="stFileUploader"] {
        background: white !important;
        border: 1px dashed #e2e8f0 !important;
        border-radius: 8px !important;
        padding: 0.5rem !important;
    }
    
    [data-testid="stSidebar"] [data-testid="stFileUploader"] section {
        background: transparent !important;
    }
    
    [data-testid="stSidebar"] [data-testid="stFileUploader"] button {
        background: #6366f1 !important;
        color: white !important;
        border-radius: 6px !important;
    }
    
    /* 侧边栏 number_input 样式 */
    [data-testid="stSidebar"] .stNumberInput > div > div > input {
        background: white !important;
        color: #1e293b !important;
        border: 1px solid #e2e8f0 !important;
        border-radius: 8px !important;
    }
    
    /* 侧边栏链接样式 */
    [data-testid="stSidebar"] a {
        color: #6366f1 !important;
    }
    
    /* 侧边栏分隔线 */
    [data-testid="stSidebar"] hr {
        border-color: #e2e8f0 !important;
    }
    </style>
    """, unsafe_allow_html=True)
    
    # 初始化导航状态
    if 'current_page' not in st.session_state:
        st.session_state.current_page = "response_list"
    
    # Logo
    st.markdown("""
    <div class="sidebar-logo">
        <div class="sidebar-logo-icon">📊</div>
        <span class="sidebar-logo-text">Survey Insights</span>
    </div>
    """, unsafe_allow_html=True)
    
    # Survey Insights 模块
    st.markdown('<div class="nav-section-title">ANALYSIS</div>', unsafe_allow_html=True)
    
    if st.button("📋 Response List", key="nav_response_list", use_container_width=True, 
                 type="primary" if st.session_state.current_page == "response_list" else "secondary"):
        st.session_state.current_page = "response_list"
        st.rerun()
    
    if st.button("🗂️ Segmentation Kanban", key="nav_kanban", use_container_width=True,
                 type="primary" if st.session_state.current_page == "kanban" else "secondary"):
        st.session_state.current_page = "kanban"
        st.rerun()
    
    if st.button("📈 Trends Dashboard", key="nav_trends", use_container_width=True,
                 type="primary" if st.session_state.current_page == "trends" else "secondary"):
        st.session_state.current_page = "trends"
        st.rerun()
    
    # Reporting 模块
    st.markdown('<div class="nav-section-title">REPORTING</div>', unsafe_allow_html=True)
    
    if st.button("📝 Summary Timeline", key="nav_timeline", use_container_width=True,
                 type="primary" if st.session_state.current_page == "timeline" else "secondary"):
        st.session_state.current_page = "timeline"
        st.rerun()
    
    if st.button("👤 Respondent Gallery", key="nav_gallery", use_container_width=True,
                 type="primary" if st.session_state.current_page == "gallery" else "secondary"):
        st.session_state.current_page = "gallery"
        st.rerun()
    
    # Deep Analysis 模块
    st.markdown('<div class="nav-section-title">DEEP ANALYSIS</div>', unsafe_allow_html=True)
    
    if st.button("📊 Distribution Analysis", key="nav_distribution", use_container_width=True,
                 type="primary" if st.session_state.current_page == "distribution" else "secondary"):
        st.session_state.current_page = "distribution"
        st.rerun()
    
    if st.button("🔀 Cross Analysis", key="nav_cross", use_container_width=True,
                 type="primary" if st.session_state.current_page == "cross_analysis" else "secondary"):
        st.session_state.current_page = "cross_analysis"
        st.rerun()
    
    if st.button("🤖 AI Summary", key="nav_ai_summary", use_container_width=True,
                 type="primary" if st.session_state.current_page == "ai_summary" else "secondary"):
        st.session_state.current_page = "ai_summary"
        st.rerun()
    
    if st.button("📤 Export Reports", key="nav_export", use_container_width=True,
                 type="primary" if st.session_state.current_page == "export" else "secondary"):
        st.session_state.current_page = "export"
        st.rerun()
    
    st.divider()
    
    # AI 配置
    st.markdown('<div class="sidebar-section-title">🤖 AI 配置</div>', unsafe_allow_html=True)
    
    # AI模型选择
    ai_provider = st.selectbox(
        "选择AI服务商",
        ["规则分析(无需API)", "OpenAI", "Claude (Anthropic)", "自定义API"],
        help="选择要使用的AI服务",
        key="sidebar_ai_provider"
    )
    
    # API配置
    if ai_provider != "规则分析(无需API)":
        api_key = st.text_input(
            "API Key",
            type="password",
            help="输入你的API密钥",
            key="api_key_input"
        )
        
        if ai_provider == "OpenAI":
            model = st.selectbox(
                "模型",
                ["gpt-4o", "gpt-4o-mini", "gpt-4-turbo", "gpt-3.5-turbo"],
                help="选择OpenAI模型",
                key="sidebar_openai_model"
            )
            api_base = "https://api.openai.com/v1"
        elif ai_provider == "Claude (Anthropic)":
            model = st.selectbox(
                "模型",
                ["claude-3-5-sonnet-20241022", "claude-3-opus-20240229", "claude-3-sonnet-20240229", "claude-3-haiku-20240307"],
                help="选择Claude模型",
                key="sidebar_claude_model"
            )
            api_base = "https://api.anthropic.com/v1"
        else:  # 自定义API
            api_base = st.text_input(
                "API Base URL",
                placeholder="https://api.example.com/v1",
                help="输入API基础URL",
                key="sidebar_custom_api_base"
            )
            model = st.text_input(
                "模型名称",
                placeholder="gpt-4",
                help="输入模型名称",
                key="sidebar_custom_model"
            )
        
        # 保存配置到session_state
        if api_key:
            st.session_state.ai_config = {
                'provider': ai_provider,
                'api_key': api_key,
                'model': model,
                'api_base': api_base
            }
            st.success("✅ API配置已保存")
        else:
            st.warning("⚠️ 请输入API Key")
    else:
        st.session_state.ai_config = {'provider': '规则分析(无需API)'}
        st.info("💡 使用内置规则分析,无需API")
    
    st.markdown("---")
    
    # 导出功能区域 - 简化版
    st.markdown("### 📤 导出报告")
    st.markdown('<p style="color: #71717a; font-size: 0.75rem; margin-top: -0.5rem; margin-bottom: 1rem;">一键导出当前分析结果</p>', unsafe_allow_html=True)
    
    st.markdown("---")
    st.markdown("""
    <div style="text-align: center; padding: 1.5rem 0;">
        <div style="color: #6c757d; font-size: 0.75rem; margin-bottom: 0.5rem;">
            Powered by
        </div>
        <div style="display: flex; justify-content: center; gap: 0.8rem; flex-wrap: wrap;">
            <span style="background: white; padding: 0.3rem 0.8rem; border-radius: 15px; font-size: 0.75rem; font-weight: 600; color: #667eea; border: 1px solid #e9ecef;">
                Streamlit
            </span>
            <span style="background: white; padding: 0.3rem 0.8rem; border-radius: 15px; font-size: 0.75rem; font-weight: 600; color: #667eea; border: 1px solid #e9ecef;">
                Plotly
            </span>
            <span style="background: white; padding: 0.3rem 0.8rem; border-radius: 15px; font-size: 0.75rem; font-weight: 600; color: #667eea; border: 1px solid #e9ecef;">
                Pandas
            </span>
        </div>
        <div style="color: #adb5bd; font-size: 0.7rem; margin-top: 0.8rem;">
            © 2025 DataInsight Pro | v2.0
        </div>
    </div>
    """, unsafe_allow_html=True)

# --- 辅助函数 ---
def read_pdf(file):
    pdf_reader = PyPDF2.PdfReader(file)
    text = ""
    for page in pdf_reader.pages:
        txt = page.extract_text()
        if txt: text += txt
    return text

def read_docx(file):
    doc = Document(file)
    return "\n".join([para.text for para in doc.paragraphs])

def generate_quick_summary(df):
    """生成数据快速总结"""
    summary = {}
    
    # 基础信息
    summary['rows'] = df.shape[0]
    summary['cols'] = df.shape[1]
    summary['numeric_cols'] = df.select_dtypes(include=['number']).shape[1]
    summary['cat_cols'] = df.select_dtypes(include=['object']).shape[1]
    
    # 数据质量
    missing_total = df.isnull().sum().sum()
    summary['missing_pct'] = (missing_total / (df.shape[0] * df.shape[1])) * 100
    summary['completeness'] = 100 - summary['missing_pct']
    
    # 质量评级
    if summary['completeness'] >= 95:
        summary['quality_level'] = '优秀'
        summary['quality_color'] = '#28a745'
    elif summary['completeness'] >= 80:
        summary['quality_level'] = '良好'
        summary['quality_color'] = '#ffc107'
    else:
        summary['quality_level'] = '需改进'
        summary['quality_color'] = '#dc3545'
    
    # 关键发现
    findings = []
    
    # 检测数值型变量
    numeric_cols = df.select_dtypes(include=['number']).columns
    if len(numeric_cols) > 0:
        for col in numeric_cols[:2]:
            mean_val = df[col].mean()
            std_val = df[col].std()
            cv = (std_val / mean_val * 100) if mean_val != 0 else 0
            if cv > 50:
                findings.append(f"📊 {col}数据波动较大(CV={cv:.1f}%)")
            else:
                findings.append(f"📊 {col}数据较稳定(均值={mean_val:.1f})")
    
    # 检测类别型变量
    cat_cols = df.select_dtypes(include=['object']).columns
    if len(cat_cols) > 0:
        for col in cat_cols[:2]:
            top_val = df[col].value_counts().index[0]
            top_pct = df[col].value_counts().values[0] / len(df) * 100
            findings.append(f"🏷️ {col}中'{top_val}'占比最高({top_pct:.1f}%)")
    
    summary['findings'] = findings[:4]
    
    return summary

def generate_chart_insight(chart_type, data_info, df):
    """生成图表解读"""
    insights = []
    
    if chart_type == 'distribution':
        col = data_info.get('column')
        if col and col in df.columns:
            if pd.api.types.is_numeric_dtype(df[col]):
                mean_val = df[col].mean()
                median_val = df[col].median()
                std_val = df[col].std()
                skew = df[col].skew()
                
                if abs(mean_val - median_val) < std_val * 0.1:
                    insights.append("数据分布较为对称")
                elif mean_val > median_val:
                    insights.append("数据呈右偏分布,存在较大值拉高均值")
                else:
                    insights.append("数据呈左偏分布,存在较小值拉低均值")
                
                if std_val / mean_val > 0.5 if mean_val != 0 else False:
                    insights.append("数据离散程度较高,波动较大")
                else:
                    insights.append("数据较为集中,波动较小")
            else:
                value_counts = df[col].value_counts()
                top_val = value_counts.index[0]
                top_pct = value_counts.values[0] / len(df) * 100
                
                if top_pct > 50:
                    insights.append(f"'{top_val}'占主导地位({top_pct:.1f}%)")
                elif len(value_counts) > 10:
                    insights.append(f"类别较多({len(value_counts)}种),分布较分散")
                else:
                    insights.append(f"共{len(value_counts)}种类别,分布相对均匀")
    
    elif chart_type == 'correlation':
        x_col = data_info.get('x')
        y_col = data_info.get('y')
        if x_col and y_col and x_col in df.columns and y_col in df.columns:
            if pd.api.types.is_numeric_dtype(df[x_col]) and pd.api.types.is_numeric_dtype(df[y_col]):
                corr = df[x_col].corr(df[y_col])
                if abs(corr) > 0.7:
                    direction = "正" if corr > 0 else "负"
                    insights.append(f"两变量呈强{direction}相关(r={corr:.2f})")
                elif abs(corr) > 0.4:
                    direction = "正" if corr > 0 else "负"
                    insights.append(f"两变量呈中等{direction}相关(r={corr:.2f})")
                else:
                    insights.append(f"两变量相关性较弱(r={corr:.2f})")
    
    return " | ".join(insights) if insights else "点击查看详细分析"

def call_openai_api(messages, api_key, model, api_base):
    """调用OpenAI兼容的API"""
    try:
        headers = {
            "Authorization": f"Bearer {api_key}",
            "Content-Type": "application/json"
        }
        
        data = {
            "model": model,
            "messages": messages,
            "temperature": 0.7,
            "max_tokens": 2000
        }
        
        response = requests.post(
            f"{api_base}/chat/completions",
            headers=headers,
            json=data,
            timeout=30
        )
        
        if response.status_code == 200:
            return response.json()['choices'][0]['message']['content']
        else:
            return f"❌ API调用失败: {response.status_code} - {response.text}"
    except Exception as e:
        return f"❌ API调用出错: {str(e)}"

def call_claude_api(messages, api_key, model, api_base):
    """调用Claude API"""
    try:
        headers = {
            "x-api-key": api_key,
            "anthropic-version": "2023-06-01",
            "Content-Type": "application/json"
        }
        
        # 转换消息格式
        claude_messages = []
        for msg in messages:
            if msg['role'] != 'system':
                claude_messages.append({
                    "role": msg['role'],
                    "content": msg['content']
                })
        
        data = {
            "model": model,
            "messages": claude_messages,
            "max_tokens": 2000,
            "temperature": 0.7
        }
        
        # 如果有system消息,添加到data中
        system_msg = next((msg['content'] for msg in messages if msg['role'] == 'system'), None)
        if system_msg:
            data['system'] = system_msg
        
        response = requests.post(
            f"{api_base}/messages",
            headers=headers,
            json=data,
            timeout=30
        )
        
        if response.status_code == 200:
            return response.json()['content'][0]['text']
        else:
            return f"❌ API调用失败: {response.status_code} - {response.text}"
    except Exception as e:
        return f"❌ API调用出错: {str(e)}"

def generate_ai_response(question, df):
    """生成AI响应(支持真实API或基于规则的分析)"""
    
    # 准备数据摘要信息
    data_summary = f"""
数据集信息:
- 总样本数: {df.shape[0]}
- 变量数量: {df.shape[1]}
- 数值型变量: {df.select_dtypes(include=['number']).shape[1]}
- 类别型变量: {df.select_dtypes(include=['object']).shape[1]}
- 变量列表: {', '.join(df.columns.tolist()[:10])}{'...' if len(df.columns) > 10 else ''}

数据统计:
"""
    
    # 添加数值型变量统计
    numeric_cols = df.select_dtypes(include=['number']).columns
    if len(numeric_cols) > 0:
        data_summary += "\n数值型变量统计:\n"
        for col in numeric_cols[:5]:
            data_summary += f"- {col}: 均值={df[col].mean():.2f}, 标准差={df[col].std():.2f}, 范围=[{df[col].min():.2f}, {df[col].max():.2f}]\n"
    
    # 添加类别型变量统计
    cat_cols = df.select_dtypes(include=['object']).columns
    if len(cat_cols) > 0:
        data_summary += "\n类别型变量统计:\n"
        for col in cat_cols[:5]:
            top_values = df[col].value_counts().head(3)
            data_summary += f"- {col}: 唯一值数={df[col].nunique()}, 最常见值={top_values.index[0]}({top_values.values[0]}次)\n"
    
    # 检查是否配置了API
    ai_config = st.session_state.get('ai_config', {'provider': '规则分析(无需API)'})
    
    if ai_config['provider'] != '规则分析(无需API)' and 'api_key' in ai_config:
        # 使用真实API
        system_prompt = f"""你是一个专业的数据分析助手。用户上传了一个数据集,你需要根据数据信息回答用户的问题。

{data_summary}

请用中文回答,提供专业、清晰、有洞察力的分析。使用markdown格式美化输出。"""
        
        messages = [
            {"role": "system", "content": system_prompt},
            {"role": "user", "content": question}
        ]
        
        if ai_config['provider'] == 'Claude (Anthropic)':
            return call_claude_api(messages, ai_config['api_key'], ai_config['model'], ai_config['api_base'])
        else:  # OpenAI或自定义API
            return call_openai_api(messages, ai_config['api_key'], ai_config['model'], ai_config['api_base'])
    
    # 使用规则分析(原有逻辑)
    question_lower = question.lower()
    
    # 异常值检测
    if "异常" in question_lower or "outlier" in question_lower:
        response = "**异常值分析:**\n\n"
        numeric_cols = df.select_dtypes(include=['number']).columns
        if len(numeric_cols) > 0:
            for col in numeric_cols[:3]:
                Q1 = df[col].quantile(0.25)
                Q3 = df[col].quantile(0.75)
                IQR = Q3 - Q1
                outliers = df[(df[col] < Q1 - 1.5*IQR) | (df[col] > Q3 + 1.5*IQR)]
                if len(outliers) > 0:
                    response += f"- **{col}**: 发现 {len(outliers)} 个异常值 ({len(outliers)/len(df)*100:.1f}%)\n"
                else:
                    response += f"- **{col}**: 未发现明显异常值\n"
        else:
            response += "数据集中没有数值型变量,无法进行异常值检测。"
        return response
    
    # 相关性分析
    elif "相关" in question_lower or "correlation" in question_lower:
        response = "**变量相关性分析:**\n\n"
        numeric_cols = df.select_dtypes(include=['number']).columns
        if len(numeric_cols) >= 2:
            corr_matrix = df[numeric_cols].corr()
            # 找出最强的相关性
            strong_corr = []
            for i in range(len(corr_matrix.columns)):
                for j in range(i+1, len(corr_matrix.columns)):
                    corr_value = corr_matrix.iloc[i, j]
                    if abs(corr_value) > 0.5:
                        strong_corr.append((corr_matrix.columns[i], corr_matrix.columns[j], corr_value))
            
            if strong_corr:
                response += "发现以下强相关关系:\n"
                for var1, var2, corr in sorted(strong_corr, key=lambda x: abs(x[2]), reverse=True)[:5]:
                    strength = "强正相关" if corr > 0.7 else "正相关" if corr > 0 else "强负相关" if corr < -0.7 else "负相关"
                    response += f"- **{var1}** 与 **{var2}**: {strength} (r={corr:.3f})\n"
            else:
                response += "未发现显著的强相关关系(|r| > 0.5)"
        else:
            response += "数值型变量不足2个,无法进行相关性分析。"
        return response
    
    # 数据质量
    elif "质量" in question_lower or "quality" in question_lower:
        response = "**数据质量评估:**\n\n"
        total_cells = df.shape[0] * df.shape[1]
        missing_cells = df.isnull().sum().sum()
        completeness = (1 - missing_cells / total_cells) * 100
        
        response += f"**完整性:** {completeness:.2f}%\n"
        if completeness >= 95:
            response += "✅ 数据完整性很好\n\n"
        elif completeness >= 80:
            response += "⚠️ 数据完整性一般,建议处理缺失值\n\n"
        else:
            response += "❌ 数据完整性较差,需要重点处理缺失值\n\n"
        
        # 缺失值分析
        missing_by_col = df.isnull().sum()
        if missing_by_col.sum() > 0:
            response += "**缺失值分布:**\n"
            for col in missing_by_col[missing_by_col > 0].index[:5]:
                missing_pct = missing_by_col[col] / len(df) * 100
                response += f"- {col}: {missing_by_col[col]} 个 ({missing_pct:.1f}%)\n"
        
        # 重复值检测
        duplicates = df.duplicated().sum()
        response += f"\n**重复记录:** {duplicates} 条"
        if duplicates > 0:
            response += f" ({duplicates/len(df)*100:.1f}%)"
        
        return response
    
    # 有趣的发现
    elif "发现" in question_lower or "insight" in question_lower or "有趣" in question_lower:
        response = "**数据洞察:**\n\n"
        
        # 样本量
        response += f"📊 **样本规模:** 共 {df.shape[0]:,} 条记录\n\n"
        
        # 类别型变量的分布
        cat_cols = df.select_dtypes(include=['object']).columns
        if len(cat_cols) > 0:
            response += "**类别分布特征:**\n"
            for col in cat_cols[:2]:
                value_counts = df[col].value_counts()
                top_value = value_counts.index[0]
                top_pct = value_counts.values[0] / len(df) * 100
                response += f"- {col}: '{top_value}' 占比最高 ({top_pct:.1f}%)\n"
            response += "\n"
        
        # 数值型变量的特征
        numeric_cols = df.select_dtypes(include=['number']).columns
        if len(numeric_cols) > 0:
            response += "**数值特征:**\n"
            for col in numeric_cols[:2]:
                mean_val = df[col].mean()
                std_val = df[col].std()
                cv = (std_val / mean_val * 100) if mean_val != 0 else 0
                if cv > 50:
                    response += f"- {col}: 变异系数 {cv:.1f}%,数据波动较大\n"
                else:
                    response += f"- {col}: 平均值 {mean_val:.2f},数据较为集中\n"
        
        return response
    
    # 默认响应
    else:
        response = f"**关于 '{question}' 的分析:**\n\n"
        response += f"数据集包含 {df.shape[0]:,} 条记录和 {df.shape[1]} 个变量。\n\n"
        response += "**建议的分析方向:**\n"
        response += "- 使用'单变量分布'查看各变量的分布情况\n"
        response += "- 使用'交叉分析'探索变量之间的关系\n"
        response += "- 点击左侧的快速问题获取更多洞察\n\n"
        response += "💡 你可以问我:\n"
        response += "- 数据中有哪些异常值?\n"
        response += "- 变量之间有什么相关性?\n"
        response += "- 如何提高数据质量?\n"
        return response

# --- 页面标题 ---
st.markdown("""
<div class="page-header">
    <div class="page-title">Survey Insights</div>
    <div class="page-subtitle">调研数据分析工具 · 支持 Excel/CSV 交叉分析与文本挖掘</div>
</div>
""", unsafe_allow_html=True)

# --- 主逻辑 ---
# 判断数据来源：上传文件优先，其次是已保存数据，最后是 URL 获取的数据
has_url_data = 'url_df' in st.session_state and st.session_state['url_df'] is not None
has_saved_data = 'saved_df' in st.session_state and st.session_state['saved_df'] is not None

if uploaded_file or has_saved_data or has_url_data:
    # 确定文件类型和数据来源
    if uploaded_file:
        file_type = uploaded_file.name.split('.')[-1].lower()
        data_source = 'file'
    elif has_saved_data:
        file_type = 'csv'  # 已保存数据视为 CSV 格式
        data_source = 'saved'
    else:
        file_type = 'csv'  # URL 数据默认为 CSV
        data_source = 'url'
    
    # ==========================================
    # 模块 A: 结构化数据分析 (Excel/CSV)
    # ==========================================
    if file_type in ['csv', 'xlsx']:
        try:
            # 根据数据来源获取 DataFrame
            if data_source == 'saved':
                # 使用已保存的数据
                df = st.session_state['saved_df']
            elif data_source == 'url':
                # 使用从 URL 获取的数据
                df = st.session_state['url_df']
            else:
                # 使用 session_state 缓存数据，彻底避免重复读取
                cache_key = f"df_cache_{uploaded_file.name}_{uploaded_file.size}"
                
                if cache_key not in st.session_state:
                    # 只在第一次上传时读取数据
                    if file_type == 'csv':
                        uploaded_file.seek(0)
                        file_content = uploaded_file.read().decode('utf-8')
                        st.session_state[cache_key] = load_csv_data(file_content, uploaded_file.name)
                    else:
                        uploaded_file.seek(0)
                        file_content = uploaded_file.read()
                        st.session_state[cache_key] = load_excel_data(file_content, uploaded_file.name)
                
                # 从 session_state 获取数据
                df = st.session_state[cache_key]
            
            # ==========================================
            # 数据来源提示
            # ==========================================
            if data_source == 'saved':
                saved_name = st.session_state.get('saved_df_name', '未命名')
                st.markdown(f"""
                <div style="background: linear-gradient(90deg, #667eea 0%, #764ba2 100%); padding: 0.75rem 1rem; border-radius: 8px; margin-bottom: 1rem; display: flex; align-items: center; gap: 0.75rem;">
                    <span style="font-size: 1.25rem;">💾</span>
                    <div>
                        <div style="font-size: 0.9rem; font-weight: 600; color: white;">已加载保存的数据: {saved_name}</div>
                        <div style="font-size: 0.75rem; color: rgba(255,255,255,0.8);">{len(df)} 条记录 · {len(df.columns)} 个字段</div>
                    </div>
                </div>
                """, unsafe_allow_html=True)
            elif data_source == 'url':
                st.markdown(f"""
                <div style="background: linear-gradient(90deg, #11998e 0%, #38ef7d 100%); padding: 0.75rem 1rem; border-radius: 8px; margin-bottom: 1rem; display: flex; align-items: center; gap: 0.75rem;">
                    <span style="font-size: 1.25rem;">🔗</span>
                    <div>
                        <div style="font-size: 0.9rem; font-weight: 600; color: white;">已从链接加载数据</div>
                        <div style="font-size: 0.75rem; color: rgba(255,255,255,0.8);">{len(df)} 条记录 · {len(df.columns)} 个字段</div>
                    </div>
                </div>
                """, unsafe_allow_html=True)
            
            # ==========================================
            # 顶部 KPI 仪表盘 - 核心指标一览
            # ==========================================
            st.markdown("---")
            
            # 计算核心指标
            total_responses = df.shape[0]
            total_questions = df.shape[1]
            completeness = (1 - df.isnull().sum().sum() / (df.shape[0] * df.shape[1])) * 100
            
            # 尝试检测 NPS 相关列
            nps_score = None
            nps_col = None
            for col in df.columns:
                col_lower = str(col).lower()
                if 'nps' in col_lower or 'recommend' in col_lower or '推荐' in col_lower:
                    if df[col].dtype in ['int64', 'float64']:
                        nps_col = col
                        # 计算 NPS
                        promoters = (df[col] >= 9).sum()
                        detractors = (df[col] <= 6).sum()
                        nps_score = int((promoters - detractors) / len(df[col].dropna()) * 100)
                        break
            
            # KPI 卡片行
            kpi_cols = st.columns(4)
            
            with kpi_cols[0]:
                st.markdown(f"""
                <div class="kpi-card">
                    <div class="kpi-label">总回收量</div>
                    <div class="kpi-value">{total_responses:,}</div>
                    <span class="kpi-delta positive">↑ 数据已加载</span>
                </div>
                """, unsafe_allow_html=True)
            
            with kpi_cols[1]:
                if nps_score is not None:
                    delta_class = "positive" if nps_score > 0 else ("negative" if nps_score < 0 else "neutral")
                    st.markdown(f"""
                    <div class="kpi-card">
                        <div class="kpi-label">NPS 净推荐值</div>
                        <div class="kpi-value">{nps_score}</div>
                        <span class="kpi-delta {delta_class}">{'↑' if nps_score > 0 else '↓' if nps_score < 0 else '→'} 已检测</span>
                    </div>
                    """, unsafe_allow_html=True)
                else:
                    st.markdown(f"""
                    <div class="kpi-card">
                        <div class="kpi-label">问题数量</div>
                        <div class="kpi-value">{total_questions}</div>
                        <span class="kpi-delta neutral">个字段</span>
                    </div>
                    """, unsafe_allow_html=True)
            
            with kpi_cols[2]:
                st.markdown(f"""
                <div class="kpi-card">
                    <div class="kpi-label">数据完整度</div>
                    <div class="kpi-value">{completeness:.1f}%</div>
                    <span class="kpi-delta {'positive' if completeness > 90 else 'warning' if completeness > 70 else 'negative'}">
                        {'优秀' if completeness > 90 else '良好' if completeness > 70 else '需关注'}
                    </span>
                </div>
                """, unsafe_allow_html=True)
            
            with kpi_cols[3]:
                numeric_cols = df.select_dtypes(include=['number']).shape[1]
                text_cols = df.select_dtypes(include=['object']).shape[1]
                st.markdown(f"""
                <div class="kpi-card">
                    <div class="kpi-label">字段类型</div>
                    <div class="kpi-value">{numeric_cols}/{text_cols}</div>
                    <span class="kpi-delta neutral">数值/文本</span>
                </div>
                """, unsafe_allow_html=True)
            
            # 快速总结面板 - 新样式
            quick_summary = generate_quick_summary(df)
            
            st.markdown(f"""
            <div class="chart-container" style="border-left: 4px solid {quick_summary['quality_color']};">
                <div class="chart-title">
                    <span>📋</span> 数据快速诊断
                    <div style="margin-left: auto; display: flex; gap: 0.75rem;">
                        <span class="badge" style="background: {quick_summary['quality_color']}20; color: {quick_summary['quality_color']};">
                            质量: {quick_summary['quality_level']}
                        </span>
                        <span class="badge badge-primary">
                            完整度: {quick_summary['completeness']:.1f}%
                        </span>
                    </div>
                </div>
                <div style="display: grid; grid-template-columns: repeat(auto-fit, minmax(220px, 1fr)); gap: 0.75rem; margin-top: 1rem;">
                    {''.join([f'<div style="background: var(--gray-50); padding: 0.75rem 1rem; border-radius: var(--radius); font-size: 0.85rem; color: var(--gray-800); border: 1px solid var(--gray-200);">{finding}</div>' for finding in quick_summary['findings']])}
                </div>
            </div>
            """, unsafe_allow_html=True)

            # ========== 基于导航状态渲染不同页面 ==========
            current_page = st.session_state.get('current_page', 'response_list')
            
            # 页面标题映射 - 匹配 Airtable 结构
            page_titles = {
                'response_list': ('📋 Response List', '查看所有调研回复数据'),
                'kanban': ('🗂️ Segmentation Kanban', '按状态分组查看调研回复'),
                'trends': ('📈 Trends Dashboard', '调研数据趋势与洞察'),
                'timeline': ('📝 Summary Timeline', '调研项目时间线概览'),
                'gallery': ('👤 Respondent Gallery', '受访者信息画廊'),
                'distribution': ('📊 Distribution Analysis', '单变量分布分析'),
                'cross_analysis': ('🔀 Cross Analysis', '多维交叉分析'),
                'ai_summary': ('🤖 AI Summary', 'AI 智能分析与洞察'),
                'export': ('📤 Export Reports', '导出分析报告')
            }
            
            # 显示当前页面标题
            page_title, page_desc = page_titles.get(current_page, ('📋 Response List', ''))
            st.markdown(f"""
            <div style="margin-bottom: 1.5rem;">
                <h2 style="margin: 0; font-size: 1.5rem; font-weight: 600; color: #18181b;">{page_title}</h2>
                <p style="margin: 0.25rem 0 0; font-size: 0.875rem; color: #71717a;">{page_desc}</p>
            </div>
            """, unsafe_allow_html=True)
            
            # ==================== Response List 页面 ====================
            if current_page == 'response_list':
                st.markdown("#### 📄 原始数据预览")
                
                # 获取问题映射
                question_map = st.session_state.get('question_map', {})
                
                # 工具栏：切换显示模式 + 下载按钮
                toolbar_col1, toolbar_col2, toolbar_col3 = st.columns([2, 1, 1])
                
                with toolbar_col1:
                    if question_map:
                        display_mode = st.radio(
                            "列名显示",
                            ["简写字段名", "完整问题"],
                            horizontal=True,
                            key="display_mode_radio"
                        )
                    else:
                        display_mode = "简写字段名"
                
                with toolbar_col3:
                    # 下载按钮
                    csv_data = df.to_csv(index=False).encode('utf-8-sig')
                    file_prefix = uploaded_file.name.split('.')[0] if uploaded_file else "url_data"
                    st.download_button(
                        label="📥 下载 CSV",
                        data=csv_data,
                        file_name=f"survey_data_{file_prefix}.csv",
                        mime="text/csv",
                        use_container_width=True
                    )
                
                # 根据显示模式处理数据
                display_df = df.copy()
                if question_map and display_mode == "完整问题":
                    # 将列名替换为完整问题
                    new_columns = []
                    for col in df.columns:
                        if col in question_map:
                            new_columns.append(question_map[col][:50] + "..." if len(question_map[col]) > 50 else question_map[col])
                        else:
                            new_columns.append(col)
                    display_df.columns = new_columns
                
                # 显示字段映射（折叠）
                if question_map:
                    mapped_cols = [col for col in df.columns if col in question_map]
                    if mapped_cols:
                        with st.expander(f"📋 字段映射 ({len(mapped_cols)} 个)", expanded=False):
                            for col in mapped_cols:
                                st.caption(f"**{col}** → {question_map[col]}")
                
                st.dataframe(display_df, use_container_width=True, height=400)
                
                st.markdown("#### 📈 数据统计摘要")
                col_a, col_b = st.columns(2)
                
                with col_a:
                    st.markdown("**数值型变量统计**")
                    numeric_df = df.describe()
                    if not numeric_df.empty:
                        st.dataframe(numeric_df, use_container_width=True)
                    else:
                        st.info("暂无数值型变量")
                
                with col_b:
                    st.markdown("**数据质量检查**")
                    quality_df = pd.DataFrame({
                        '字段名': df.columns,
                        '缺失值': df.isnull().sum().values,
                        '缺失率': (df.isnull().sum() / len(df) * 100).round(2).astype(str) + '%',
                        '唯一值数': [df[col].nunique() for col in df.columns]
                    })
                    st.dataframe(quality_df, use_container_width=True)
                
                # 智能分析建议
                st.markdown("#### 💡 智能分析建议")
                
                suggestions = []
                
                # 基于数据特征生成建议
                numeric_cols = df.select_dtypes(include=['number']).columns.tolist()
                cat_cols = df.select_dtypes(include=['object']).columns.tolist()
                
                # 检查缺失值
                high_missing = df.columns[df.isnull().sum() / len(df) > 0.1].tolist()
                if high_missing:
                    suggestions.append({
                        'icon': '⚠️',
                        'type': '数据质量',
                        'title': f'{len(high_missing)}个字段缺失率>10%',
                        'desc': f'建议检查: {", ".join(high_missing[:3])}...' if len(high_missing) > 3 else f'建议检查: {", ".join(high_missing)}',
                        'color': '#ffc107'
                    })
                
                # 推荐分析
                if len(numeric_cols) >= 2:
                    suggestions.append({
                        'icon': '📊',
                        'type': '推荐分析',
                        'title': '可进行相关性分析',
                        'desc': f'发现{len(numeric_cols)}个数值变量,可使用散点图探索变量关系',
                        'color': '#667eea'
                    })
                
                if len(cat_cols) >= 1 and len(numeric_cols) >= 1:
                    suggestions.append({
                        'icon': '🔀',
                        'type': '推荐分析',
                        'title': '可进行分组对比',
                        'desc': f'可按{cat_cols[0]}分组,对比{numeric_cols[0]}的分布差异',
                        'color': '#764ba2'
                    })
                
                if len(cat_cols) >= 2:
                    suggestions.append({
                        'icon': '🔥',
                        'type': '推荐分析',
                        'title': '可进行交叉分析',
                        'desc': f'可分析{cat_cols[0]}与{cat_cols[1]}的关联关系',
                        'color': '#28a745'
                    })
                
                # 检查异常值
                for col in numeric_cols[:2]:
                    q1 = df[col].quantile(0.25)
                    q3 = df[col].quantile(0.75)
                    iqr = q3 - q1
                    outliers = ((df[col] < q1 - 1.5*iqr) | (df[col] > q3 + 1.5*iqr)).sum()
                    if outliers > 0:
                        suggestions.append({
                            'icon': '🔍',
                            'type': '异常检测',
                            'title': f'{col}存在{outliers}个异常值',
                            'desc': '建议在箱线图中查看具体分布',
                            'color': '#dc3545'
                        })
                        break
                
                if suggestions:
                    # 使用Streamlit原生组件展示建议
                    cols = st.columns(min(len(suggestions), 2))
                    for i, sug in enumerate(suggestions[:4]):
                        with cols[i % 2]:
                            st.markdown(f"""
<div style="background: white; padding: 1rem; border-radius: 10px; border-left: 4px solid {sug['color']}; box-shadow: 0 2px 8px rgba(0,0,0,0.05); margin-bottom: 0.5rem;">
<div style="display: flex; align-items: center; gap: 0.5rem; margin-bottom: 0.3rem;">
<span style="font-size: 1.2rem;">{sug['icon']}</span>
<span style="background: {sug['color']}20; color: {sug['color']}; padding: 0.2rem 0.6rem; border-radius: 10px; font-size: 0.75rem; font-weight: 600;">{sug['type']}</span>
<span style="font-weight: 600; color: #333;">{sug['title']}</span>
</div>
<div style="color: #6c757d; font-size: 0.85rem; padding-left: 1.8rem;">{sug['desc']}</div>
</div>
""", unsafe_allow_html=True)
                else:
                    st.success("✅ 数据质量良好,可以开始分析!")

            # ==================== Distribution Analysis 页面 ====================
            elif current_page == 'distribution':
                st.markdown("#### 📊 单变量分布分析")
                st.markdown("选择一个或多个变量,系统将自动生成适合的可视化图表")
                
                # 初始化默认选择（只在第一次时设置）
                data_key = uploaded_file.name if uploaded_file else "url_data"
                default_key = f"default_cols_{data_key}"
                if default_key not in st.session_state:
                    st.session_state[default_key] = [df.columns[0]] if len(df.columns) > 0 else []
                
                # 支持多选问题 - 使用固定 key 避免重建
                col_select_list = st.multiselect(
                    "🔍 选择分析变量:", 
                    df.columns.tolist(),
                    default=st.session_state[default_key],
                    help="可以选择多个变量进行批量分析",
                    key="var_multiselect"
                )
                
                if not col_select_list:
                    st.warning("⚠️ 请至少选择一个变量进行分析")
                else:
                    # 获取问题映射
                    question_map = st.session_state.get('question_map', {})
                    
                    # 为每个选中的列生成图表
                    for idx, col_select in enumerate(col_select_list):
                        # 使用容器创建更好的视觉分隔
                        with st.container():
                            # 显示字段名和完整问题（如果有映射）
                            full_question = question_map.get(col_select, "")
                            if full_question:
                                st.markdown(f"### 📌 {col_select}")
                                st.caption(f"📝 {full_question}")
                            else:
                                st.markdown(f"### 📌 {col_select}")
                
                            # 智能判断图表类型
                            is_numeric = pd.api.types.is_numeric_dtype(df[col_select])
                            unique_values = df[col_select].nunique()
                            
                            # 判断是否为评分型数据（数值型、唯一值<=10、且为整数）
                            is_rating_data = False
                            if is_numeric and unique_values <= 10:
                                # 检查是否都是整数
                                non_null_values = df[col_select].dropna()
                                if len(non_null_values) > 0:
                                    is_rating_data = (non_null_values == non_null_values.astype(int)).all()
                            
                            if is_numeric and not is_rating_data:
                                # 纯数值型（连续变量）-> 直方图和箱线图
                                chart_col1, chart_col2 = st.columns(2)
                                
                                with chart_col1:
                                    fig = px.histogram(
                                        df, x=col_select, 
                                        title=f"📊 {col_select} 分布图", 
                                        nbins=20,
                                        color_discrete_sequence=['#667eea']
                                    )
                                    fig.update_layout(
                                        plot_bgcolor='rgba(0,0,0,0)',
                                        paper_bgcolor='rgba(0,0,0,0)',
                                        font=dict(size=12)
                                    )
                                    st.plotly_chart(fig, use_container_width=True)
                                
                                with chart_col2:
                                    fig_box = px.box(
                                        df, y=col_select,
                                        title=f"📦 {col_select} 箱线图",
                                        color_discrete_sequence=['#764ba2']
                                    )
                                    fig_box.update_layout(
                                        plot_bgcolor='rgba(0,0,0,0)',
                                        paper_bgcolor='rgba(0,0,0,0)',
                                        font=dict(size=12)
                                    )
                                    st.plotly_chart(fig_box, use_container_width=True)
                                
                                # 自动生成数值型变量解读
                                mean_val = df[col_select].mean()
                                median_val = df[col_select].median()
                                std_val = df[col_select].std()
                                min_val = df[col_select].min()
                                max_val = df[col_select].max()
                                skew = df[col_select].skew()
                                
                                # 判断分布特征
                                if abs(skew) < 0.5:
                                    dist_desc = "近似正态分布"
                                    dist_icon = "✅"
                                elif skew > 0:
                                    dist_desc = "右偏分布(存在较大值)"
                                    dist_icon = "📈"
                                else:
                                    dist_desc = "左偏分布(存在较小值)"
                                    dist_icon = "📉"
                                
                                cv = (std_val / mean_val * 100) if mean_val != 0 else 0
                                if cv > 50:
                                    var_desc = "波动较大"
                                    var_color = "#ffc107"
                                else:
                                    var_desc = "较为稳定"
                                    var_color = "#28a745"
                                
                                st.markdown(f"""
                                <div style="background: linear-gradient(135deg, #f8f9fa 0%, #e9ecef 100%); padding: 1rem; border-radius: 10px; margin-top: 0.5rem; border-left: 4px solid #667eea;">
                                    <div style="font-weight: 600; color: #333; margin-bottom: 0.5rem;">💡 自动解读</div>
                                    <div style="display: grid; grid-template-columns: repeat(auto-fit, minmax(150px, 1fr)); gap: 0.5rem; font-size: 0.9rem;">
                                        <div>📊 均值: <strong>{mean_val:.2f}</strong></div>
                                        <div>📍 中位数: <strong>{median_val:.2f}</strong></div>
                                        <div>📏 范围: <strong>{min_val:.2f} ~ {max_val:.2f}</strong></div>
                                        <div>{dist_icon} {dist_desc}</div>
                                        <div style="color: {var_color};">📈 {var_desc} (CV={cv:.1f}%)</div>
                                    </div>
                                </div>
                                """, unsafe_allow_html=True)
                            
                            elif is_rating_data:
                                # 评分型数据 -> 专用评分分布图
                                mean_val = df[col_select].mean()
                                median_val = df[col_select].median()
                                std_val = df[col_select].std()
                                min_val = int(df[col_select].min())
                                max_val = int(df[col_select].max())
                                total_responses = df[col_select].count()
                                
                                # 统计各评分的频次
                                rating_counts = df[col_select].value_counts().sort_index()
                                all_ratings = list(range(min_val, max_val + 1))
                                rating_df = pd.DataFrame({
                                    '评分': all_ratings,
                                    '人数': [rating_counts.get(r, 0) for r in all_ratings]
                                })
                                rating_df['占比'] = (rating_df['人数'] / total_responses * 100).round(1)
                                
                                # 评分分布柱状图
                                colors = ['#ef4444', '#f97316', '#eab308', '#84cc16', '#22c55e']  # 红到绿渐变
                                if max_val - min_val + 1 <= len(colors):
                                    bar_colors = colors[-(max_val - min_val + 1):]
                                else:
                                    bar_colors = px.colors.qualitative.Set2[:len(all_ratings)]
                                
                                fig = go.Figure()
                                
                                for i, row in rating_df.iterrows():
                                    fig.add_trace(go.Bar(
                                        x=[str(int(row['评分']))],
                                        y=[row['人数']],
                                        name=f"{int(row['评分'])}分",
                                        marker_color=bar_colors[i % len(bar_colors)],
                                        text=f"{row['人数']}<br>({row['占比']}%)",
                                        textposition='outside',
                                        hovertemplate=f"评分: {int(row['评分'])}分<br>人数: {row['人数']}<br>占比: {row['占比']}%<extra></extra>"
                                    ))
                                
                                fig.update_layout(
                                    title=dict(
                                        text=f"⭐ 评分分布 (1-{max_val}分)",
                                        font=dict(size=16)
                                    ),
                                    xaxis_title="评分",
                                    yaxis_title="人数",
                                    showlegend=False,
                                    plot_bgcolor='rgba(0,0,0,0)',
                                    paper_bgcolor='rgba(0,0,0,0)',
                                    font=dict(size=12),
                                    height=400,
                                    bargap=0.3
                                )
                                
                                st.plotly_chart(fig, use_container_width=True)
                                
                                # 评分解读卡片
                                # 计算满意度指标
                                high_score_count = df[col_select][df[col_select] >= (max_val - 1)].count()  # 高分（最高两档）
                                low_score_count = df[col_select][df[col_select] <= (min_val + 1)].count()   # 低分（最低两档）
                                high_score_pct = (high_score_count / total_responses * 100) if total_responses > 0 else 0
                                low_score_pct = (low_score_count / total_responses * 100) if total_responses > 0 else 0
                                
                                # NPS 类似评分（假设满分为5时）
                                if max_val == 5:
                                    promoters = df[col_select][df[col_select] >= 4].count()
                                    detractors = df[col_select][df[col_select] <= 2].count()
                                    nps_score = ((promoters - detractors) / total_responses * 100) if total_responses > 0 else 0
                                    nps_display = f"<div>📊 NPS得分: <strong style='color: {'#22c55e' if nps_score > 0 else '#ef4444'};'>{nps_score:.0f}</strong></div>"
                                else:
                                    nps_display = ""
                                
                                # 评分等级判断
                                score_ratio = (mean_val - min_val) / (max_val - min_val) if max_val > min_val else 0
                                if score_ratio >= 0.8:
                                    rating_level = "优秀"
                                    rating_color = "#22c55e"
                                    rating_icon = "🌟"
                                elif score_ratio >= 0.6:
                                    rating_level = "良好"
                                    rating_color = "#84cc16"
                                    rating_icon = "👍"
                                elif score_ratio >= 0.4:
                                    rating_level = "一般"
                                    rating_color = "#eab308"
                                    rating_icon = "😐"
                                else:
                                    rating_level = "待改进"
                                    rating_color = "#ef4444"
                                    rating_icon = "⚠️"
                                
                                st.markdown(f"""
                                <div style="background: linear-gradient(135deg, #fefce8 0%, #fef9c3 100%); padding: 1.25rem; border-radius: 12px; margin-top: 0.5rem; border-left: 4px solid #eab308;">
                                    <div style="font-weight: 600; color: #854d0e; margin-bottom: 0.75rem; font-size: 1rem;">
                                        {rating_icon} 评分分析 · <span style="color: {rating_color};">{rating_level}</span>
                                    </div>
                                    <div style="display: grid; grid-template-columns: repeat(auto-fit, minmax(140px, 1fr)); gap: 0.75rem; font-size: 0.875rem; color: #713f12;">
                                        <div>⭐ 平均分: <strong>{mean_val:.2f}</strong> / {max_val}</div>
                                        <div>📍 中位数: <strong>{median_val:.1f}</strong></div>
                                        <div>👍 高分率: <strong style="color: #22c55e;">{high_score_pct:.1f}%</strong></div>
                                        <div>👎 低分率: <strong style="color: #ef4444;">{low_score_pct:.1f}%</strong></div>
                                        <div>📊 样本数: <strong>{total_responses}</strong></div>
                                        {nps_display}
                                    </div>
                                </div>
                                """, unsafe_allow_html=True)
                                
                            else:
                                # 类别型 -> 提供多种图表选择
                                chart_col1, chart_col2 = st.columns([1, 4])
                                with chart_col1:
                                    st.markdown("**图表类型**")
                                    chart_option = st.radio(
                                        "选择展示方式", 
                                        ["📊 柱状图", "🥧 饼状图", "📈 条形图"],
                                        key=f"chart_type_{col_select}",
                                        label_visibility="collapsed"
                                    )
                                    
                                    st.markdown("**配色方案**")
                                    color_scheme = st.selectbox(
                                        "选择颜色",
                                        ["紫色 Purples", "蓝色 Blues", "绿色 Greens", "橙色 Oranges", 
                                         "红色 Reds", "粉色 Pinkyl", "青色 Teal", "彩虹 Rainbow"],
                                        key=f"color_scheme_{col_select}",
                                        label_visibility="collapsed"
                                    )
                                
                                with chart_col2:
                                    value_counts_df = df[col_select].value_counts().reset_index()
                                    value_counts_df.columns = [col_select, 'count']
                                    
                                    # 颜色方案映射
                                    color_map = {
                                        "紫色 Purples": ("Purples", px.colors.sequential.Purples_r),
                                        "蓝色 Blues": ("Blues", px.colors.sequential.Blues_r),
                                        "绿色 Greens": ("Greens", px.colors.sequential.Greens_r),
                                        "橙色 Oranges": ("Oranges", px.colors.sequential.Oranges_r),
                                        "红色 Reds": ("Reds", px.colors.sequential.Reds_r),
                                        "粉色 Pinkyl": ("Pinkyl", px.colors.sequential.Pinkyl_r),
                                        "青色 Teal": ("Teal", px.colors.sequential.Teal_r),
                                        "彩虹 Rainbow": ("Rainbow", px.colors.qualitative.Vivid)
                                    }
                                    
                                    color_scale, color_discrete = color_map.get(color_scheme, ("Purples", px.colors.sequential.Purples_r))
                                    
                                    if "柱状图" in chart_option:
                                        # 图表设置面板
                                        with st.expander("⚙️ 图表显示设置", expanded=False):
                                            bar_col1, bar_col2, bar_col3 = st.columns(3)
                                            with bar_col1:
                                                bar_height = st.slider("图表高度", 300, 800, 450, 50, key=f"bar_height_{col_select}")
                                                show_values = st.checkbox("显示数值", value=True, key=f"bar_values_{col_select}")
                                            with bar_col2:
                                                bar_top_n = st.slider("显示前N项", 3, 30, min(15, len(value_counts_df)), 1, key=f"bar_topn_{col_select}")
                                                sort_order = st.selectbox("排序方式", ["按数量降序", "按数量升序", "按名称"], index=0, key=f"bar_sort_{col_select}")
                                            with bar_col3:
                                                x_angle = st.slider("X轴标签角度", -90, 0, -45, 15, key=f"bar_angle_{col_select}")
                                                max_label_len = st.slider("标签最大长度", 10, 50, 20, 5, key=f"bar_label_{col_select}")
                                        
                                        # 处理数据
                                        bar_data = value_counts_df.copy()
                                        if sort_order == "按数量降序":
                                            bar_data = bar_data.sort_values('count', ascending=False)
                                        elif sort_order == "按数量升序":
                                            bar_data = bar_data.sort_values('count', ascending=True)
                                        else:
                                            bar_data = bar_data.sort_values(col_select)
                                        
                                        bar_data = bar_data.head(bar_top_n).reset_index(drop=True)
                                        
                                        # 计算百分比
                                        total_count = value_counts_df['count'].sum()
                                        bar_data['percentage'] = (bar_data['count'] / total_count * 100).round(1)
                                        
                                        # 截断过长的标签
                                        bar_data['display_label'] = bar_data[col_select].apply(
                                            lambda x: str(x)[:max_label_len] + "..." if len(str(x)) > max_label_len else str(x)
                                        )
                                        
                                        # 准备数据明细
                                        display_df = value_counts_df.copy()
                                        display_total = display_df['count'].sum()
                                        display_df['percentage'] = (display_df['count'] / display_total * 100).round(1)
                                        display_df = display_df.sort_values('count', ascending=False).reset_index(drop=True)
                                        
                                        # 显示柱状图
                                        fig = px.bar(
                                            bar_data, 
                                            x='display_label', y='count', 
                                            labels={'display_label': col_select, 'count': '数量'},
                                            title=f"📊 {col_select} 频次统计",
                                            color='count',
                                            color_continuous_scale=color_scale,
                                            text='count' if show_values else None,
                                            custom_data=[col_select]
                                        )
                                        
                                        fig.update_traces(
                                            textposition='outside' if show_values else 'none',
                                            textfont_size=10,
                                            marker_line_color='white',
                                            marker_line_width=1,
                                            hovertemplate='<b>%{customdata[0]}</b><br>数量: %{y}<extra></extra>'
                                        )
                                        fig.update_layout(
                                            plot_bgcolor='rgba(0,0,0,0)',
                                            paper_bgcolor='rgba(0,0,0,0)',
                                            font=dict(size=11),
                                            showlegend=False,
                                            height=bar_height,
                                            xaxis=dict(tickangle=x_angle, title_font_size=11, tickfont_size=9),
                                            yaxis=dict(title_font_size=11, gridcolor='rgba(128,128,128,0.1)'),
                                            margin=dict(t=50, b=120, l=50, r=30)
                                        )
                                        st.plotly_chart(fig, use_container_width=True)
                                        
                                        # 数据明细放在图表下方，使用表格形式
                                        with st.expander("📋 查看数据明细", expanded=False):
                                            detail_df = display_df[[col_select, 'count', 'percentage']].copy()
                                            detail_df.columns = ['选项', '数量', '占比(%)']
                                            st.dataframe(detail_df, use_container_width=True, hide_index=True)
                                        
                                    elif "饼状图" in chart_option:
                                        # 图表设置面板
                                        with st.expander("⚙️ 图表显示设置", expanded=False):
                                            setting_col1, setting_col2, setting_col3 = st.columns(3)
                                            with setting_col1:
                                                chart_height = st.slider("图表高度", 300, 800, 450, 50, key=f"pie_height_{col_select}")
                                                show_legend = st.checkbox("显示图例", value=False, key=f"pie_legend_{col_select}")
                                            with setting_col2:
                                                text_display = st.selectbox(
                                                    "图内显示",
                                                    ["仅百分比", "仅数值", "百分比+数值", "不显示"],
                                                    index=0,
                                                    key=f"pie_text_{col_select}"
                                                )
                                                hole_size = st.slider("环形大小", 0.0, 0.7, 0.4, 0.1, key=f"pie_hole_{col_select}")
                                            with setting_col3:
                                                top_n = st.slider("显示前N项(其余合并为'其他')", 3, 20, min(10, len(value_counts_df)), 1, key=f"pie_topn_{col_select}")
                                                legend_pos = st.selectbox(
                                                    "图例位置",
                                                    ["底部横向", "右侧竖向", "左侧竖向"],
                                                    index=0,
                                                    key=f"pie_legend_pos_{col_select}"
                                                ) if show_legend else "底部横向"
                                        
                                        # 处理数据：合并小类别
                                        pie_data = value_counts_df.copy()
                                        pie_data = pie_data.sort_values('count', ascending=False).reset_index(drop=True)
                                        
                                        if len(pie_data) > top_n:
                                            top_data = pie_data.head(top_n).copy()
                                            other_count = pie_data.iloc[top_n:]['count'].sum()
                                            other_row = pd.DataFrame({col_select: ['其他'], 'count': [other_count]})
                                            pie_data = pd.concat([top_data, other_row], ignore_index=True)
                                        
                                        # 计算真实百分比（基于原始总数）
                                        original_total = value_counts_df['count'].sum()
                                        pie_data['percentage'] = (pie_data['count'] / original_total * 100).round(1)
                                        
                                        # 创建自定义文本标签
                                        if text_display == "仅百分比":
                                            pie_data['text_label'] = pie_data['percentage'].apply(lambda x: f"{x}%")
                                        elif text_display == "仅数值":
                                            pie_data['text_label'] = pie_data['count'].astype(str)
                                        elif text_display == "百分比+数值":
                                            pie_data['text_label'] = pie_data.apply(lambda r: f"{r['percentage']}%<br>({r['count']})", axis=1)
                                        else:
                                            pie_data['text_label'] = ""
                                        
                                        # 准备数据明细
                                        display_df = value_counts_df.copy()
                                        display_total = display_df['count'].sum()
                                        display_df['percentage'] = (display_df['count'] / display_total * 100).round(1)
                                        display_df = display_df.sort_values('count', ascending=False).reset_index(drop=True)
                                        
                                        # 只显示饼图，数据明细放在下方表格
                                        import plotly.graph_objects as go
                                        
                                        fig = go.Figure(data=[go.Pie(
                                            labels=pie_data[col_select],
                                            values=pie_data['count'],
                                            hole=hole_size,
                                            text=pie_data['text_label'],
                                            textposition='inside' if text_display != "不显示" else 'none',
                                            textinfo='text' if text_display != "不显示" else 'none',
                                            textfont=dict(size=11, color='white'),
                                            hovertemplate='<b>%{label}</b><br>数量: %{value}<br>占比: %{percent}<extra></extra>',
                                            marker=dict(
                                                colors=color_discrete[:len(pie_data)],
                                                line=dict(color='white', width=2)
                                            )
                                        )])
                                        
                                        fig.update_layout(
                                            title=f"🥧 {col_select} 占比分布",
                                            showlegend=show_legend,
                                            legend=dict(orientation="h", yanchor="top", y=-0.1, xanchor="center", x=0.5, font=dict(size=9)) if show_legend else dict(font=dict(size=9)),
                                            plot_bgcolor='rgba(0,0,0,0)',
                                            paper_bgcolor='rgba(0,0,0,0)',
                                            font=dict(size=12),
                                            height=chart_height,
                                            margin=dict(t=50, b=30 if not show_legend else 80, l=20, r=20)
                                        )
                                        
                                        st.plotly_chart(fig, use_container_width=True)
                                        
                                        # 数据明细放在图表下方，使用表格形式
                                        with st.expander("📋 查看数据明细", expanded=False):
                                            detail_df = display_df[[col_select, 'count', 'percentage']].copy()
                                            detail_df.columns = ['选项', '数量', '占比(%)']
                                            st.dataframe(detail_df, use_container_width=True, hide_index=True)
                                        
                                    else:  # 条形图(横向)
                                        # 图表设置面板
                                        with st.expander("⚙️ 图表显示设置", expanded=False):
                                            hbar_col1, hbar_col2, hbar_col3 = st.columns(3)
                                            with hbar_col1:
                                                hbar_height = st.slider("图表高度", 300, 1000, 500, 50, key=f"hbar_height_{col_select}")
                                                show_hbar_values = st.checkbox("显示数值", value=True, key=f"hbar_values_{col_select}")
                                            with hbar_col2:
                                                hbar_top_n = st.slider("显示前N项", 3, 30, min(15, len(value_counts_df)), 1, key=f"hbar_topn_{col_select}")
                                                hbar_sort = st.selectbox("排序方式", ["按数量降序", "按数量升序", "按名称"], index=0, key=f"hbar_sort_{col_select}")
                                            with hbar_col3:
                                                hbar_max_label = st.slider("标签最大长度", 15, 80, 40, 5, key=f"hbar_label_{col_select}")
                                                bar_thickness = st.slider("条形粗细", 10, 50, 25, 5, key=f"hbar_thick_{col_select}")
                                        
                                        # 处理数据
                                        hbar_data = value_counts_df.copy()
                                        if hbar_sort == "按数量降序":
                                            hbar_data = hbar_data.sort_values('count', ascending=True)  # 横向图需要反转
                                        elif hbar_sort == "按数量升序":
                                            hbar_data = hbar_data.sort_values('count', ascending=False)
                                        else:
                                            hbar_data = hbar_data.sort_values(col_select, ascending=False)
                                        
                                        hbar_data = hbar_data.head(hbar_top_n).reset_index(drop=True)
                                        
                                        # 截断过长的标签
                                        hbar_data['display_label'] = hbar_data[col_select].apply(
                                            lambda x: str(x)[:hbar_max_label] + "..." if len(str(x)) > hbar_max_label else str(x)
                                        )
                                        
                                        # 动态计算高度
                                        auto_height = max(hbar_height, len(hbar_data) * bar_thickness + 100)
                                        
                                        # 准备数据明细
                                        display_df = value_counts_df.copy()
                                        display_total = display_df['count'].sum()
                                        display_df['percentage'] = (display_df['count'] / display_total * 100).round(1)
                                        display_df = display_df.sort_values('count', ascending=False).reset_index(drop=True)
                                        
                                        # 显示条形图
                                        fig = px.bar(
                                            hbar_data, 
                                            x='count', y='display_label', 
                                            orientation='h',
                                            labels={'display_label': col_select, 'count': '数量'},
                                            title=f"📈 {col_select} 频次统计",
                                            color='count',
                                            color_continuous_scale=color_scale,
                                            text='count' if show_hbar_values else None,
                                            custom_data=[col_select]
                                        )
                                        
                                        fig.update_traces(
                                            textposition='outside' if show_hbar_values else 'none',
                                            textfont_size=10,
                                            marker_line_color='white',
                                            marker_line_width=1,
                                            hovertemplate='<b>%{customdata[0]}</b><br>数量: %{x}<extra></extra>'
                                        )
                                        
                                        fig.update_layout(
                                            plot_bgcolor='rgba(0,0,0,0)',
                                            paper_bgcolor='rgba(0,0,0,0)',
                                            font=dict(size=11),
                                            showlegend=False,
                                            height=auto_height,
                                            bargap=0.3,
                                            xaxis=dict(title_font_size=11, gridcolor='rgba(128,128,128,0.1)'),
                                            yaxis=dict(title_font_size=11, tickfont_size=9, automargin=True),
                                            margin=dict(t=50, b=40, l=10, r=50)
                                        )
                                        st.plotly_chart(fig, use_container_width=True)
                                        
                                        # 数据明细放在图表下方，使用表格形式
                                        with st.expander("📋 查看数据明细", expanded=False):
                                            detail_df = display_df[[col_select, 'count', 'percentage']].copy()
                                            detail_df.columns = ['选项', '数量', '占比(%)']
                                            st.dataframe(detail_df, use_container_width=True, hide_index=True)
                                
                                # 类别型变量自动解读
                                value_counts = df[col_select].value_counts()
                                total = len(df)
                                unique_count = len(value_counts)
                                top_val = value_counts.index[0]
                                top_count = value_counts.values[0]
                                top_pct = top_count / total * 100
                                
                                # 判断分布特征
                                if top_pct > 50:
                                    dist_desc = f"'{top_val}'占主导地位"
                                    dist_icon = "🎯"
                                elif unique_count <= 5:
                                    dist_desc = "类别较少,分布清晰"
                                    dist_icon = "✅"
                                elif unique_count > 20:
                                    dist_desc = "类别较多,建议合并分析"
                                    dist_icon = "⚠️"
                                else:
                                    dist_desc = "分布相对均匀"
                                    dist_icon = "📊"
                                
                                st.markdown(f"""
                                <div style="background: linear-gradient(135deg, #f8f9fa 0%, #e9ecef 100%); padding: 1rem; border-radius: 10px; margin-top: 0.5rem; border-left: 4px solid #764ba2;">
                                    <div style="font-weight: 600; color: #333; margin-bottom: 0.5rem;">💡 自动解读</div>
                                    <div style="display: grid; grid-template-columns: repeat(auto-fit, minmax(150px, 1fr)); gap: 0.5rem; font-size: 0.9rem;">
                                        <div>🏷️ 唯一值: <strong>{unique_count}</strong> 种</div>
                                        <div>🥇 最常见: <strong>{top_val}</strong></div>
                                        <div>📊 占比: <strong>{top_pct:.1f}%</strong></div>
                                        <div>{dist_icon} {dist_desc}</div>
                                    </div>
                                </div>
                                """, unsafe_allow_html=True)
                                
                                # 多选题分析功能
                                # 检测是否可能是多选题（包含分隔符的答案，或者有类似 "A. xxx,B. xxx" 的格式）
                                sample_values = df[col_select].dropna().astype(str).head(100)
                                # 检测是否有多选格式（字母+点号开头的多个选项，或普通分隔符）
                                has_letter_format = sample_values.str.contains(r'[A-Za-z]\.\s*[^,]+,[A-Za-z]\.', regex=True).any()
                                has_separator = sample_values.str.contains(r'[,;，；、\|]', regex=True).any()
                                
                                if has_separator or has_letter_format:
                                    with st.expander("🔀 多选题深度分析", expanded=False):
                                        st.markdown("""
                                        <div style="background: linear-gradient(135deg, #eff6ff 0%, #dbeafe 100%); padding: 1rem; border-radius: 8px; margin-bottom: 1rem; border-left: 4px solid #3b82f6;">
                                            <div style="font-weight: 600; color: #1e40af; font-size: 0.9rem;">📊 检测到多选题格式</div>
                                            <div style="color: #1e3a8a; font-size: 0.8rem; margin-top: 0.25rem;">系统将自动拆分答案，分析各选项的选择情况和组合关系</div>
                                        </div>
                                        """, unsafe_allow_html=True)
                                        
                                        # 智能检测分隔模式
                                        import re
                                        
                                        # 检测是否是 "A. xxx,B. xxx" 格式
                                        sample_str = str(sample_values.iloc[0]) if len(sample_values) > 0 else ""
                                        is_letter_option_format = bool(re.search(r'[A-Za-z]\.\s*[^,]+,[A-Za-z]\.', sample_str))
                                        
                                        sep_col1, sep_col2 = st.columns([1, 3])
                                        with sep_col1:
                                            split_mode = st.radio(
                                                "拆分模式",
                                                ["智能识别(推荐)", "按分隔符拆分"],
                                                index=0,
                                                key=f"split_mode_{col_select}_{idx}",
                                                help="智能识别会自动检测'A. xxx,B. xxx'格式"
                                            )
                                        
                                        with sep_col2:
                                            if split_mode == "按分隔符拆分":
                                                separator = st.selectbox(
                                                    "选择分隔符",
                                                    [",", ";", "，", "；", "、", "|"],
                                                    index=0,
                                                    key=f"multi_sep_{col_select}_{idx}",
                                                    help="选择用于分隔多选答案的符号"
                                                )
                                            else:
                                                st.info("智能模式将自动识别 'A. xxx,B. xxx' 或 '选项1,选项2' 格式")
                                        
                                        # 拆分多选答案
                                        all_choices = []
                                        respondent_choices = []  # 每个受访者的选择列表
                                        
                                        def smart_split_options(val_str):
                                            """智能拆分多选答案，正确处理括号内的逗号"""
                                            # 方法：先用正则找到所有 "字母. " 开头的选项位置，然后按位置切分
                                            # 这样可以正确处理括号内的逗号，如 "(e.g. iPhone 17 Pro, Sony A7 IV)"
                                            
                                            # 查找所有选项的起始位置（匹配 ",A. " 或字符串开头的 "A. "）
                                            option_pattern = r'(?:^|,\s*)([A-Za-z])\.\s*'
                                            
                                            # 找到所有匹配的位置
                                            matches = list(re.finditer(option_pattern, val_str))
                                            
                                            if len(matches) >= 1:
                                                options = []
                                                for i, match in enumerate(matches):
                                                    # 选项开始位置（从字母开始）
                                                    start = match.start(1)
                                                    # 选项结束位置（下一个选项开始前的逗号，或字符串结尾）
                                                    if i + 1 < len(matches):
                                                        # 下一个匹配的完整起始位置（包括逗号）
                                                        end = matches[i + 1].start()
                                                        # 去掉末尾的逗号
                                                        option_text = val_str[start:end].rstrip(',').strip()
                                                    else:
                                                        option_text = val_str[start:].strip()
                                                    
                                                    if option_text:
                                                        options.append(option_text)
                                                
                                                if options:
                                                    return options
                                            
                                            # 如果没有找到字母选项格式，尝试按逗号分隔（但要排除括号内的逗号）
                                            # 使用简单方法：如果没有括号，直接按逗号分隔
                                            if '(' not in val_str and '（' not in val_str:
                                                return [c.strip() for c in val_str.split(',') if c.strip()]
                                            
                                            # 有括号的情况，手动解析，跳过括号内的逗号
                                            result = []
                                            current = ""
                                            depth = 0  # 括号深度
                                            for char in val_str:
                                                if char in '(（[':
                                                    depth += 1
                                                    current += char
                                                elif char in ')）]':
                                                    depth -= 1
                                                    current += char
                                                elif char == ',' and depth == 0:
                                                    if current.strip():
                                                        result.append(current.strip())
                                                    current = ""
                                                else:
                                                    current += char
                                            if current.strip():
                                                result.append(current.strip())
                                            
                                            return result if result else [val_str.strip()]
                                        
                                        for val in df[col_select].dropna():
                                            val_str = str(val)
                                            
                                            if split_mode == "智能识别(推荐)":
                                                choices = smart_split_options(val_str)
                                            else:
                                                choices = [c.strip() for c in val_str.split(separator) if c.strip()]
                                            
                                            all_choices.extend(choices)
                                            if choices:
                                                respondent_choices.append(set(choices))
                                        
                                        # 统计各选项被选次数
                                        from collections import Counter
                                        choice_counts = Counter(all_choices)
                                        total_respondents = len(respondent_choices)
                                        
                                        # 创建选项频次表
                                        choice_df = pd.DataFrame([
                                            {'选项': k, '选择人数': v, '选择率': f"{v/total_respondents*100:.1f}%"} 
                                            for k, v in choice_counts.most_common()
                                        ])
                                        
                                        st.markdown("#### 📈 各选项选择情况")
                                        st.markdown(f"**总样本数:** {total_respondents} 人")
                                        
                                        # 选项选择率柱状图
                                        if len(choice_df) > 0:
                                            choice_df['选择率数值'] = choice_df['选择人数'] / total_respondents * 100
                                            
                                            fig_choice = go.Figure()
                                            fig_choice.add_trace(go.Bar(
                                                x=choice_df['选项'],
                                                y=choice_df['选择人数'],
                                                text=choice_df.apply(lambda r: f"{r['选择人数']}<br>({r['选择率数值']:.1f}%)", axis=1),
                                                textposition='outside',
                                                marker_color=px.colors.qualitative.Set2[:len(choice_df)],
                                                hovertemplate='<b>%{x}</b><br>选择人数: %{y}<br>选择率: %{text}<extra></extra>'
                                            ))
                                            
                                            # 处理长标签
                                            max_label_len = max(len(str(x)) for x in choice_df['选项'])
                                            if max_label_len > 30:
                                                # 标签太长，截断显示
                                                short_labels = [str(x)[:25] + '...' if len(str(x)) > 25 else str(x) for x in choice_df['选项']]
                                                fig_choice.update_traces(x=short_labels)
                                            
                                            fig_choice.update_layout(
                                                title="各选项被选择次数（可多选）",
                                                xaxis_title="",
                                                yaxis_title="选择人数",
                                                height=500,
                                                plot_bgcolor='rgba(0,0,0,0)',
                                                paper_bgcolor='rgba(0,0,0,0)',
                                                xaxis_tickangle=-45,
                                                margin=dict(b=180, t=50, l=50, r=20),
                                                xaxis=dict(tickfont=dict(size=9))
                                            )
                                            st.plotly_chart(fig_choice, use_container_width=True)
                                            
                                            # 显示详细数据表
                                            st.dataframe(choice_df[['选项', '选择人数', '选择率']], use_container_width=True, hide_index=True)
                                        
                                        # 选项组合分析
                                        st.markdown("---")
                                        st.markdown("#### 🔗 选项组合分析")
                                        st.markdown("分析哪些选项经常被一起选择")
                                        
                                        # 获取所有唯一选项
                                        unique_choices = list(choice_counts.keys())
                                        
                                        if len(unique_choices) >= 2:
                                            # 让用户选择要分析的选项
                                            combo_col1, combo_col2 = st.columns(2)
                                            with combo_col1:
                                                selected_option = st.selectbox(
                                                    "选择一个选项查看组合情况",
                                                    unique_choices,
                                                    key=f"combo_option_{col_select}_{idx}"
                                                )
                                            
                                            # 计算选择了该选项的人中，同时选择了其他选项的比例
                                            selected_respondents = [r for r in respondent_choices if selected_option in r]
                                            selected_count = len(selected_respondents)
                                            
                                            if selected_count > 0:
                                                st.markdown(f"**选择了「{selected_option}」的人数:** {selected_count} 人 ({selected_count/total_respondents*100:.1f}%)")
                                                
                                                # 统计这些人同时选择的其他选项
                                                co_selection_counts = Counter()
                                                for r in selected_respondents:
                                                    for choice in r:
                                                        if choice != selected_option:
                                                            co_selection_counts[choice] += 1
                                                
                                                if co_selection_counts:
                                                    co_df = pd.DataFrame([
                                                        {
                                                            '同时选择的选项': k, 
                                                            '人数': v, 
                                                            '占比': f"{v/selected_count*100:.1f}%",
                                                            '占比数值': v/selected_count*100
                                                        } 
                                                        for k, v in co_selection_counts.most_common()
                                                    ])
                                                    
                                                    # 组合选择率柱状图
                                                    fig_combo = go.Figure()
                                                    fig_combo.add_trace(go.Bar(
                                                        x=co_df['同时选择的选项'],
                                                        y=co_df['人数'],
                                                        text=co_df.apply(lambda r: f"{r['人数']}<br>({r['占比数值']:.1f}%)", axis=1),
                                                        textposition='outside',
                                                        marker_color='#3b82f6',
                                                        hovertemplate='<b>%{x}</b><br>人数: %{y}<extra></extra>'
                                                    ))
                                                    
                                                    # 处理长标签
                                                    short_labels = [str(x)[:20] + '...' if len(str(x)) > 20 else str(x) for x in co_df['同时选择的选项']]
                                                    fig_combo.update_traces(x=short_labels)
                                                    
                                                    fig_combo.update_layout(
                                                        title=f"选择「{selected_option[:20]}...」的人同时还选了什么？" if len(selected_option) > 20 else f"选择「{selected_option}」的人同时还选了什么？",
                                                        xaxis_title="",
                                                        yaxis_title="人数",
                                                        height=400,
                                                        plot_bgcolor='rgba(0,0,0,0)',
                                                        paper_bgcolor='rgba(0,0,0,0)',
                                                        xaxis_tickangle=-45,
                                                        margin=dict(b=150, t=60, l=50, r=20),
                                                        xaxis=dict(tickfont=dict(size=9)),
                                                        uniformtext_minsize=8,
                                                        uniformtext_mode='hide'
                                                    )
                                                    st.plotly_chart(fig_combo, use_container_width=True)
                                                    
                                                    st.dataframe(co_df[['同时选择的选项', '人数', '占比']], use_container_width=True, hide_index=True)
                                                else:
                                                    st.info("选择该选项的人都没有选择其他选项")
                                            
                                            # 热门组合分析
                                            st.markdown("---")
                                            st.markdown("#### 🔥 热门选项组合 TOP 10")
                                            
                                            # 统计所有两两组合的频次
                                            from itertools import combinations
                                            combo_counts = Counter()
                                            
                                            for r in respondent_choices:
                                                if len(r) >= 2:
                                                    for combo in combinations(sorted(r), 2):
                                                        combo_counts[combo] += 1
                                            
                                            if combo_counts:
                                                top_combos = combo_counts.most_common(10)
                                                combo_table = pd.DataFrame([
                                                    {
                                                        '选项组合': f"{c[0]} + {c[1]}",
                                                        '同时选择人数': count,
                                                        '占总样本比例': f"{count/total_respondents*100:.1f}%"
                                                    }
                                                    for c, count in top_combos
                                                ])
                                                
                                                st.dataframe(combo_table, use_container_width=True, hide_index=True)
                                                
                                                # 组合热力图（如果选项不太多）
                                                if len(unique_choices) <= 10:
                                                    st.markdown("#### 🗺️ 选项组合热力图")
                                                    
                                                    # 创建共现矩阵
                                                    matrix_data = []
                                                    for opt1 in unique_choices:
                                                        row = []
                                                        for opt2 in unique_choices:
                                                            if opt1 == opt2:
                                                                row.append(choice_counts[opt1])
                                                            else:
                                                                combo_key = tuple(sorted([opt1, opt2]))
                                                                row.append(combo_counts.get(combo_key, 0))
                                                        matrix_data.append(row)
                                                    
                                                    # 截断过长的标签
                                                    short_labels = [opt[:15] + '...' if len(opt) > 15 else opt for opt in unique_choices]
                                                    
                                                    fig_heatmap = go.Figure(data=go.Heatmap(
                                                        z=matrix_data,
                                                        x=short_labels,
                                                        y=short_labels,
                                                        colorscale='Blues',
                                                        hovertemplate='%{x} + %{y}<br>共同选择: %{z}人<extra></extra>'
                                                    ))
                                                    
                                                    fig_heatmap.update_layout(
                                                        title="选项共现热力图（对角线为单选人数）",
                                                        height=450,
                                                        xaxis_tickangle=-45
                                                    )
                                                    st.plotly_chart(fig_heatmap, use_container_width=True)
                                            else:
                                                st.info("没有发现两个及以上选项的组合")
                        
                            # 显示统计信息
                            with st.expander(f"📊 查看 {col_select} 详细统计数据"):
                                if pd.api.types.is_numeric_dtype(df[col_select]):
                                    stat_col1, stat_col2, stat_col3 = st.columns(3)
                                    
                                    with stat_col1:
                                        st.metric("平均值", f"{df[col_select].mean():.2f}")
                                        st.metric("最小值", f"{df[col_select].min():.2f}")
                                    
                                    with stat_col2:
                                        st.metric("中位数", f"{df[col_select].median():.2f}")
                                        st.metric("最大值", f"{df[col_select].max():.2f}")
                                    
                                    with stat_col3:
                                        st.metric("标准差", f"{df[col_select].std():.2f}")
                                        st.metric("有效样本", f"{df[col_select].count()}")
                                else:
                                    freq_df = df[col_select].value_counts().reset_index()
                                    freq_df.columns = [col_select, '频次']
                                    freq_df['占比'] = (freq_df['频次'] / freq_df['频次'].sum() * 100).round(2).astype(str) + '%'
                                    freq_df['累计占比'] = (freq_df['频次'].cumsum() / freq_df['频次'].sum() * 100).round(2).astype(str) + '%'
                                    st.dataframe(freq_df, use_container_width=True, hide_index=True)
                            
                            if idx < len(col_select_list) - 1:
                                st.markdown("---")  # 分隔线
                    
                    # 导出功能区域 - shadcn UI 风格
                    st.markdown("---")
                    st.markdown("""
                    <div style="margin: 1.5rem 0 1rem 0;">
                        <div style="display: flex; align-items: center; gap: 0.5rem; margin-bottom: 0.75rem;">
                            <div style="width: 2rem; height: 2rem; background: linear-gradient(135deg, #18181b 0%, #3f3f46 100%); border-radius: 0.5rem; display: flex; align-items: center; justify-content: center;">
                                <span style="color: white; font-size: 0.875rem;">📤</span>
                            </div>
                            <div>
                                <h3 style="margin: 0; font-size: 1rem; font-weight: 600; color: #18181b;">导出分析报告</h3>
                                <p style="margin: 0; font-size: 0.75rem; color: #71717a;">将当前选中的 {len(col_select_list)} 个变量分析结果导出为文档</p>
                            </div>
                        </div>
                    </div>
                    """.replace("{len(col_select_list)}", str(len(col_select_list))), unsafe_allow_html=True)
                    
                    # 导出卡片容器
                    st.markdown("""
                    <style>
                    .export-card {
                        background: #fafafa;
                        border: 1px solid #e4e4e7;
                        border-radius: 0.75rem;
                        padding: 1.25rem;
                        margin-bottom: 1rem;
                    }
                    .export-option-group {
                        display: flex;
                        gap: 0.5rem;
                        margin-bottom: 1rem;
                    }
                    .export-option {
                        flex: 1;
                        padding: 0.75rem 1rem;
                        background: white;
                        border: 1px solid #e4e4e7;
                        border-radius: 0.5rem;
                        cursor: pointer;
                        transition: all 0.2s ease;
                        text-align: center;
                    }
                    .export-option:hover {
                        border-color: #18181b;
                        background: #f4f4f5;
                    }
                    .export-option.selected {
                        border-color: #18181b;
                        background: #18181b;
                        color: white;
                    }
                    .export-option-icon {
                        font-size: 1.5rem;
                        margin-bottom: 0.25rem;
                    }
                    .export-option-label {
                        font-size: 0.75rem;
                        font-weight: 500;
                    }
                    </style>
                    """, unsafe_allow_html=True)
                    
                    export_col1, export_col2 = st.columns([3, 1])
                    
                    with export_col1:
                        export_format = st.selectbox(
                            "📄 导出格式",
                            ["Word 文档 (.docx)", "PPT 演示文稿 (.pptx)", "PDF 文档 (.pdf)"],
                            key="single_var_export_format",
                            help="选择报告导出的文件格式"
                        )
                        include_charts = st.checkbox("📊 包含图表图片", value=True, key="single_var_include_charts", 
                                                    help="勾选后报告中将包含可视化图表")
                    
                    with export_col2:
                        st.markdown("<div style='height: 0.5rem;'></div>", unsafe_allow_html=True)
                        export_btn = st.button("📥 生成报告", key="single_var_export_btn", use_container_width=True, type="primary")
                    
                    if export_btn:
                        with st.spinner("正在生成报告，请稍候..."):
                            try:
                                # 显示 kaleido 状态
                                if include_charts:
                                    if KALEIDO_AVAILABLE:
                                        st.info("📊 正在生成图表图片...")
                                    else:
                                        st.warning("⚠️ kaleido 库不可用，报告将不包含图表图片")
                                
                                # 自动生成所有选中变量的导出数据
                                sections = generate_all_export_sections(df, col_select_list)
                                title = "用户调研分析报告"
                                
                                # 统计图表生成情况
                                charts_count = sum(1 for s in sections if s.get('chart_image') is not None)
                                if include_charts:
                                    st.info(f"📈 成功生成 {charts_count}/{len(sections)} 个图表")
                                
                                if len(sections) == 0:
                                    st.warning("没有可导出的内容")
                                else:
                                    # 生成报告
                                    report_bytes = None
                                    file_ext = ""
                                    mime_type = ""
                                    format_name = ""
                                    format_icon = ""
                                    
                                    if "Word" in export_format:
                                        report_bytes = create_word_report(title, sections, include_charts)
                                        file_ext = "docx"
                                        mime_type = "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
                                        format_name = "Word"
                                        format_icon = "📝"
                                    elif "PPT" in export_format:
                                        if PPTX_AVAILABLE:
                                            report_bytes = create_ppt_report(title, sections, include_charts)
                                            file_ext = "pptx"
                                            mime_type = "application/vnd.openxmlformats-officedocument.presentationml.presentation"
                                            format_name = "PPT"
                                            format_icon = "📊"
                                        else:
                                            st.error("PPT 导出需要安装 python-pptx 库")
                                    else:  # PDF
                                        if PDF_AVAILABLE:
                                            report_bytes = create_pdf_report(title, sections, include_charts)
                                            file_ext = "pdf"
                                            mime_type = "application/pdf"
                                            format_name = "PDF"
                                            format_icon = "📄"
                                        else:
                                            st.error("PDF 导出需要安装 reportlab 库")
                                    
                                    if report_bytes:
                                        # shadcn 风格的成功提示卡片
                                        st.markdown(f"""
                                        <div style="
                                            background: linear-gradient(135deg, #f0fdf4 0%, #dcfce7 100%);
                                            border: 1px solid #86efac;
                                            border-radius: 0.75rem;
                                            padding: 1.25rem;
                                            margin: 1rem 0;
                                        ">
                                            <div style="display: flex; align-items: flex-start; gap: 0.75rem;">
                                                <div style="
                                                    width: 2.5rem; height: 2.5rem;
                                                    background: #22c55e;
                                                    border-radius: 50%;
                                                    display: flex; align-items: center; justify-content: center;
                                                    flex-shrink: 0;
                                                ">
                                                    <span style="color: white; font-size: 1.25rem;">✓</span>
                                                </div>
                                                <div style="flex: 1;">
                                                    <h4 style="margin: 0 0 0.25rem 0; color: #166534; font-size: 0.95rem; font-weight: 600;">
                                                        报告生成成功
                                                    </h4>
                                                    <p style="margin: 0; color: #15803d; font-size: 0.8rem;">
                                                        已生成包含 <strong>{len(sections)}</strong> 个分析问题的 {format_name} 报告
                                                        {'（含图表）' if include_charts else '（纯文本）'}
                                                    </p>
                                                </div>
                                            </div>
                                        </div>
                                        """, unsafe_allow_html=True)
                                        
                                        # 下载按钮
                                        st.download_button(
                                            label=f"{format_icon} 下载 {format_name} 报告",
                                            data=report_bytes,
                                            file_name=f"调研报告_{datetime.now().strftime('%Y%m%d_%H%M%S')}.{file_ext}",
                                            mime=mime_type,
                                            key=f"download_{file_ext}_single",
                                            use_container_width=True
                                        )
                            except Exception as e:
                                st.markdown(f"""
                                <div style="
                                    background: linear-gradient(135deg, #fef2f2 0%, #fee2e2 100%);
                                    border: 1px solid #fca5a5;
                                    border-radius: 0.75rem;
                                    padding: 1.25rem;
                                    margin: 1rem 0;
                                ">
                                    <div style="display: flex; align-items: flex-start; gap: 0.75rem;">
                                        <div style="
                                            width: 2.5rem; height: 2.5rem;
                                            background: #ef4444;
                                            border-radius: 50%;
                                            display: flex; align-items: center; justify-content: center;
                                            flex-shrink: 0;
                                        ">
                                            <span style="color: white; font-size: 1.25rem;">✕</span>
                                        </div>
                                        <div style="flex: 1;">
                                            <h4 style="margin: 0 0 0.25rem 0; color: #991b1b; font-size: 0.95rem; font-weight: 600;">
                                                导出失败
                                            </h4>
                                            <p style="margin: 0; color: #b91c1c; font-size: 0.8rem;">
                                                {str(e)}
                                            </p>
                                        </div>
                                    </div>
                                </div>
                                """, unsafe_allow_html=True)
                                import traceback
                                with st.expander("查看详细错误信息"):
                                    st.code(traceback.format_exc())

            # ==================== Cross Analysis 页面 ====================
            elif current_page == 'cross_analysis':
                st.markdown("#### 🔀 多维交叉分析")
                st.markdown("选择变量组合，探索数据之间的关联关系")
                
                # ========== 智能推荐分析 ==========
                st.markdown("##### 💡 智能推荐分析")
                
                # 基于数据特征生成推荐
                recommendations = []
                numeric_cols = df.select_dtypes(include=['number']).columns.tolist()
                categorical_cols = df.select_dtypes(include=['object', 'category']).columns.tolist()
                
                # 推荐1: 如果有问题列，推荐问题之间的交叉
                question_cols = [c for c in df.columns if c.startswith('Q') or '?' in c or '？' in c]
                if len(question_cols) >= 2:
                    recommendations.append({
                        'name': f"🎯 {question_cols[0][:30]}... × {question_cols[1][:30]}...",
                        'x': question_cols[0],
                        'y': question_cols[1],
                        'chart': '🔥 热力图',
                        'desc': '查看两个问题答案的交叉分布'
                    })
                
                # 推荐2: 分类变量 × 数值变量 (箱线图)
                if categorical_cols and numeric_cols:
                    cat_col = categorical_cols[0]
                    num_col = numeric_cols[0]
                    recommendations.append({
                        'name': f"📦 {cat_col[:20]}... 下的 {num_col[:20]}... 分布",
                        'x': cat_col,
                        'y': num_col,
                        'chart': '📦 箱线图',
                        'desc': '对比不同类别的数值分布差异'
                    })
                
                # 推荐3: 时间相关分析
                time_cols = [c for c in df.columns if '时间' in c or 'time' in c.lower() or 'date' in c.lower()]
                if time_cols and categorical_cols:
                    recommendations.append({
                        'name': f"📊 {time_cols[0][:20]}... 趋势分析",
                        'x': time_cols[0],
                        'y': categorical_cols[0] if categorical_cols else df.columns[1],
                        'chart': '📊 柱状图',
                        'desc': '查看数据随时间的变化趋势'
                    })
                
                # 推荐4: 状态/完成度分析
                status_cols = [c for c in df.columns if 'state' in c.lower() or 'status' in c.lower() or '状态' in c]
                if status_cols and len(df.columns) > 2:
                    other_col = [c for c in df.columns if c not in status_cols][0]
                    recommendations.append({
                        'name': f"🎯 {status_cols[0][:20]}... × {other_col[:20]}...",
                        'x': status_cols[0],
                        'y': other_col,
                        'chart': '🔥 热力图',
                        'desc': '分析不同状态下的数据分布'
                    })
                
                # 如果没有足够推荐，添加默认推荐
                if len(recommendations) < 2 and len(df.columns) >= 2:
                    recommendations.append({
                        'name': f"🔥 {df.columns[0][:20]}... × {df.columns[1][:20]}...",
                        'x': df.columns[0],
                        'y': df.columns[1],
                        'chart': '🔥 热力图',
                        'desc': '查看两个变量的交叉分布'
                    })
                
                # 显示推荐卡片（静态显示，不使用按钮避免状态循环）
                if recommendations:
                    st.markdown("**推荐分析组合:**")
                    for i, rec in enumerate(recommendations[:3]):
                        st.markdown(f"""
                        <div style="background: #f0f4ff; padding: 0.6rem 0.8rem; border-radius: 6px; margin-bottom: 0.4rem; border-left: 3px solid #667eea;">
                            <strong style="color: #333; font-size: 0.85rem;">{rec['name'][:40]}...</strong>
                            <div style="font-size: 0.75rem; color: #6c757d; margin-top: 0.2rem;">{rec['desc']}</div>
                        </div>
                        """, unsafe_allow_html=True)
                
                st.markdown("---")
                
                # ========== 手动配置区域 ==========
                st.markdown("##### ⚙️ 自定义分析参数")
                
                # 直接使用固定默认值，不从 session_state 读取
                chart_options = ["📍 散点图", "📦 箱线图", "🔥 热力图", "📊 柱状图"]
                x_index = 0
                y_index = 1 if len(df.columns) > 1 else 0
                chart_index = 2  # 默认热力图
                
                form_col1, form_col2, form_col3 = st.columns(3)
                with form_col1:
                    new_x_axis = st.selectbox(
                        "📊 X 轴 (自变量)", 
                        df.columns.tolist(), 
                        index=x_index, 
                        key="cross_x_manual"
                    )
                with form_col2:
                    new_y_axis = st.selectbox(
                        "📈 Y 轴 (因变量)", 
                        df.columns.tolist(), 
                        index=y_index, 
                        key="cross_y_manual"
                    )
                with form_col3:
                    new_chart_type = st.selectbox(
                        "🎨 图表类型", 
                        chart_options,
                        index=chart_index,
                        key="cross_chart_manual"
                    )
                
                # 可选参数
                opt_col1, opt_col2 = st.columns(2)
                with opt_col1:
                    color_options = ["无"] + df.columns.tolist()
                    new_color_col = st.selectbox(
                        "🎨 颜色分组 (可选)", 
                        color_options,
                        index=0,
                        help="为数据点添加颜色分组",
                        key="cross_color_manual"
                    )
                    if new_color_col == "无":
                        new_color_col = None
                
                with opt_col2:
                    new_color_scheme = st.selectbox(
                        "🎨 配色方案",
                        ["蓝色 Blues", "紫色 Purples", "绿色 Greens", "橙色 Oranges", 
                         "红色 Reds", "粉色 Pinkyl", "青色 Teal", "彩虹 Rainbow"],
                        key="cross_color_scheme_manual"
                    )
                
                # 聚合方式（仅柱状图需要）
                new_agg_func = "🔢 计数"
                if "柱状图" in new_chart_type:
                    y_is_numeric = pd.api.types.is_numeric_dtype(df[new_y_axis])
                    if y_is_numeric:
                        new_agg_func = st.radio(
                            "📊 聚合方式", 
                            ["📊 平均值", "➕ 总和", "🔢 计数"], 
                            horizontal=True,
                            key="cross_agg_manual"
                        )
                
                st.markdown("---")
                
                # 直接生成图表（实时响应，无需点击按钮）
                
                # 颜色方案映射
                color_map = {
                    "紫色 Purples": ("Purples", px.colors.qualitative.Pastel),
                    "蓝色 Blues": ("Blues", px.colors.qualitative.Safe),
                    "绿色 Greens": ("Greens", px.colors.qualitative.Prism),
                    "橙色 Oranges": ("Oranges", px.colors.qualitative.Bold),
                    "红色 Reds": ("Reds", px.colors.qualitative.Vivid),
                    "粉色 Pinkyl": ("Pinkyl", px.colors.qualitative.Pastel),
                    "青色 Teal": ("Teal", px.colors.qualitative.Set2),
                    "彩虹 Rainbow": ("Rainbow", px.colors.qualitative.Vivid)
                }
                
                color_scale, color_discrete_seq = color_map.get(new_color_scheme, ("Purples", px.colors.qualitative.Pastel))
                
                try:
                    fig = None
                    agg_label = "计数"
                    
                    if "散点图" in new_chart_type:
                        st.info("💡 **适用场景**: 查看两个数值变量的相关性")
                        fig = px.scatter(
                            df, x=new_x_axis, y=new_y_axis, 
                            color=new_color_col if new_color_col else None, 
                            title=f"📍 {new_x_axis} 与 {new_y_axis} 散点图",
                            color_discrete_sequence=color_discrete_seq
                        )
                        fig.update_traces(marker=dict(size=8, opacity=0.7))
                    
                    elif "箱线图" in new_chart_type:
                        st.info("💡 **适用场景**: 对比不同类别的数值分布")
                        fig = px.box(
                            df, x=new_x_axis, y=new_y_axis, 
                            color=new_color_col if new_color_col else None, 
                            title=f"📦 {new_x_axis} 下的 {new_y_axis} 分布",
                            color_discrete_sequence=color_discrete_seq
                        )

                    elif "柱状图" in new_chart_type:
                        st.info("💡 **适用场景**: 对比不同类别的数值总和或平均值")
                        
                        y_is_numeric = pd.api.types.is_numeric_dtype(df[new_y_axis])
                        
                        if y_is_numeric and new_agg_func and "平均值" in str(new_agg_func):
                            grouped = df.groupby(new_x_axis)[new_y_axis].mean().reset_index()
                            agg_label = "平均值"
                            y_axis_display = new_y_axis
                        elif y_is_numeric and new_agg_func and "总和" in str(new_agg_func):
                            grouped = df.groupby(new_x_axis)[new_y_axis].sum().reset_index()
                            agg_label = "总和"
                            y_axis_display = new_y_axis
                        else:
                            grouped = df.groupby(new_x_axis).size().reset_index(name='count')
                            y_axis_display = 'count'
                            agg_label = "计数"
                        
                        fig = px.bar(
                            grouped, x=new_x_axis, y=y_axis_display, 
                            title=f"📊 {new_x_axis} vs {new_y_axis} ({agg_label})",
                            color=y_axis_display,
                            color_continuous_scale=color_scale
                        )

                    elif "热力图" in new_chart_type:
                        st.info("💡 **适用场景**: 查看两个分类变量的交叉密度")
                        crosstab = pd.crosstab(df[new_x_axis], df[new_y_axis])
                        fig = px.imshow(
                            crosstab, 
                            text_auto=True, 
                            title=f"🔥 {new_x_axis} 与 {new_y_axis} 热力分布",
                            color_continuous_scale=color_scale,
                            aspect="auto"
                        )
                    
                    else:
                        st.warning(f"未识别的图表类型: {new_chart_type}，使用默认散点图")
                        fig = px.scatter(
                            df, x=new_x_axis, y=new_y_axis,
                            title=f"📍 {new_x_axis} 与 {new_y_axis} 散点图"
                        )

                    # 统一图表样式并显示
                    if fig is not None:
                        fig.update_layout(
                            plot_bgcolor='rgba(0,0,0,0)',
                            paper_bgcolor='rgba(0,0,0,0)',
                            font=dict(size=12),
                            title_font_size=16,
                            height=450
                        )
                        st.plotly_chart(fig, use_container_width=True, key="cross_chart_main")
                    
                    # 自动解读
                    insight_text = ""
                    if "散点图" in new_chart_type:
                        if pd.api.types.is_numeric_dtype(df[new_x_axis]) and pd.api.types.is_numeric_dtype(df[new_y_axis]):
                            corr = df[new_x_axis].corr(df[new_y_axis])
                            if abs(corr) > 0.7:
                                direction = "正" if corr > 0 else "负"
                                insight_text = f"🔗 两变量呈<strong>强{direction}相关</strong>(r={corr:.2f})"
                            elif abs(corr) > 0.4:
                                direction = "正" if corr > 0 else "负"
                                insight_text = f"🔗 两变量呈<strong>中等{direction}相关</strong>(r={corr:.2f})"
                            else:
                                insight_text = f"🔗 两变量相关性较弱(r={corr:.2f})"
                    elif "箱线图" in new_chart_type:
                        if pd.api.types.is_numeric_dtype(df[new_y_axis]):
                            groups = df.groupby(new_x_axis)[new_y_axis].agg(['mean', 'std'])
                            if not groups.empty:
                                max_group = groups['mean'].idxmax()
                                min_group = groups['mean'].idxmin()
                                insight_text = f"📊 <strong>{max_group}</strong>的{new_y_axis}均值最高,<strong>{min_group}</strong>最低"
                    elif "柱状图" in new_chart_type:
                        insight_text = f"📊 展示了{new_x_axis}各类别下{new_y_axis}的{agg_label}对比"
                    elif "热力图" in new_chart_type:
                        crosstab_temp = pd.crosstab(df[new_x_axis], df[new_y_axis])
                        if not crosstab_temp.empty:
                            max_cell = crosstab_temp.stack().idxmax()
                            insight_text = f"🔥 <strong>{max_cell[0]}</strong>与<strong>{max_cell[1]}</strong>组合出现频率最高"
                    
                    if insight_text:
                        st.markdown(f"""
                        <div style="background: linear-gradient(135deg, #f0f4ff 0%, #f5f0ff 100%); padding: 0.8rem 1rem; border-radius: 8px; margin-top: 0.5rem; border-left: 4px solid #667eea;">
                            <span style="font-weight: 600; color: #333;">💡 快速解读:</span> {insight_text}
                        </div>
                        """, unsafe_allow_html=True)
                    
                except Exception as e:
                    st.error(f"❌ 生成图表时出错: {e}")

            # ==================== AI Summary 页面 ====================
            elif current_page == 'ai_summary':
                # 差异化价值展示
                st.markdown("""
                <div style="background: linear-gradient(135deg, #667eea 0%, #764ba2 100%); padding: 1.5rem; border-radius: 15px; color: white; margin-bottom: 2rem; box-shadow: 0 4px 20px rgba(102, 126, 234, 0.3);">
                    <h3 style="margin: 0 0 1rem 0; font-size: 1.5rem;">🚀 数据驱动的AI分析</h3>
                    <div style="display: grid; grid-template-columns: repeat(auto-fit, minmax(200px, 1fr)); gap: 1rem; margin-top: 1rem;">
                        <div style="background: rgba(255,255,255,0.15); padding: 1rem; border-radius: 10px; backdrop-filter: blur(10px);">
                            <div style="font-size: 1.5rem; margin-bottom: 0.5rem;">📊</div>
                            <div style="font-weight: 600; margin-bottom: 0.3rem;">基于真实数据</div>
                            <div style="font-size: 0.85rem; opacity: 0.9;">分析基于您上传的实际数据,而非通用建议</div>
                        </div>
                        <div style="background: rgba(255,255,255,0.15); padding: 1rem; border-radius: 10px; backdrop-filter: blur(10px);">
                            <div style="font-size: 1.5rem; margin-bottom: 0.5rem;">🎯</div>
                            <div style="font-weight: 600; margin-bottom: 0.3rem;">智能问题推荐</div>
                            <div style="font-size: 0.85rem; opacity: 0.9;">根据数据特征自动生成针对性问题</div>
                        </div>
                        <div style="background: rgba(255,255,255,0.15); padding: 1rem; border-radius: 10px; backdrop-filter: blur(10px);">
                            <div style="font-size: 1.5rem; margin-bottom: 0.5rem;">📈</div>
                            <div style="font-weight: 600; margin-bottom: 0.3rem;">可视化+解读</div>
                            <div style="font-size: 0.85rem; opacity: 0.9;">图表与AI分析完美结合</div>
                        </div>
                        <div style="background: rgba(255,255,255,0.15); padding: 1rem; border-radius: 10px; backdrop-filter: blur(10px);">
                            <div style="font-size: 1.5rem; margin-bottom: 0.5rem;">⚡</div>
                            <div style="font-weight: 600; margin-bottom: 0.3rem;">一键生成报告</div>
                            <div style="font-size: 0.85rem; opacity: 0.9;">结构化的完整分析报告</div>
                        </div>
                    </div>
                </div>
                """, unsafe_allow_html=True)
                
                # 初始化对话历史
                if 'chat_history' not in st.session_state:
                    st.session_state.chat_history = []
                
                # 三栏布局
                col_insights, col_actions, col_chat = st.columns([1, 1, 1.5])
                
                with col_insights:
                    st.markdown("##### 🔍 智能洞察")
                    
                    # 自动检测数据特征并生成洞察
                    insights = []
                    
                    # 检测1: 样本量
                    sample_size = df.shape[0]
                    if sample_size < 30:
                        insights.append(("⚠️", "样本量较小", f"当前仅{sample_size}条数据,建议增加样本"))
                    elif sample_size > 10000:
                        insights.append(("✅", "大样本数据", f"{sample_size:,}条数据,统计结果可靠"))
                    
                    # 检测2: 缺失值
                    missing_pct = (df.isnull().sum().sum() / (df.shape[0] * df.shape[1])) * 100
                    if missing_pct > 10:
                        insights.append(("❌", "缺失值较多", f"{missing_pct:.1f}%的数据缺失"))
                    elif missing_pct > 0:
                        insights.append(("⚠️", "存在缺失值", f"{missing_pct:.1f}%的数据缺失"))
                    else:
                        insights.append(("✅", "数据完整", "无缺失值"))
                    
                    # 检测3: 变量类型
                    num_vars = df.select_dtypes(include=['number']).shape[1]
                    cat_vars = df.select_dtypes(include=['object']).shape[1]
                    if num_vars > 0 and cat_vars > 0:
                        insights.append(("💡", "混合数据", f"{num_vars}个数值+{cat_vars}个类别变量"))
                    
                    # 检测4: 数据分布
                    if num_vars > 0:
                        numeric_cols = df.select_dtypes(include=['number']).columns
                        for col in numeric_cols[:2]:
                            skew = df[col].skew()
                            if abs(skew) > 1:
                                insights.append(("📊", f"{col}分布偏斜", f"偏度={skew:.2f}"))
                    
                    # 显示洞察卡片
                    for emoji, title, desc in insights[:5]:
                        st.markdown(f"""
                        <div style="background: white; padding: 0.8rem; border-radius: 8px; margin-bottom: 0.5rem; border-left: 3px solid #667eea; box-shadow: 0 2px 4px rgba(0,0,0,0.05);">
                            <div style="font-weight: 600; color: #333; margin-bottom: 0.2rem;">{emoji} {title}</div>
                            <div style="font-size: 0.85rem; color: #6c757d;">{desc}</div>
                        </div>
                        """, unsafe_allow_html=True)
                
                with col_actions:
                    st.markdown("##### ⚡ 快速操作")
                    
                    # 智能问题推荐(基于数据特征) - 静态显示，不使用按钮
                    st.markdown("**💡 智能问题推荐**")
                    st.markdown('<p style="font-size: 0.85rem; color: #6c757d; margin-bottom: 0.8rem;">基于您的数据特征生成</p>', unsafe_allow_html=True)
                    
                    smart_questions = []
                    
                    # 根据数据特征生成问题
                    if missing_pct > 5:
                        smart_questions.append("如何处理缺失值?")
                    
                    if num_vars >= 2:
                        smart_questions.append("哪些变量之间存在相关性?")
                    
                    if cat_vars > 0:
                        cat_col = df.select_dtypes(include=['object']).columns[0]
                        smart_questions.append(f"{cat_col}的分布有什么特点?")
                    
                    if num_vars > 0:
                        num_col = df.select_dtypes(include=['number']).columns[0]
                        smart_questions.append(f"{num_col}有异常值吗?")
                    
                    smart_questions.append("有什么业务建议?")
                    
                    # 显示推荐问题列表（静态显示）
                    for i, question in enumerate(smart_questions[:5]):
                        st.markdown(f"""
                        <div style="background: #f8f9fa; padding: 0.5rem 0.75rem; border-radius: 6px; margin-bottom: 0.4rem; font-size: 0.85rem; color: #495057; border-left: 3px solid #667eea;">
                            💬 {question}
                        </div>
                        """, unsafe_allow_html=True)
                    
                    st.markdown("---")
                    st.caption("💡 复制上方问题到对话框中提问")
                
                with col_chat:
                    st.markdown("##### 💬 对话区")
                    
                    # 显示当前AI模式
                    ai_config = st.session_state.get('ai_config', {'provider': '规则分析(无需API)'})
                    if ai_config['provider'] == '规则分析(无需API)':
                        st.info("🤖 当前模式: 规则分析 | 💡 在侧边栏配置API使用真实AI")
                    else:
                        model_name = ai_config.get('model', 'Unknown')
                        st.success(f"🤖 当前模式: {ai_config['provider']} ({model_name})")
                    
                    st.markdown("---")
                    
                    # 显示对话历史
                    chat_container = st.container()
                    with chat_container:
                        if len(st.session_state.chat_history) == 0:
                            st.info("👋 你好!我是AI分析助手。你可以问我关于数据的任何问题!")
                        else:
                            for msg in st.session_state.chat_history:
                                if msg['role'] == 'user':
                                    st.markdown(f"""
                                    <div style="background: linear-gradient(135deg, #667eea 0%, #764ba2 100%); padding: 1rem; border-radius: 10px; margin: 0.5rem 0; color: white;">
                                        <strong>👤 你:</strong><br>{msg['content']}
                                    </div>
                                    """, unsafe_allow_html=True)
                                else:
                                    st.markdown(f"""
                                    <div style="background: #f8f9fa; padding: 1rem; border-radius: 10px; margin: 0.5rem 0; border-left: 4px solid #667eea;">
                                        <strong>🤖 AI:</strong><br>{msg['content']}
                                    </div>
                                    """, unsafe_allow_html=True)
                    
                    # 输入框 - 使用 form 避免刷新
                    st.markdown("---")
                    with st.form(key="chat_form", clear_on_submit=True):
                        user_input = st.text_area(
                            "输入你的问题:",
                            placeholder="例如: 这个数据集的主要特征是什么?",
                            height=80,
                            key="user_input_form"
                        )
                        
                        col_send, col_clear = st.columns([3, 1])
                        with col_send:
                            submit_btn = st.form_submit_button("📤 发送", use_container_width=True)
                        with col_clear:
                            clear_btn = st.form_submit_button("🗑️ 清空", use_container_width=True)
                        
                        # 在 form 内处理提交，避免状态循环
                        if submit_btn and user_input.strip():
                            st.session_state.chat_history.append({
                                'role': 'user',
                                'content': user_input.strip()
                            })
                            answer = generate_ai_response(user_input.strip(), df)
                            st.session_state.chat_history.append({
                                'role': 'assistant',
                                'content': answer
                            })
                        
                        if clear_btn:
                            st.session_state.chat_history = []

            # ==================== Segmentation Kanban 页面 ====================
            elif current_page == 'kanban':
                st.markdown("""
                <style>
                .kanban-container {
                    display: flex;
                    gap: 1rem;
                    overflow-x: auto;
                    padding: 1rem 0;
                }
                .kanban-column {
                    min-width: 300px;
                    max-width: 350px;
                    background: #f8f9fa;
                    border-radius: 12px;
                    padding: 1rem;
                }
                .kanban-header {
                    display: flex;
                    align-items: center;
                    gap: 0.5rem;
                    margin-bottom: 1rem;
                    padding-bottom: 0.75rem;
                    border-bottom: 2px solid #e9ecef;
                }
                .kanban-badge {
                    padding: 0.25rem 0.75rem;
                    border-radius: 20px;
                    font-size: 0.75rem;
                    font-weight: 600;
                }
                .kanban-badge.complete { background: #d4edda; color: #155724; }
                .kanban-badge.partial { background: #fff3cd; color: #856404; }
                .kanban-badge.disqualified { background: #f8d7da; color: #721c24; }
                .kanban-count {
                    margin-left: auto;
                    background: #e9ecef;
                    padding: 0.25rem 0.5rem;
                    border-radius: 10px;
                    font-size: 0.75rem;
                    font-weight: 600;
                }
                .kanban-card {
                    background: white;
                    border-radius: 8px;
                    padding: 1rem;
                    margin-bottom: 0.75rem;
                    border: 1px solid #e9ecef;
                    box-shadow: 0 1px 3px rgba(0,0,0,0.05);
                    transition: all 0.2s ease;
                }
                .kanban-card:hover {
                    box-shadow: 0 4px 12px rgba(0,0,0,0.1);
                    transform: translateY(-2px);
                }
                .kanban-card-title {
                    font-weight: 600;
                    font-size: 0.9rem;
                    color: #1a1a2e;
                    margin-bottom: 0.5rem;
                }
                .kanban-card-field {
                    font-size: 0.75rem;
                    color: #6c757d;
                    margin-bottom: 0.25rem;
                }
                .kanban-card-value {
                    font-size: 0.85rem;
                    color: #333;
                }
                .kanban-card-tag {
                    display: inline-block;
                    padding: 0.2rem 0.5rem;
                    border-radius: 4px;
                    font-size: 0.7rem;
                    font-weight: 500;
                    margin-right: 0.25rem;
                    margin-top: 0.5rem;
                }
                .tag-desktop { background: #e3f2fd; color: #1565c0; }
                .tag-mobile { background: #f3e5f5; color: #7b1fa2; }
                .tag-tablet { background: #e8f5e9; color: #2e7d32; }
                </style>
                """, unsafe_allow_html=True)
                
                # 检测状态列
                status_col = None
                for col in df.columns:
                    col_lower = str(col).lower()
                    if any(kw in col_lower for kw in ['status', '状态', 'state', 'completion']):
                        status_col = col
                        break
                
                if status_col is None:
                    # 如果没有状态列，创建一个基于完整度的虚拟状态
                    st.info("💡 未检测到状态字段，将基于数据完整度自动分组")
                    
                    # 计算每行的完整度
                    df_kanban = df.copy()
                    df_kanban['_completeness'] = df_kanban.notna().sum(axis=1) / len(df_kanban.columns) * 100
                    df_kanban['_status'] = df_kanban['_completeness'].apply(
                        lambda x: 'Complete' if x >= 90 else ('Partial' if x >= 50 else 'Incomplete')
                    )
                    status_col = '_status'
                else:
                    df_kanban = df.copy()
                
                # 获取唯一状态值
                status_values = df_kanban[status_col].fillna('Unknown').unique()
                
                # 状态颜色映射
                status_colors = {
                    'Complete': 'complete', 'complete': 'complete', '完成': 'complete', 'Completed': 'complete',
                    'Partial': 'partial', 'partial': 'partial', '部分': 'partial', 'Incomplete': 'partial',
                    'Disqualified': 'disqualified', 'disqualified': 'disqualified', '不合格': 'disqualified',
                    'Unknown': 'partial'
                }
                
                # Airtable 风格筛选器
                st.markdown("""
                <style>
                .filter-container {
                    display: flex;
                    gap: 0.5rem;
                    margin-bottom: 1.5rem;
                    flex-wrap: wrap;
                }
                .filter-btn {
                    padding: 0.5rem 1rem;
                    border: 1px solid #e4e4e7;
                    border-radius: 6px;
                    background: white;
                    font-size: 0.85rem;
                    cursor: pointer;
                    display: flex;
                    align-items: center;
                    gap: 0.5rem;
                }
                .filter-btn:hover {
                    border-color: #a1a1aa;
                }
                </style>
                """, unsafe_allow_html=True)
                
                # 显示独立的状态筛选器（类似 Airtable）
                st.markdown("##### 筛选条件")
                filter_col1, filter_col2, filter_col3 = st.columns(3)
                
                with filter_col1:
                    status_filter = st.selectbox(
                        "Status",
                        options=["全部"] + list(status_values),
                        key="kanban_status_filter_single"
                    )
                
                # 检测设备类型列
                device_col = None
                for col in df.columns:
                    col_lower = str(col).lower()
                    if any(kw in col_lower for kw in ['device', '设备', 'platform']):
                        device_col = col
                        break
                
                with filter_col2:
                    if device_col:
                        device_options = ["全部"] + list(df_kanban[device_col].dropna().unique())
                        device_filter = st.selectbox(
                            "Device Type",
                            options=device_options,
                            key="kanban_device_filter"
                        )
                    else:
                        device_filter = "全部"
                        st.selectbox("Device Type", options=["全部"], disabled=True, key="kanban_device_disabled")
                
                # 检测调研标题列
                survey_col = None
                for col in df.columns:
                    col_lower = str(col).lower()
                    if any(kw in col_lower for kw in ['survey', 'title', '调研', '标题']):
                        survey_col = col
                        break
                
                with filter_col3:
                    if survey_col:
                        survey_options = ["全部"] + list(df_kanban[survey_col].dropna().unique())
                        survey_filter = st.selectbox(
                            "Survey",
                            options=survey_options,
                            key="kanban_survey_filter"
                        )
                    else:
                        survey_filter = "全部"
                        st.selectbox("Survey", options=["全部"], disabled=True, key="kanban_survey_disabled")
                
                # 应用筛选
                filtered_df = df_kanban.copy()
                if status_filter != "全部":
                    filtered_df = filtered_df[filtered_df[status_col] == status_filter]
                if device_col and device_filter != "全部":
                    filtered_df = filtered_df[filtered_df[device_col] == device_filter]
                if survey_col and survey_filter != "全部":
                    filtered_df = filtered_df[filtered_df[survey_col] == survey_filter]
                
                # 获取筛选后的状态值
                if status_filter != "全部":
                    selected_statuses = [status_filter]
                else:
                    selected_statuses = list(filtered_df[status_col].dropna().unique())
                
                st.markdown("---")
                
                # 创建看板列
                kanban_cols = st.columns(len(selected_statuses) if selected_statuses else 1)
                
                for idx, status in enumerate(selected_statuses):
                    status_df = filtered_df[filtered_df[status_col] == status]
                    color_class = status_colors.get(str(status), 'partial')
                    
                    with kanban_cols[idx]:
                        # 列标题
                        st.markdown(f"""
                        <div class="kanban-header">
                            <span class="kanban-badge {color_class}">{status}</span>
                            <span class="kanban-count">{len(status_df)}</span>
                        </div>
                        """, unsafe_allow_html=True)
                        
                        # 卡片
                        for row_idx, (_, row) in enumerate(status_df.head(10).iterrows()):
                            # 获取显示字段
                            card_title = str(row.iloc[0]) if len(row) > 0 else "Record"
                            
                            # 构建卡片内容 - 使用更清晰的格式
                            card_html = f'<div class="kanban-card"><div class="kanban-card-title">{card_title[:40]}</div>'
                            
                            # 显示前4个字段
                            field_count = 0
                            for col_name in df.columns:
                                if col_name != status_col and col_name not in ['_completeness', '_status'] and pd.notna(row.get(col_name)):
                                    val = str(row[col_name])[:40]
                                    card_html += f'<div style="margin-bottom: 0.5rem;"><div class="kanban-card-field">{col_name}</div><div class="kanban-card-value">{val}</div></div>'
                                    field_count += 1
                                    if field_count >= 4:
                                        break
                            
                            card_html += '</div>'
                            st.markdown(card_html, unsafe_allow_html=True)
                        
                        if len(status_df) > 10:
                            st.caption(f"... 还有 {len(status_df) - 10} 条记录")

            # ==================== Trends Dashboard 页面 (Airtable 风格) ====================
            elif current_page == 'trends':
                # Airtable 风格的 KPI 卡片
                st.markdown("""
                <style>
                .trends-kpi-row {
                    display: grid;
                    grid-template-columns: repeat(4, 1fr);
                    gap: 1rem;
                    margin-bottom: 2rem;
                }
                .trends-kpi-card {
                    background: linear-gradient(135deg, #667eea 0%, #764ba2 100%);
                    border-radius: 12px;
                    padding: 1.25rem;
                    color: white;
                    position: relative;
                    overflow: hidden;
                }
                .trends-kpi-card::after {
                    content: '↗';
                    position: absolute;
                    top: 1rem;
                    right: 1rem;
                    font-size: 1rem;
                    opacity: 0.7;
                }
                .trends-kpi-card.green { background: linear-gradient(135deg, #11998e 0%, #38ef7d 100%); }
                .trends-kpi-card.blue { background: linear-gradient(135deg, #4facfe 0%, #00f2fe 100%); }
                .trends-kpi-card.orange { background: linear-gradient(135deg, #f093fb 0%, #f5576c 100%); }
                .trends-kpi-label {
                    font-size: 0.85rem;
                    opacity: 0.9;
                    margin-bottom: 0.5rem;
                }
                .trends-kpi-value {
                    font-size: 2.5rem;
                    font-weight: 700;
                    line-height: 1;
                }
                .trends-kpi-unit {
                    font-size: 1rem;
                    opacity: 0.8;
                    margin-left: 0.25rem;
                }
                </style>
                """, unsafe_allow_html=True)
                
                # 计算 KPI
                total_responses = len(df)
                completion_rate = (1 - df.isnull().sum().sum() / (df.shape[0] * df.shape[1])) * 100
                
                # 检测完成时间列
                time_col = None
                for col in df.columns:
                    col_lower = str(col).lower()
                    if any(kw in col_lower for kw in ['time', 'duration', '时间', '时长', 'minutes']):
                        if df[col].dtype in ['int64', 'float64']:
                            time_col = col
                            break
                
                avg_time = df[time_col].mean() if time_col else 0
                
                # 检测质量列
                quality_count = 0
                for col in df.columns:
                    col_lower = str(col).lower()
                    if any(kw in col_lower for kw in ['quality', '质量', 'good', 'high']):
                        quality_count = df[col].notna().sum()
                        break
                
                # KPI 卡片行
                kpi_col1, kpi_col2, kpi_col3, kpi_col4 = st.columns(4)
                
                with kpi_col1:
                    st.markdown(f"""
                    <div class="trends-kpi-card">
                        <div class="trends-kpi-label">Total Survey Responses</div>
                        <div class="trends-kpi-value">{total_responses}</div>
                    </div>
                    """, unsafe_allow_html=True)
                
                with kpi_col2:
                    st.markdown(f"""
                    <div class="trends-kpi-card green">
                        <div class="trends-kpi-label">Response Completion Rate</div>
                        <div class="trends-kpi-value">{completion_rate:.0f}<span class="trends-kpi-unit">%</span></div>
                    </div>
                    """, unsafe_allow_html=True)
                
                with kpi_col3:
                    st.markdown(f"""
                    <div class="trends-kpi-card blue">
                        <div class="trends-kpi-label">Average Completion Time</div>
                        <div class="trends-kpi-value">{avg_time:.0f}<span class="trends-kpi-unit">min</span></div>
                    </div>
                    """, unsafe_allow_html=True)
                
                with kpi_col4:
                    st.markdown(f"""
                    <div class="trends-kpi-card orange">
                        <div class="trends-kpi-label">High Quality Responses</div>
                        <div class="trends-kpi-value">{quality_count}</div>
                    </div>
                    """, unsafe_allow_html=True)
                
                # 图表区域
                chart_col1, chart_col2 = st.columns(2)
                
                with chart_col1:
                    st.markdown("##### Responses by Status")
                    # 检测状态列
                    status_col = None
                    for col in df.columns:
                        col_lower = str(col).lower()
                        if any(kw in col_lower for kw in ['status', '状态']):
                            status_col = col
                            break
                    
                    if status_col:
                        status_counts = df[status_col].value_counts()
                        fig_status = px.bar(
                            x=status_counts.index,
                            y=status_counts.values,
                            color=status_counts.index,
                            color_discrete_sequence=['#667eea', '#f5576c', '#ffc107']
                        )
                        fig_status.update_layout(
                            showlegend=False,
                            template='plotly_white',
                            xaxis_title='Status',
                            yaxis_title='Number of Responses',
                            height=300
                        )
                        st.plotly_chart(fig_status, use_container_width=True)
                    else:
                        # 基于完整度创建状态
                        completeness_per_row = df.notna().sum(axis=1) / len(df.columns) * 100
                        status_data = pd.cut(completeness_per_row, bins=[0, 50, 90, 100], labels=['Incomplete', 'Partial', 'Complete'])
                        status_counts = status_data.value_counts()
                        fig_status = px.bar(
                            x=status_counts.index,
                            y=status_counts.values,
                            color=status_counts.index,
                            color_discrete_sequence=['#f5576c', '#ffc107', '#667eea']
                        )
                        fig_status.update_layout(
                            showlegend=False,
                            template='plotly_white',
                            xaxis_title='Status',
                            yaxis_title='Number of Responses',
                            height=300
                        )
                        st.plotly_chart(fig_status, use_container_width=True)
                
                with chart_col2:
                    st.markdown("##### Device Type Usage")
                    # 检测设备类型列
                    device_col = None
                    for col in df.columns:
                        col_lower = str(col).lower()
                        if any(kw in col_lower for kw in ['device', '设备', 'platform', 'type']):
                            device_col = col
                            break
                    
                    if device_col:
                        device_counts = df[device_col].value_counts()
                        fig_device = px.pie(
                            values=device_counts.values,
                            names=device_counts.index,
                            color_discrete_sequence=['#667eea', '#11998e', '#f5576c']
                        )
                        fig_device.update_layout(height=300)
                        st.plotly_chart(fig_device, use_container_width=True)
                    else:
                        st.info("未检测到设备类型字段")
                
                # 时间趋势图
                st.markdown("##### Survey Responses Over Time")
                date_cols = []
                for col in df.columns:
                    col_lower = str(col).lower()
                    if any(kw in col_lower for kw in ['date', 'time', '日期', '时间', 'created', 'submitted']):
                        date_cols.append(col)
                
                if date_cols:
                    selected_date_col = st.selectbox("选择时间字段", date_cols, key="trends_date_col")
                    try:
                        df_trends = df.copy()
                        df_trends[selected_date_col] = pd.to_datetime(df_trends[selected_date_col], errors='coerce')
                        df_trends = df_trends.dropna(subset=[selected_date_col])
                        
                        if len(df_trends) > 0:
                            df_trends['date_only'] = df_trends[selected_date_col].dt.date
                            daily_counts = df_trends.groupby('date_only').size().reset_index(name='count')
                            daily_counts['date_only'] = pd.to_datetime(daily_counts['date_only'])
                            
                            fig_trend = px.area(
                                daily_counts, 
                                x='date_only', 
                                y='count',
                                labels={'date_only': 'Submission Date', 'count': 'Number of Responses'}
                            )
                            fig_trend.update_layout(
                                template='plotly_white',
                                hovermode='x unified',
                                height=250
                            )
                            fig_trend.update_traces(fill='tozeroy', line_color='#667eea')
                            st.plotly_chart(fig_trend, use_container_width=True)
                    except Exception as e:
                        st.warning(f"日期解析错误: {e}")
                else:
                    st.info("未检测到时间字段，无法显示时间趋势")

            # ==================== Summary Timeline 页面 ====================
            elif current_page == 'timeline':
                st.markdown("""
                <style>
                .timeline-container {
                    position: relative;
                    padding: 1rem 0;
                }
                .timeline-item {
                    display: flex;
                    align-items: flex-start;
                    gap: 1rem;
                    padding: 1rem;
                    margin-bottom: 0.5rem;
                    background: white;
                    border-radius: 8px;
                    border-left: 4px solid #667eea;
                    box-shadow: 0 1px 3px rgba(0,0,0,0.1);
                }
                .timeline-item.active { border-left-color: #11998e; }
                .timeline-item.closed { border-left-color: #6c757d; }
                .timeline-dot {
                    width: 12px;
                    height: 12px;
                    border-radius: 50%;
                    background: #667eea;
                    margin-top: 0.25rem;
                    flex-shrink: 0;
                }
                .timeline-dot.active { background: #11998e; }
                .timeline-dot.closed { background: #6c757d; }
                .timeline-content {
                    flex: 1;
                }
                .timeline-title {
                    font-weight: 600;
                    font-size: 1rem;
                    color: #1a1a2e;
                    margin-bottom: 0.25rem;
                }
                .timeline-desc {
                    font-size: 0.85rem;
                    color: #6c757d;
                    margin-bottom: 0.5rem;
                }
                .timeline-meta {
                    display: flex;
                    gap: 1rem;
                    font-size: 0.75rem;
                    color: #adb5bd;
                }
                .timeline-status {
                    padding: 0.2rem 0.5rem;
                    border-radius: 4px;
                    font-size: 0.7rem;
                    font-weight: 600;
                }
                .status-active { background: #d4edda; color: #155724; }
                .status-closed { background: #e9ecef; color: #495057; }
                </style>
                """, unsafe_allow_html=True)
                
                st.markdown("#### 📝 调研项目时间线")
                st.markdown("查看所有调研项目的时间安排和状态")
                
                # 尝试检测调研标题列
                title_col = None
                for col in df.columns:
                    col_lower = str(col).lower()
                    if any(kw in col_lower for kw in ['title', 'survey', 'name', '标题', '名称', '调研']):
                        title_col = col
                        break
                
                if title_col:
                    # 按调研标题分组
                    survey_groups = df.groupby(title_col).agg({
                        df.columns[0]: 'count'
                    }).reset_index()
                    survey_groups.columns = [title_col, 'response_count']
                    
                    for idx, row in survey_groups.iterrows():
                        status = 'Active' if idx % 2 == 0 else 'Closed'
                        status_class = 'active' if status == 'Active' else 'closed'
                        
                        st.markdown(f"""
                        <div class="timeline-item {status_class}">
                            <div class="timeline-dot {status_class}"></div>
                            <div class="timeline-content">
                                <div class="timeline-title">{row[title_col]}</div>
                                <div class="timeline-desc">共收集 {row['response_count']} 条回复</div>
                                <div class="timeline-meta">
                                    <span class="timeline-status status-{status_class}">{status}</span>
                                </div>
                            </div>
                        </div>
                        """, unsafe_allow_html=True)
                else:
                    # 如果没有标题列，显示数据摘要时间线
                    st.markdown(f"""
                    <div class="timeline-item active">
                        <div class="timeline-dot active"></div>
                        <div class="timeline-content">
                            <div class="timeline-title">当前调研数据</div>
                            <div class="timeline-desc">共 {len(df)} 条回复，{len(df.columns)} 个字段</div>
                            <div class="timeline-meta">
                                <span class="timeline-status status-active">Active</span>
                            </div>
                        </div>
                    </div>
                    """, unsafe_allow_html=True)
                    
                    st.info("💡 提示：上传包含调研标题字段的数据可以显示更详细的时间线")

            # ==================== Respondent Gallery 页面 ====================
            elif current_page == 'gallery':
                st.markdown("""
                <style>
                .gallery-grid {
                    display: grid;
                    grid-template-columns: repeat(auto-fill, minmax(280px, 1fr));
                    gap: 1.5rem;
                    padding: 1rem 0;
                }
                .gallery-card {
                    background: white;
                    border-radius: 12px;
                    overflow: hidden;
                    box-shadow: 0 2px 8px rgba(0,0,0,0.08);
                    transition: all 0.3s ease;
                }
                .gallery-card:hover {
                    box-shadow: 0 8px 24px rgba(0,0,0,0.12);
                    transform: translateY(-4px);
                }
                .gallery-image {
                    height: 160px;
                    background: linear-gradient(135deg, #667eea 0%, #764ba2 100%);
                    display: flex;
                    align-items: center;
                    justify-content: center;
                    font-size: 4rem;
                }
                .gallery-image.alt1 { background: linear-gradient(135deg, #f093fb 0%, #f5576c 100%); }
                .gallery-image.alt2 { background: linear-gradient(135deg, #4facfe 0%, #00f2fe 100%); }
                .gallery-image.alt3 { background: linear-gradient(135deg, #11998e 0%, #38ef7d 100%); }
                .gallery-content {
                    padding: 1.25rem;
                }
                .gallery-name {
                    font-size: 1.1rem;
                    font-weight: 600;
                    color: #1a1a2e;
                    margin-bottom: 0.75rem;
                }
                .gallery-field {
                    display: flex;
                    justify-content: space-between;
                    padding: 0.4rem 0;
                    border-bottom: 1px solid #f0f0f0;
                }
                .gallery-field:last-child { border-bottom: none; }
                .gallery-field-label {
                    font-size: 0.75rem;
                    color: #6c757d;
                }
                .gallery-field-value {
                    font-size: 0.85rem;
                    color: #333;
                    font-weight: 500;
                }
                .gallery-tag {
                    display: inline-block;
                    padding: 0.25rem 0.75rem;
                    border-radius: 20px;
                    font-size: 0.75rem;
                    font-weight: 500;
                    margin-top: 0.5rem;
                }
                .tag-female { background: #fce4ec; color: #c2185b; }
                .tag-male { background: #e3f2fd; color: #1565c0; }
                .tag-other { background: #f3e5f5; color: #7b1fa2; }
                </style>
                """, unsafe_allow_html=True)
                
                st.markdown("#### 👤 受访者画廊")
                st.markdown("以卡片形式浏览受访者信息")
                
                # 筛选器
                filter_col1, filter_col2 = st.columns([1, 1])
                
                # 检测可能的筛选字段
                gender_col = None
                for col in df.columns:
                    col_lower = str(col).lower()
                    if any(kw in col_lower for kw in ['gender', '性别', 'sex']):
                        gender_col = col
                        break
                
                with filter_col1:
                    if gender_col:
                        gender_options = ['全部'] + list(df[gender_col].dropna().unique())
                        selected_gender = st.selectbox("筛选性别", gender_options, key="gallery_gender_filter")
                
                # 应用筛选
                display_df = df.copy()
                if gender_col and 'selected_gender' in dir() and selected_gender != '全部':
                    display_df = display_df[display_df[gender_col] == selected_gender]
                
                # 分页
                items_per_page = 9
                total_pages = (len(display_df) - 1) // items_per_page + 1
                
                page_col1, page_col2, page_col3 = st.columns([1, 2, 1])
                with page_col2:
                    current_gallery_page = st.number_input(
                        f"页码 (共 {total_pages} 页)",
                        min_value=1,
                        max_value=max(1, total_pages),
                        value=1,
                        key="gallery_page"
                    )
                
                start_idx = (current_gallery_page - 1) * items_per_page
                end_idx = start_idx + items_per_page
                page_df = display_df.iloc[start_idx:end_idx]
                
                # 检测姓名列
                name_col = None
                for col in df.columns:
                    col_lower = str(col).lower()
                    if any(kw in col_lower for kw in ['name', 'respondent', '姓名', '名字', '受访者']):
                        name_col = col
                        break
                
                # 渐变色列表
                gradient_colors = [
                    "#667eea, #764ba2",
                    "#f093fb, #f5576c",
                    "#4facfe, #00f2fe",
                    "#11998e, #38ef7d"
                ]
                
                # 每行显示3个卡片
                for row_start in range(0, len(page_df), 3):
                    cols = st.columns(3, gap="medium")
                    row_data = page_df.iloc[row_start:row_start + 3]
                    
                    for col_idx, (_, row) in enumerate(row_data.iterrows()):
                        gradient = gradient_colors[(row_start + col_idx) % 4]
                        
                        # 获取显示名称
                        if name_col and pd.notna(row.get(name_col)):
                            display_name = str(row[name_col])[:30]
                        else:
                            display_name = f"Respondent {start_idx + row_start + col_idx + 1}"
                        
                        with cols[col_idx]:
                            # 使用容器包裹整个卡片
                            with st.container(border=True):
                                # 渐变色头部
                                st.markdown(f'''
                                <div style="background: linear-gradient(135deg, {gradient}); 
                                            height: 100px; 
                                            margin: -1rem -1rem 1rem -1rem; 
                                            border-radius: 0.5rem 0.5rem 0 0;
                                            display: flex; 
                                            align-items: center; 
                                            justify-content: center;">
                                    <span style="font-size: 2.5rem;">👤</span>
                                </div>
                                ''', unsafe_allow_html=True)
                                
                                # 姓名
                                st.markdown(f"**{display_name}**")
                                
                                # 性别标签
                                if gender_col and pd.notna(row.get(gender_col)):
                                    gender_val = str(row[gender_col])
                                    st.caption(f"🏷️ {gender_val}")
                                
                                # 显示字段（最多4个）
                                shown_cols = [c for c in df.columns if c not in [name_col, gender_col]][:4]
                                for field_col in shown_cols:
                                    if pd.notna(row.get(field_col)):
                                        val = str(row[field_col])[:40]
                                        st.markdown(f"<small style='color: #6c757d;'>{field_col}</small><br/><span style='font-size: 0.9rem;'>{val}</span>", unsafe_allow_html=True)
                
                st.caption(f"显示 {start_idx + 1}-{min(end_idx, len(display_df))} / 共 {len(display_df)} 条记录")

            # ==================== Export Reports 页面 ====================
            elif current_page == 'export':
                st.markdown("#### 📤 导出分析报告")
                st.markdown("将分析结果导出为专业报告")
                
                # 导出选项
                export_col1, export_col2 = st.columns([2, 1])
                
                with export_col1:
                    st.markdown("##### 选择导出内容")
                    export_data_overview = st.checkbox("📋 数据概览", value=True, key="export_data_overview")
                    export_distribution = st.checkbox("📊 分布分析", value=True, key="export_distribution")
                    export_cross = st.checkbox("🔀 交叉分析", value=False, key="export_cross")
                    export_ai = st.checkbox("🤖 AI 洞察", value=False, key="export_ai")
                
                with export_col2:
                    st.markdown("##### 导出格式")
                    export_format_choice = st.radio(
                        "选择格式",
                        ["Word (.docx)", "PPT (.pptx)", "PDF (.pdf)", "Excel (.xlsx)"],
                        key="export_format_choice"
                    )
                
                st.markdown("---")
                
                # 快速导出按钮
                col1, col2, col3 = st.columns(3)
                
                with col1:
                    # CSV 下载
                    csv_data = df.to_csv(index=False).encode('utf-8-sig')
                    st.download_button(
                        label="📥 下载原始数据 (CSV)",
                        data=csv_data,
                        file_name=f"survey_data_{datetime.now().strftime('%Y%m%d')}.csv",
                        mime="text/csv",
                        use_container_width=True
                    )
                
                with col2:
                    # Excel 下载
                    try:
                        excel_buffer = BytesIO()
                        with pd.ExcelWriter(excel_buffer, engine='openpyxl') as writer:
                            df.to_excel(writer, index=False, sheet_name='Survey Data')
                        excel_data = excel_buffer.getvalue()
                        st.download_button(
                            label="📥 下载数据 (Excel)",
                            data=excel_data,
                            file_name=f"survey_data_{datetime.now().strftime('%Y%m%d')}.xlsx",
                            mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
                            use_container_width=True
                        )
                    except Exception as e:
                        st.button("📥 下载数据 (Excel)", disabled=True, use_container_width=True)
                        st.caption(f"Excel 导出不可用: {e}")
                
                with col3:
                    if st.button("📊 生成完整报告", type="primary", use_container_width=True, key="generate_full_report"):
                        st.info("🔄 报告生成功能正在开发中...")
                
                # 数据预览
                st.markdown("##### 📋 数据预览")
                st.dataframe(df.head(10), use_container_width=True)

        except Exception as e:
            st.error(f"表格读取错误: {e}")

    # ==========================================
    # 模块 B: 非结构化文本分析 (PDF/Word)
    # ==========================================
    elif file_type in ['pdf', 'docx']:
        try:
            text = read_pdf(uploaded_file) if file_type == 'pdf' else read_docx(uploaded_file)
            
            st.markdown("### 📑 文档智能分析")
            
            # 文本统计卡片
            words_list = text.split()
            col1, col2, col3, col4 = st.columns(4)
            
            with col1:
                st.markdown(f"""
                <div class="metric-card">
                    <div class="metric-label">总字符数</div>
                    <div class="metric-value">{len(text):,}</div>
                </div>
                """, unsafe_allow_html=True)
            
            with col2:
                st.markdown(f"""
                <div class="metric-card">
                    <div class="metric-label">总词数</div>
                    <div class="metric-value">{len(words_list):,}</div>
                </div>
                """, unsafe_allow_html=True)
            
            with col3:
                st.markdown(f"""
                <div class="metric-card">
                    <div class="metric-label">段落数</div>
                    <div class="metric-value">{len(text.split(chr(10))):,}</div>
                </div>
                """, unsafe_allow_html=True)
            
            with col4:
                unique_words = len(set([w for w in words_list if len(w) > 1]))
                st.markdown(f"""
                <div class="metric-card">
                    <div class="metric-label">唯一词数</div>
                    <div class="metric-value">{unique_words:,}</div>
                </div>
                """, unsafe_allow_html=True)
            
            st.markdown("<br>", unsafe_allow_html=True)
            
            # 内容展示
            col1, col2 = st.columns([1, 2])
            with col1:
                st.markdown("#### 📄 文本内容预览")
                preview_text = text[:1000] + "..." if len(text) > 1000 else text
                st.text_area("", preview_text, height=400, label_visibility="collapsed")
            
            with col2:
                if text.strip():
                    st.markdown("#### ☁️ 智能词云分析")
                    try:
                        wc = WordCloud(
                            font_path='simhei.ttf', 
                            width=800, 
                            height=400, 
                            background_color='white',
                            colormap='Purples',
                            max_words=100
                        ).generate(text)
                    
                        # 词云显示
                        fig, ax = plt.subplots(figsize=(10, 5))
                        ax.imshow(wc, interpolation='bilinear')
                        ax.axis("off")
                        st.pyplot(fig)
                    except:
                        st.warning("⚠️ 词云生成失败，可能缺少中文字体文件")
                    
                    # 词频条形图
                    st.markdown("#### 📊 高频词统计")
                    words = [w for w in text.split() if len(w) > 1]
                    word_counts = pd.Series(words).value_counts().head(15).reset_index()
                    word_counts.columns = ['词汇', '频率']
                    
                    fig_bar = px.bar(
                        word_counts, 
                        x='频率', y='词汇', 
                        orientation='h', 
                        title="Top 15 高频词汇",
                        color='频率',
                        color_continuous_scale='Purples'
                    )
                    fig_bar.update_layout(
                        plot_bgcolor='rgba(0,0,0,0)',
                        paper_bgcolor='rgba(0,0,0,0)',
                        font=dict(size=12),
                        showlegend=False,
                        height=400
                    )
                    st.plotly_chart(fig_bar, use_container_width=True)
                else:
                    st.warning("⚠️ 文档内容为空")
        except Exception as e:
            st.error(f"❌ 文档解析错误: {e}")

else:
    # 空状态展示
    st.markdown("<br>", unsafe_allow_html=True)
    
    col1, col2, col3 = st.columns([1, 2, 1])
    with col2:
        st.markdown("""
        <div style="text-align: center; padding: 4rem 2rem; background: linear-gradient(135deg, #667eea 0%, #764ba2 100%); border-radius: 20px; color: white; box-shadow: 0 10px 40px rgba(102, 126, 234, 0.4); position: relative; overflow: hidden;">
            <div style="position: absolute; top: -50%; right: -50%; width: 200%; height: 200%; background: radial-gradient(circle, rgba(255,255,255,0.1) 0%, transparent 70%); animation: pulse 3s ease-in-out infinite;"></div>
            <div style="position: relative; z-index: 1;">
                <div style="font-size: 4rem; margin-bottom: 1rem;">📊</div>
                <h2 style="margin-bottom: 1rem; font-size: 2rem; font-weight: 700;">开始您的数据分析之旅</h2>
                <p style="font-size: 1.1rem; opacity: 0.95; margin-bottom: 1.5rem;">上传文件,解锁强大的数据洞察能力</p>
                <div style="display: inline-block; background: rgba(255,255,255,0.2); padding: 0.8rem 2rem; border-radius: 25px; backdrop-filter: blur(10px);">
                    <span style="font-size: 0.9rem; font-weight: 600;">👈 点击左侧上传按钮开始</span>
                </div>
            </div>
        </div>
        """, unsafe_allow_html=True)
    
    st.markdown("<br><br>", unsafe_allow_html=True)
    
    # 功能介绍
    st.markdown('<h3 style="text-align: center; color: #667eea; font-weight: 700; margin-bottom: 2rem;">✨ 核心功能</h3>', unsafe_allow_html=True)
    
    col1, col2 = st.columns(2, gap="large")
    
    with col1:
        st.markdown("""
        <div style="background: white; padding: 2rem; border-radius: 15px; box-shadow: 0 4px 20px rgba(0,0,0,0.08); height: 100%; border-top: 4px solid #667eea;">
            <h4 style="color: #667eea; margin-bottom: 1.5rem; font-size: 1.3rem;">📊 结构化数据分析</h4>
            <div style="margin-bottom: 1rem;">
                <div style="display: flex; align-items: start; margin-bottom: 0.8rem;">
                    <span style="color: #667eea; margin-right: 0.5rem; font-size: 1.2rem;">✓</span>
                    <div>
                        <strong>数据总览</strong><br>
                        <small style="color: #6c757d;">快速了解数据概况和质量</small>
                    </div>
                </div>
                <div style="display: flex; align-items: start; margin-bottom: 0.8rem;">
                    <span style="color: #667eea; margin-right: 0.5rem; font-size: 1.2rem;">✓</span>
                    <div>
                        <strong>单变量分析</strong><br>
                        <small style="color: #6c757d;">柱状图、饼图、直方图等多种可视化</small>
                    </div>
                </div>
                <div style="display: flex; align-items: start; margin-bottom: 0.8rem;">
                    <span style="color: #667eea; margin-right: 0.5rem; font-size: 1.2rem;">✓</span>
                    <div>
                        <strong>交叉分析</strong><br>
                        <small style="color: #6c757d;">散点图、箱线图、热力图等高级分析</small>
                    </div>
                </div>
                <div style="display: flex; align-items: start;">
                    <span style="color: #667eea; margin-right: 0.5rem; font-size: 1.2rem;">✓</span>
                    <div>
                        <strong>AI智能对话</strong><br>
                        <small style="color: #6c757d;">与AI对话获取数据洞察</small>
                    </div>
                </div>
            </div>
            <div style="margin-top: 1.5rem; padding-top: 1rem; border-top: 1px solid #e9ecef;">
                <span style="background: #e7f3ff; color: #0066cc; padding: 0.3rem 0.8rem; border-radius: 15px; font-size: 0.85rem; font-weight: 600;">
                    📄 CSV, Excel (.xlsx)
                </span>
            </div>
        </div>
        """, unsafe_allow_html=True)
    
    with col2:
        st.markdown("""
        <div style="background: white; padding: 2rem; border-radius: 15px; box-shadow: 0 4px 20px rgba(0,0,0,0.08); height: 100%; border-top: 4px solid #764ba2;">
            <h4 style="color: #764ba2; margin-bottom: 1.5rem; font-size: 1.3rem;">📑 文本智能分析</h4>
            <div style="margin-bottom: 1rem;">
                <div style="display: flex; align-items: start; margin-bottom: 0.8rem;">
                    <span style="color: #764ba2; margin-right: 0.5rem; font-size: 1.2rem;">✓</span>
                    <div>
                        <strong>内容提取</strong><br>
                        <small style="color: #6c757d;">自动提取PDF/Word文档内容</small>
                    </div>
                </div>
                <div style="display: flex; align-items: start; margin-bottom: 0.8rem;">
                    <span style="color: #764ba2; margin-right: 0.5rem; font-size: 1.2rem;">✓</span>
                    <div>
                        <strong>词云生成</strong><br>
                        <small style="color: #6c757d;">可视化展示文本关键词</small>
                    </div>
                </div>
                <div style="display: flex; align-items: start; margin-bottom: 0.8rem;">
                    <span style="color: #764ba2; margin-right: 0.5rem; font-size: 1.2rem;">✓</span>
                    <div>
                        <strong>词频统计</strong><br>
                        <small style="color: #6c757d;">识别高频词汇和主题</small>
                    </div>
                </div>
                <div style="display: flex; align-items: start;">
                    <span style="color: #764ba2; margin-right: 0.5rem; font-size: 1.2rem;">✓</span>
                    <div>
                        <strong>文本度量</strong><br>
                        <small style="color: #6c757d;">字符数、词数、段落数统计</small>
                    </div>
                </div>
            </div>
            <div style="margin-top: 1.5rem; padding-top: 1rem; border-top: 1px solid #e9ecef;">
                <span style="background: #f3e7ff; color: #6a0dad; padding: 0.3rem 0.8rem; border-radius: 15px; font-size: 0.85rem; font-weight: 600;">
                    📄 PDF, Word (.docx)
                </span>
            </div>
        </div>
        """, unsafe_allow_html=True)
    
    st.markdown("<br>", unsafe_allow_html=True)
    st.markdown("---")
    
    # 使用提示
    st.markdown('<h3 style="text-align: center; color: #667eea; font-weight: 700; margin: 2rem 0;">💡 使用提示</h3>', unsafe_allow_html=True)
    
    tips_col1, tips_col2, tips_col3 = st.columns(3, gap="large")
    
    with tips_col1:
        st.markdown("""
        <div style="background: linear-gradient(135d, #d1ecf1 0%, #bee5eb 100%); padding: 1.5rem; border-radius: 12px; border-left: 4px solid #17a2b8; height: 100%;">
            <h5 style="color: #0c5460; margin-bottom: 1rem; font-weight: 700;">📈 数据准备</h5>
            <ul style="color: #0c5460; font-size: 0.9rem; margin: 0; padding-left: 1.2rem;">
                <li style="margin-bottom: 0.5rem;">确保Excel/CSV第一行为列名</li>
                <li style="margin-bottom: 0.5rem;">数据格式保持一致</li>
                <li>避免空行和特殊字符</li>
            </ul>
        </div>
        """, unsafe_allow_html=True)
    
    with tips_col2:
        st.markdown("""
        <div style="background: linear-gradient(135deg, #d4edda 0%, #c3e6cb 100%); padding: 1.5rem; border-radius: 12px; border-left: 4px solid #28a745; height: 100%;">
            <h5 style="color: #155724; margin-bottom: 1rem; font-weight: 700;">🎨 可视化技巧</h5>
            <ul style="color: #155724; font-size: 0.9rem; margin: 0; padding-left: 1.2rem;">
                <li style="margin-bottom: 0.5rem;">数值型数据适合散点图</li>
                <li style="margin-bottom: 0.5rem;">类别型数据适合柱状图</li>
                <li>自定义配色方案</li>
            </ul>
        </div>
        """, unsafe_allow_html=True)
    
    with tips_col3:
        st.markdown("""
        <div style="background: linear-gradient(135deg, #fff3cd 0%, #ffeaa7 100%); padding: 1.5rem; border-radius: 12px; border-left: 4px solid #ffc107; height: 100%;">
            <h5 style="color: #856404; margin-bottom: 1rem; font-weight: 700;">⚡ 性能优化</h5>
            <ul style="color: #856404; font-size: 0.9rem; margin: 0; padding-left: 1.2rem;">
                <li style="margin-bottom: 0.5rem;">建议文件 < 50MB</li>
                <li style="margin-bottom: 0.5rem;">大数据集需较长时间</li>
                <li>选择性分析关键变量</li>
            </ul>
        </div>
        """, unsafe_allow_html=True)
    
    st.markdown("<br><br>", unsafe_allow_html=True)